from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT = OUT_DIR / "iegp_revmed_reference_template_editable.pptx"
OUT_MAIN = OUT_DIR / "iegp_revmed_reference_template.pptx"

WIDE = Inches(16)
HIGH = Inches(9)

PAPER = RGBColor(255, 255, 255)
INK = RGBColor(5, 41, 39)
TEAL_950 = RGBColor(0, 63, 59)
TEAL_900 = RGBColor(0, 85, 79)
TEAL_800 = RGBColor(8, 114, 103)
TEAL_700 = RGBColor(11, 134, 120)
MINT = RGBColor(21, 210, 137)
AQUA = RGBColor(78, 229, 194)
PURPLE = RGBColor(154, 43, 131)
ORANGE = RGBColor(255, 107, 53)
AMBER = RGBColor(242, 183, 5)
GRAY_050 = RGBColor(246, 248, 247)
GRAY_100 = RGBColor(238, 242, 241)
GRAY_200 = RGBColor(221, 229, 227)
GRAY_300 = RGBColor(199, 210, 207)
GRAY_500 = RGBColor(102, 115, 111)
GRAY_700 = RGBColor(49, 64, 61)
RISK = RGBColor(216, 58, 46)


def rgb(hex_string):
    h = hex_string.strip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def set_line(shape, color=GRAY_200, width=0.75, transparency=0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if transparency:
        shape.line.transparency = transparency


def no_line(shape):
    shape.line.fill.background()


def add_box(slide, x, y, w, h, fill=GRAY_100, line=None, radius=False, transparency=0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    set_fill(shp, fill, transparency)
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line)
    return shp


def add_line(slide, x1, y1, x2, y2, color=GRAY_200, width=1):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    return shp


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=12,
    color=INK,
    bold=False,
    align=None,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_rich_heading(slide, strong, rest, x=0.72, y=0.72, w=14.3, h=0.45, size=20):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = strong
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = TEAL_950
    r = p.add_run()
    r.text = rest
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.color.rgb = TEAL_900
    return tb


def add_bullets(slide, items, x, y, w, h, size=9.5, color=INK, bullet_color=MINT, gap=0.28):
    top = y
    for item in items:
        add_box(slide, x, top + 0.08, 0.06, 0.06, fill=bullet_color)
        add_text(slide, item, x + 0.18, top, w - 0.18, gap, size=size, color=color)
        top += gap


def add_top_bars(slide, num=None):
    add_box(slide, 0.72, 0.70, 0.45, 0.035, fill=MINT)
    add_box(slide, 1.26, 0.70, 0.67, 0.035, fill=AQUA)
    if num is not None:
        add_text(slide, f"{num:02d}", 15.00, 0.55, 0.25, 0.16, size=6, color=TEAL_700, bold=True, align=PP_ALIGN.CENTER)


def add_footer(slide, right):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(0.72), Inches(8.57), Inches(0.13), Inches(0.13))
    d.fill.background()
    set_line(d, TEAL_700, width=1.2)
    add_text(slide, "IEGP template", 0.90, 8.52, 1.2, 0.18, size=5.5, color=TEAL_700, bold=True)
    add_text(slide, right, 14.55, 8.52, 0.7, 0.18, size=5.2, color=GRAY_500, align=PP_ALIGN.RIGHT)


def add_base(slide, num, footer_text):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_top_bars(slide, num)
    add_footer(slide, footer_text)


def add_dark_card(slide, title, bullets, x, y, w, h, accent=MINT):
    add_box(slide, x, y, w, h, fill=TEAL_950, radius=True)
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.35, size=11, color=PAPER, bold=True)
    add_box(slide, x + 0.18, y + 0.63, 0.52, 0.035, fill=accent)
    add_bullets(slide, bullets, x + 0.18, y + 0.82, w - 0.36, h - 0.95, size=8.4, color=PAPER, bullet_color=accent, gap=0.30)


def add_light_panel(slide, title, bullets, x, y, w, h, accent=TEAL_700):
    add_box(slide, x, y, w, h, fill=GRAY_100)
    add_text(slide, title, x + 0.16, y + 0.15, w - 0.32, 0.26, size=10.2, color=TEAL_950, bold=True)
    add_bullets(slide, bullets, x + 0.16, y + 0.52, w - 0.32, h - 0.58, size=8.4, color=INK, bullet_color=accent, gap=0.27)


def add_callout(slide, text, x, y, w, h, size=9.5):
    add_box(slide, x, y, w, h, fill=rgb("e8f4f1"))
    add_box(slide, x, y, 0.07, h, fill=TEAL_700)
    add_text(slide, text, x + 0.22, y + 0.10, w - 0.35, h - 0.15, size=size, color=INK)


def add_metric(slide, value, label, x, y, w, h):
    add_box(slide, x, y, w, h, fill=PAPER, line=GRAY_200)
    add_text(slide, value, x + 0.12, y + 0.12, w - 0.24, 0.32, size=18, color=TEAL_700, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.12, y + 0.50, w - 0.24, h - 0.58, size=7.4, color=GRAY_700, align=PP_ALIGN.CENTER)


def add_table(slide, headers, rows, x, y, w, h, col_widths=None, font_size=6.8):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(w * cw / total)
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        set_fill(cell, TEAL_800)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = PAPER
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            set_fill(cell, GRAY_050 if i % 2 == 0 else PAPER)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Arial"
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = INK if not (j == 0) else TEAL_950
                    r.font.bold = j == 0
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
    return table_shape


def add_section(slide, num, title, subtitle, footer):
    add_base(slide, num, footer)
    add_line(slide, 6.0, -0.3, 10.6, 3.4, color=GRAY_200, width=0.7)
    add_line(slide, 8.8, -0.2, 13.4, 3.5, color=GRAY_200, width=0.7)
    add_text(slide, title, 5.5, 3.15, 5.0, 0.35, size=18, color=TEAL_950, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 5.5, 3.50, 5.0, 0.38, size=13.5, color=TEAL_900, align=PP_ALIGN.CENTER)
    add_text(slide, f"Section {max(1, num // 4)} of 6", 6.6, 3.90, 2.8, 0.16, size=5.5, color=GRAY_500, align=PP_ALIGN.CENTER)
    add_box(slide, 0.72, 6.23, 14.55, 1.15, fill=TEAL_800)
    p1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(3.25), Inches(6.50), Inches(2.75), Inches(0.68))
    set_fill(p1, TEAL_950, 30); no_line(p1)
    p2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(5.60), Inches(6.50), Inches(2.50), Inches(0.68))
    set_fill(p2, AQUA, 72); no_line(p2)
    p3 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(12.2), Inches(6.23), Inches(2.6), Inches(1.15))
    set_fill(p3, TEAL_950, 68); no_line(p3)


def make_deck():
    prs = Presentation()
    prs.slide_width = WIDE
    prs.slide_height = HIGH
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    add_box(s, 0, 0, 5.45, 9, fill=TEAL_800)
    for x, y, sz, color, tr in [(0.4, 4.8, 3.0, MINT, 55), (1.2, 0.4, 3.9, TEAL_950, 35), (2.3, 3.0, 3.0, AQUA, 70)]:
        d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(sz), Inches(sz))
        set_fill(d, color, tr); no_line(d)
    add_top_bars(s)
    add_text(s, "Integrated Evidence", 6.45, 1.10, 7.7, 0.55, size=30, color=TEAL_950, bold=True)
    add_text(s, "Generation Plan", 6.45, 1.62, 6.9, 0.55, size=28, color=TEAL_900)
    add_text(s, "Editable template for product, indication, and geography-specific evidence strategy.", 6.45, 2.55, 6.8, 0.45, size=12.5, color=GRAY_700)
    for i, (label, val) in enumerate([("Product", "[Product / asset name]"), ("Indication", "[Therapy area / population]"), ("Version", "[Draft / final / refresh]"), ("Owner", "[Function / accountable lead]")]):
        x = 6.45 + (i % 2) * 3.9
        y = 3.70 + (i // 2) * 0.72
        add_line(s, x, y, x + 3.45, y, color=GRAY_300)
        add_text(s, label, x, y + 0.10, 1.5, 0.18, size=7.5, color=GRAY_700, bold=True)
        add_text(s, val, x, y + 0.30, 3.4, 0.22, size=7.3, color=INK)
    add_text(s, "Replace placeholders before use", 13.3, 8.48, 1.2, 0.15, size=5.2, color=GRAY_500, align=PP_ALIGN.RIGHT)

    # 2
    s = prs.slides.add_slide(blank); add_base(s, 2, "General instructions")
    add_rich_heading(s, "How to use this template", " | Build the plan around evidence decisions")
    cards = [
        ("1. Define the decision frame", ["Clarify launch, access, guideline, regulatory, and lifecycle decisions.", "Set geography, timing, and stakeholder scope."]),
        ("2. Map evidence need vs have", ["Inventory completed, ongoing, and planned evidence.", "Classify gaps by stakeholder, market, evidence type, and urgency."]),
        ("3. Convert gaps into tactics", ["Prioritize research questions and study concepts.", "Assign owner, budget, timing, readout, and refresh cadence."]),
    ]
    for i, card in enumerate(cards):
        add_dark_card(s, card[0], card[1], 0.72 + i * 4.84, 1.48, 4.55, 1.38)
    add_callout(s, "Design note: The slide system uses a biotech clinical briefing grammar: compact teal modules, hard evidence panels, visible decision signals, and sparse section pacing. Swap color tokens for company colors without changing layout.", 0.72, 3.10, 14.55, 0.55, size=8.2)

    # 3
    s = prs.slides.add_slide(blank); add_section(s, 3, "Evidence strategy", "and product context", "Section divider")

    # 4
    s = prs.slides.add_slide(blank); add_base(s, 4, "Executive briefing")
    add_rich_heading(s, "Executive summary", " | Evidence priorities for the next planning cycle")
    exec_cards = [
        ("Priority 1: Access-critical gap", ["Gap: [Comparator, outcome, or local-market evidence need].", "Tactic: [Study / analysis / synthesis].", "Decision impact: [HTA, formulary, label, guideline]."]),
        ("Priority 2: Differentiation gap", ["Gap: [Value message lacks support].", "Tactic: [RWE, HEOR, post-hoc, registry].", "Timing: [Readout needed by milestone]."]),
        ("Priority 3: Patient impact gap", ["Gap: [PRO, burden, preference, adherence].", "Tactic: [COA / preference / qualitative study].", "Owner: [Function / workstream]."]),
    ]
    for i, card in enumerate(exec_cards):
        add_dark_card(s, card[0], card[1], 0.72 + i * 4.84, 1.40, 4.55, 1.85)
    for i, metric in enumerate([("[x]", "Priority evidence gaps to close"), ("[x]", "New studies or analyses proposed"), ("[x]", "Decision gates supported"), ("[x]", "Markets requiring adaptation")]):
        add_metric(s, metric[0], metric[1], 0.72 + i * 3.68, 3.78, 3.35, 0.82)

    # 5
    s = prs.slides.add_slide(blank); add_base(s, 5, "Decision alignment")
    add_rich_heading(s, "Key decisions needed", " | Align the evidence plan to funded choices")
    add_table(s, ["Decision", "Recommendation", "Evidence basis", "Owner", "Due"], [
        ["[Fund study concept A]", "Proceed", "[Closes payer-critical comparator gap]", "[HEOR]", "[Qx]"],
        ["[Include market X in RWE design]", "Decide", "[Local HTA requirement or guideline timing]", "[Access]", "[Qx]"],
        ["[Retire low-impact tactic]", "Stop", "[No decision linkage or duplicate evidence]", "[Medical]", "[Qx]"],
        ["[Adapt claim language]", "Revise", "[Evidence strength below desired support]", "[Comms]", "[Qx]"],
        ["[Activate affiliate input cycle]", "Proceed", "[Local evidence needs not yet captured]", "[Regional]", "[Qx]"],
    ], 0.72, 1.45, 7.45, 2.55, [2.2, 1.1, 3.0, 1.0, 0.7], 5.8)
    add_box(s, 8.55, 1.48, 5.2, 2.2, fill=TEAL_800, transparency=18)
    for x, y, size in [(9.45, 1.95, 1.0), (10.08, 1.58, 1.55), (10.88, 1.90, 1.0)]:
        d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(size), Inches(size))
        d.fill.background(); set_line(d, AQUA, width=2, transparency=20)
    add_box(s, 0.72, 6.75, 14.55, 0.22, fill=TEAL_800)
    add_text(s, "Use this slide to make the IEGP a decision instrument, not a catalog of possible studies.", 1.0, 6.75, 14.0, 0.22, size=6.5, color=PAPER, bold=True, align=PP_ALIGN.CENTER)

    # 6
    s = prs.slides.add_slide(blank); add_base(s, 6, "Asset context")
    add_rich_heading(s, "Asset snapshot", " | Product context that shapes evidence needs")
    add_light_panel(s, "[Product / asset name]", ["Mechanism: [MoA / platform]", "Population: [Target patients and line of therapy]", "Lifecycle stage: [Pre-launch / launch / mature / expansion]", "Current evidence base: [Pivotal, extension, RWE, HEOR]", "Major decision windows: [Submission, HTA, guideline, launch]"], 0.72, 1.54, 4.55, 2.30)
    add_box(s, 5.65, 1.35, 8.9, 3.05, fill=GRAY_100)
    for i, (title, accent) in enumerate([("Clinical promise", TEAL_800), ("Access question", MINT), ("Patient value", PURPLE)]):
        x = 6.0 + i * 2.65
        add_box(s, x, 1.70, 2.25, 1.08, fill=PAPER, line=GRAY_200)
        add_box(s, x, 1.70, 2.25, 0.05, fill=accent)
        add_text(s, title, x + 0.12, 1.86, 2.0, 0.20, size=7.8, color=TEAL_950, bold=True)
        add_text(s, "[Efficacy, safety, unmet need, differentiation hypothesis]", x + 0.12, 2.10, 2.0, 0.35, size=5.8, color=GRAY_700)
        for j, c in enumerate([TEAL_700, MINT, GRAY_300, PURPLE]):
            add_box(s, x + 0.14 + j * 0.45, 2.55, 0.36, 0.04, fill=c)
    add_callout(s, "Replace this with the core product profile and the evidence implications that follow from it.", 5.95, 3.35, 7.95, 0.50, size=7.5)

    # 7
    s = prs.slides.add_slide(blank); add_base(s, 7, "Strategy translation")
    add_rich_heading(s, "Strategic imperatives", " | Convert product strategy into evidence requirements")
    strategic = [
        ("Launch and access", ["Secure reimbursement in priority markets", "Support payer value story and budget impact", "Prepare HTA dossier evidence inputs"]),
        ("Differentiation", ["Demonstrate value vs standard of care", "Support claims in defined subpopulations", "Clarify real-world positioning"]),
        ("Lifecycle growth", ["Enable label expansion or sequencing", "Support guideline inclusion", "Build patient-centered evidence base"]),
    ]
    for i, card in enumerate(strategic):
        add_dark_card(s, card[0], card[1], 0.72 + i * 4.84, 1.45, 4.55, 1.80)
    add_box(s, 0.72, 6.75, 14.55, 0.28, fill=TEAL_800)
    add_text(s, "Strategic imperative -> stakeholder need -> evidence gap -> study tactic -> decision impact", 0.92, 6.78, 14.1, 0.18, size=7.0, color=PAPER, bold=True, align=PP_ALIGN.CENTER)

    # 8
    s = prs.slides.add_slide(blank); add_section(s, 8, "Evidence landscape", "and stakeholder needs", "Section divider")

    # 9
    s = prs.slides.add_slide(blank); add_base(s, 9, "Stakeholder matrix")
    add_rich_heading(s, "Stakeholder evidence needs", " | Map what each audience must believe")
    headers = ["", "Regulators", "HTA / payers", "HCPs", "Patients", "Guidelines"]
    rows = [
        ["Clinical benefit", "Have", "Partial", "Have", "Partial", "Gap"],
        ["Comparative value", "Partial", "Gap", "Planned", "Partial", "Gap"],
        ["Safety in practice", "Have", "Partial", "Planned", "Partial", "Partial"],
        ["Quality of life", "Partial", "Gap", "Partial", "Gap", "Planned"],
        ["Economic impact", "N/A", "Gap", "Low", "Low", "Partial"],
    ]
    x0, y0, cw, ch = 0.72, 1.55, 2.25, 0.43
    for j, h in enumerate(headers):
        add_text(s, h, x0 + j * cw, y0, cw - 0.04, 0.25, size=7, color=GRAY_700, bold=True, align=PP_ALIGN.CENTER)
    color_map = {"Have": TEAL_700, "Partial": MINT, "Gap": ORANGE, "Planned": PURPLE, "N/A": GRAY_300, "Low": GRAY_300}
    for i, row in enumerate(rows):
        y = y0 + 0.38 + i * ch
        add_box(s, x0, y, cw, 0.34, fill=GRAY_100)
        add_text(s, row[0], x0 + 0.06, y + 0.08, cw - 0.12, 0.14, size=6.5, color=TEAL_950, bold=True)
        for j, val in enumerate(row[1:], start=1):
            add_box(s, x0 + j * cw, y, cw - 0.06, 0.34, fill=color_map[val])
            add_text(s, val, x0 + j * cw, y + 0.08, cw - 0.06, 0.14, size=6.5, color=PAPER if val not in ["N/A", "Low"] else INK, bold=True, align=PP_ALIGN.CENTER)
    add_callout(s, "Use color coding to separate evidence available, partial support, planned evidence, and decision-critical gaps.", 0.72, 5.80, 14.55, 0.45, size=7.2)

    # 10
    s = prs.slides.add_slide(blank); add_base(s, 10, "Patient journey")
    add_rich_heading(s, "Patient journey evidence map", " | Locate evidence needs by care-stage friction")
    donut = s.shapes.add_shape(MSO_SHAPE.DONUT, Inches(2.08), Inches(1.70), Inches(2.65), Inches(2.65))
    set_fill(donut, TEAL_700); no_line(donut)
    add_box(s, 2.08, 1.70, 1.32, 1.32, fill=GRAY_200, transparency=10)
    add_text(s, "Evidence\njourney", 2.65, 2.45, 1.0, 0.35, size=8, color=TEAL_950, bold=True, align=PP_ALIGN.CENTER)
    for label, x, y in [("Diagnosis and referral", 6.15, 1.50), ("Treatment selection", 6.15, 2.15), ("Access and adherence", 6.15, 2.80), ("Long-term outcomes", 6.15, 3.45)]:
        add_light_panel(s, label, ["[Evidence needs by care-stage friction]"], x, y, 7.8, 0.48)

    # 11
    s = prs.slides.add_slide(blank); add_base(s, 11, "Evidence catalog")
    add_rich_heading(s, "Current evidence inventory", " | Completed, ongoing, and planned evidence")
    add_table(s, ["Evidence source", "Type", "Population", "Endpoint / topic", "Status", "Readout", "Owner", "Linked claim"], [
        ["[Pivotal study]", "Clinical", "[Population]", "[OS, PFS, response, safety]", "Completed", "[Date]", "Clinical", "[Claim]"],
        ["[Extension / follow-up]", "Clinical", "[Population]", "[Long-term safety, durability]", "Ongoing", "[Date]", "Clinical", "[Claim]"],
        ["[SLR / NMA]", "HEOR", "[Comparator set]", "[Relative efficacy]", "Completed", "[Date]", "HEOR", "[Access message]"],
        ["[Budget impact model]", "HEOR", "[Market-specific]", "[Cost, utilization]", "Draft", "[Date]", "Access", "[Payer value]"],
        ["[Registry analysis]", "RWE", "[Real-world patients]", "[Treatment patterns, outcomes]", "Gap", "[Needed]", "RWE", "[Real-world value]"],
        ["[Patient preference study]", "Patient evidence", "[Patients / caregivers]", "[Burden, preference, QoL]", "Proposed", "[Date]", "Medical", "[Patient value]"],
    ], 0.72, 1.45, 14.45, 3.18, [1.7, 1.0, 1.4, 2.1, 1.0, 0.9, 0.9, 1.4], 5.2)
    add_callout(s, "Keep the evidence catalog traceable. Each row should connect to a source, owner, status, timing, evidence need, and value claim.", 0.72, 5.05, 14.45, 0.42, size=7)

    # 12
    s = prs.slides.add_slide(blank); add_section(s, 12, "Gap analysis", "and prioritization", "Section divider")

    # 13
    s = prs.slides.add_slide(blank); add_base(s, 13, "Evidence gap matrix")
    add_rich_heading(s, "Evidence gap matrix", " | Prioritize gaps by stakeholder, geography, and timing")
    add_table(s, ["Evidence need", "Stakeholder", "Geography", "Current support", "Gap", "Priority", "Decision consequence"], [
        ["[Comparator effectiveness]", "HTA / payers", "[EU5]", "Indirect evidence only", "No local comparator data", "High", "[May weaken reimbursement case]"],
        ["[QoL and symptom burden]", "Patients / HCPs", "[Global]", "Limited PRO collection", "Insufficient patient-centered support", "High", "[Limits value story and guidelines]"],
        ["[Long-term safety]", "Regulators / HCPs", "[US, EU]", "Trial follow-up ongoing", "Real-world durability unknown", "Medium", "[May constrain confidence]"],
        ["[Economic model inputs]", "HTA / payers", "[Priority markets]", "Partial source support", "Local utility and cost inputs missing", "Medium", "[HTA uncertainty]"],
        ["[Sequencing evidence]", "HCPs / guidelines", "[Global]", "Emerging evidence", "Clear pathway data absent", "Watch", "[May delay adoption]"],
    ], 0.72, 1.45, 14.45, 2.70, [1.8, 1.4, 1.1, 2.1, 2.0, 0.9, 2.1], 5.2)
    add_box(s, 0.72, 6.80, 14.45, 0.24, fill=TEAL_800)
    add_text(s, "The right-side decision consequence column keeps the gap matrix connected to action.", 0.90, 6.83, 14.0, 0.15, size=6.3, color=PAPER, bold=True, align=PP_ALIGN.CENTER)

    # 14
    s = prs.slides.add_slide(blank); add_base(s, 14, "Prioritization")
    add_rich_heading(s, "Gap prioritization scorecard", " | Score impact, urgency, and feasibility")
    labels = ["Comparator RWE", "QoL / symptom burden", "Local utility data", "Publication readiness", "Sequencing evidence"]
    add_text(s, "Gap / tactic", 0.95, 1.48, 2.2, 0.2, size=7.5, color=GRAY_700, bold=True)
    for i, h in enumerate(["Impact", "Urgency", "Feasibility", "Decision value", "Total"]):
        add_text(s, h, 4.1 + i * 1.7, 1.48, 1.1, 0.2, size=7.2, color=GRAY_700, bold=True, align=PP_ALIGN.CENTER)
    for r, label in enumerate(labels):
        y = 1.86 + r * 0.46
        add_text(s, label, 0.95, y, 2.2, 0.18, size=7.0, color=INK)
        for c in range(5):
            add_box(s, 4.15 + c * 1.7, y, 1.18, 0.16, fill=GRAY_200)
            width = [0.95, 0.85, 0.45, 0.75, 1.05][(r + c) % 5]
            add_box(s, 4.15 + c * 1.7, y, width, 0.16, fill=[TEAL_700, MINT, ORANGE, TEAL_800, MINT][(r + c) % 5])
    for i, note in enumerate(["Impact: Does the gap affect access, label, adoption, or guideline decisions?", "Urgency: Is the readout needed before a near-term decision gate?", "Feasibility: Can evidence be generated credibly within timeline and budget?"]):
        add_callout(s, note, 0.72 + i * 4.9, 5.65, 4.55, 0.52, size=6.7)

    # 15
    s = prs.slides.add_slide(blank); add_base(s, 15, "Research questions")
    add_rich_heading(s, "Open research questions", " | Translate priority gaps into PICOT-ready questions")
    add_table(s, ["ORQ", "Population", "Intervention", "Comparator", "Outcome", "Timing", "Decision use"], [
        ["[Does product improve outcomes vs local standard of care?]", "[Eligible patients]", "[Product]", "[SoC]", "[Clinical and economic outcomes]", "[Follow-up]", "[HTA / reimbursement]"],
        ["[Which patients derive the greatest value?]", "[Subgroups]", "[Product]", "[Relevant alternatives]", "[Response, discontinuation, QoL]", "[Milestone]", "[Positioning / guidelines]"],
        ["[What patient outcomes matter most?]", "[Patients / caregivers]", "[Treatment experience]", "[Current care]", "[Preference, burden, HRQoL]", "[Pre-launch]", "[Value story / comms]"],
        ["[What is the budget impact?]", "[Market population]", "[Product adoption]", "[Current mix]", "[Costs, offsets, utilization]", "[1-5 years]", "[Payer discussions]"],
    ], 0.72, 1.45, 14.45, 2.35, [2.2, 1.3, 1.2, 1.3, 2.0, 1.0, 1.5], 5.3)
    add_callout(s, "Use PICOT where applicable. For evidence synthesis, market access, or qualitative work, define the question, decision, stakeholder, and minimum credible evidence standard.", 0.72, 4.25, 14.45, 0.52, size=7.2)

    # 16
    s = prs.slides.add_slide(blank); add_section(s, 16, "Evidence generation", "tactics and roadmap", "Section divider")

    # 17
    s = prs.slides.add_slide(blank); add_base(s, 17, "Tactic menu")
    add_rich_heading(s, "Evidence generation options", " | Choose the right tactic for the decision")
    tactics = [
        ("Clinical studies", "Phase IV, pragmatic trials, extension, post-hoc, subgroup analysis, investigator-initiated research.", TEAL_800),
        ("RWE and data strategy", "Registry, claims, EHR, chart review, external control, prospective observational, site-level feasibility.", MINT),
        ("HEOR and access", "SLR, NMA, economic model, BIA, utility study, payer research, HTA package.", PURPLE),
        ("Patient-centered evidence", "PRO, COA, preference study, qualitative interviews, caregiver burden, adherence evidence.", ORANGE),
        ("Evidence synthesis", "Living literature review, gap assessment, landscape update, comparative evidence refresh.", TEAL_700),
        ("Communications pull-through", "Publication plan, congress plan, field medical enablement, claim readiness review.", GRAY_500),
    ]
    for i, (title, desc, accent) in enumerate(tactics):
        x = 0.72 + (i % 3) * 4.84
        y = 1.48 + (i // 3) * 1.85
        add_box(s, x, y, 4.55, 1.32, fill=PAPER, line=GRAY_200)
        add_box(s, x, y, 4.55, 0.05, fill=accent)
        add_text(s, title, x + 0.14, y + 0.16, 4.2, 0.22, size=8.5, color=TEAL_950, bold=True)
        add_text(s, desc, x + 0.14, y + 0.48, 4.2, 0.42, size=6.7, color=GRAY_700)
        add_box(s, x + 0.14, y + 1.05, 0.45, 0.03, fill=MINT)

    # 18
    s = prs.slides.add_slide(blank); add_base(s, 18, "Study concept")
    add_rich_heading(s, "Study concept one-pager", " | Convert a gap into an executable tactic")
    add_light_panel(s, "Evidence gap and research question", ["Gap: [Decision-critical evidence gap]", "ORQ: [Research question]", "Stakeholder: [Regulator / payer / HCP / patient]", "Decision use: [What this evidence will change]"], 3.25, 1.55, 4.45, 2.55)
    add_line(s, 7.95, 2.85, 8.55, 2.85, color=TEAL_700, width=2)
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.48), Inches(2.75), Inches(0.18), Inches(0.20))
    set_fill(tri, TEAL_700); no_line(tri)
    add_light_panel(s, "Proposed design and execution", ["Design: [Study / analysis type]", "Data source: [Trial, registry, claims, EHR, survey]", "Endpoints: [Outcomes and measures]", "Timing / owner: [Readout, budget, accountable function]"], 8.78, 1.55, 4.45, 2.55)

    # 19
    s = prs.slides.add_slide(blank); add_base(s, 19, "3-5 year roadmap")
    add_rich_heading(s, "Integrated evidence roadmap", " | Sequence workstreams against decision gates")
    add_box(s, 0.72, 1.45, 14.45, 4.30, fill=GRAY_100, radius=True)
    years = ["Y1 H1", "Y1 H2", "Y2 H1", "Y2 H2", "Y3"]
    for i, year in enumerate(years):
        add_text(s, year, 2.45 + i * 2.35, 1.98, 1.2, 0.18, size=7.3, color=GRAY_700, bold=True, align=PP_ALIGN.CENTER)
    lanes = ["Clinical", "RWE", "HEOR", "Patient evidence", "Communications"]
    for r, lane in enumerate(lanes):
        y = 2.62 + r * 0.62
        add_text(s, lane, 1.02, y, 1.5, 0.18, size=7.2, color=TEAL_950, bold=True)
        add_line(s, 1.02, y - 0.20, 14.85, y - 0.20, color=GRAY_200)
        for c in range(5):
            color = [TEAL_700, TEAL_800, MINT, GRAY_300, TEAL_800][(r + c) % 5]
            add_box(s, 2.55 + c * 2.35, y, 1.85, 0.16, fill=color, radius=True)
    for x, y in [(12.78, 2.60), (12.78, 3.85)]:
        d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(0.26), Inches(0.26))
        set_fill(d, ORANGE); no_line(d)
    for i, txt in enumerate(["Teal: active workstream", "Green: readout or synthesis", "Orange diamond: decision gate"]):
        add_callout(s, txt, 0.72 + i * 4.84, 6.05, 4.55, 0.42, size=7.5)

    # 20
    s = prs.slides.add_slide(blank); add_base(s, 20, "Execution tracker")
    add_rich_heading(s, "Activity tracker", " | Manage status, dependencies, and next actions")
    add_table(s, ["Activity", "Gap addressed", "Function", "Owner", "Status", "Budget", "Next action"], [
        ["[Study / analysis A]", "[Gap ID]", "RWE", "[Name]", "In design", "[Range]", "[Protocol review]"],
        ["[SLR / NMA refresh]", "[Gap ID]", "HEOR", "[Name]", "On track", "[Range]", "[Finalize search]"],
        ["[Payer advisory board]", "[Gap ID]", "Access", "[Name]", "At risk", "[Range]", "[Confirm markets]"],
        ["[Patient preference study]", "[Gap ID]", "Medical", "[Name]", "Proposed", "[Range]", "[Approve concept]"],
        ["[Publication pull-through]", "[Gap ID]", "Comms", "[Name]", "On track", "[Range]", "[Congress target]"],
    ], 0.72, 1.45, 14.45, 2.65, [1.8, 1.4, 1.1, 1.1, 1.1, 0.9, 1.9], 5.4)
    add_box(s, 0.72, 6.80, 14.45, 0.24, fill=TEAL_800)
    add_text(s, "Use the tracker as the living operating layer after the initial IEGP is approved.", 0.90, 6.83, 14.0, 0.15, size=6.3, color=PAPER, bold=True, align=PP_ALIGN.CENTER)

    # 21
    s = prs.slides.add_slide(blank); add_section(s, 21, "Functional modules", "and evidence pull-through", "Section divider")

    # 22
    s = prs.slides.add_slide(blank); add_base(s, 22, "HEOR / access")
    add_rich_heading(s, "HEOR and market access evidence", " | Prepare value, reimbursement, and HTA support")
    add_light_panel(s, "Core HEOR work products", ["SLR, targeted literature review, and evidence synthesis", "Network meta-analysis or indirect comparison", "Cost-effectiveness and budget impact models", "Utility, resource use, and burden inputs", "Payer research and value message testing"], 0.72, 1.45, 6.95, 1.75)
    add_light_panel(s, "Market access readiness questions", ["Which markets require local comparator data?", "Which endpoints drive uncertainty in HTA review?", "What economic assumptions need source support?", "What evidence is needed before launch sequencing?", "Which gaps can be closed through synthesis vs new data?"], 8.10, 1.45, 6.95, 1.75)
    add_box(s, 0.72, 3.62, 14.45, 0.38, fill=TEAL_800)
    add_text(s, "Access evidence should be planned early enough to influence submissions, not only explain outcomes after launch.", 1.0, 3.71, 13.9, 0.18, size=7.8, color=PAPER, bold=True, align=PP_ALIGN.CENTER)

    # 23
    s = prs.slides.add_slide(blank); add_base(s, 23, "RWE strategy")
    add_rich_heading(s, "RWE and data strategy", " | Match data sources to evidence use cases")
    add_table(s, ["Use case", "Candidate source", "Endpoint availability", "Population fit", "Bias / limitations", "Feasibility", "Decision use"], [
        ["[Treatment patterns]", "[Claims / EHR]", "High", "Medium", "Missing clinical detail", "High", "[Launch planning]"],
        ["[Comparative outcomes]", "[Registry]", "Medium", "High", "Confounding, sample size", "Medium", "[HTA / HCP confidence]"],
        ["[Safety in practice]", "[EHR / registry]", "High", "High", "Follow-up completeness", "High", "[Medical communication]"],
        ["[External control]", "[Registry + chart]", "Medium", "Medium", "Eligibility alignment", "Needs review", "[Regulatory / access]"],
    ], 0.72, 1.45, 14.45, 2.45, [1.6, 1.6, 1.8, 1.4, 2.1, 1.3, 1.8], 6.0)
    add_callout(s, "Document data-source limitations directly in the plan so downstream users do not overstate evidence strength.", 0.72, 4.35, 14.45, 0.62, size=10)

    # 24
    s = prs.slides.add_slide(blank); add_base(s, 24, "Patient / regulatory")
    add_rich_heading(s, "Patient-centered and regulatory evidence", " | Integrate outcomes that change confidence")
    add_dark_card(s, "Patient-centered evidence", ["PROs, COAs, HRQoL, symptom burden, and functional status.", "Patient preference, treatment burden, caregiver impact, adherence.", "Qualitative interviews to contextualize outcomes and unmet need.", "Digital endpoints or passive measures where credible and validated."], 0.72, 1.55, 6.95, 3.05)
    add_dark_card(s, "Regulatory evidence considerations", ["Confirm commitments, safety follow-up, and real-world acceptability.", "Identify label-expansion, post-marketing, and evidence standards.", "Document where RWE or external controls may support decisions.", "Keep regulatory questions separate from promotional claims."], 8.10, 1.55, 6.95, 3.05)

    # 25
    s = prs.slides.add_slide(blank); add_base(s, 25, "Governance / KPIs")
    add_rich_heading(s, "Governance and impact dashboard", " | Keep the plan current and accountable")
    for i, metric in enumerate([("[x]", "Gaps closed or downgraded this cycle"), ("[x%]", "Activities on track vs planned milestones"), ("[$x]", "Budget committed by workstream"), ("[x]", "Decision gates supported by new evidence")]):
        add_metric(s, metric[0], metric[1], 0.72 + i * 3.68, 1.48, 3.35, 0.95)
    add_light_panel(s, "Governance cadence", ["Core working team: [monthly]", "Cross-functional evidence council: [quarterly]", "Leadership decision forum: [as needed]", "Formal refresh: [annual or trigger-based]"], 0.72, 2.92, 6.95, 1.72)
    add_light_panel(s, "Refresh triggers", ["New clinical readout or safety signal", "Competitor approval, label, data, or access change", "HTA, guideline, or regulatory requirement shift", "Budget change, study delay, or evidence gap reprioritization"], 8.10, 2.92, 6.95, 1.72)

    # 26
    s = prs.slides.add_slide(blank); add_base(s, 26, "Appendix starter")
    add_rich_heading(s, "Appendix worksheet", " | Source log, assumptions, and methods")
    for i in range(9):
        add_line(s, 0.72 + i * 1.80, 1.45, 0.72 + i * 1.80, 6.90, color=GRAY_200, width=0.5)
    for i in range(11):
        add_line(s, 0.72, 1.45 + i * 0.545, 15.10, 1.45 + i * 0.545, color=GRAY_200, width=0.5)
    add_box(s, 1.0, 1.70, 4.75, 1.10, fill=PAPER)
    add_box(s, 1.0, 1.70, 0.07, 1.10, fill=TEAL_700)
    add_text(s, "Suggested appendix content:\nsource log, evidence-quality notes, interview list, workshop participants, scoring rubric, assumptions, limitations, excluded tactics, country adaptation notes, and glossary.", 1.25, 1.86, 4.05, 0.75, size=8.7, color=INK, bold=False)

    return prs


if __name__ == "__main__":
    deck = make_deck()
    deck.save(OUT)
    deck.save(OUT_MAIN)
    print(f"saved {OUT}")
    print(f"saved {OUT_MAIN}")
