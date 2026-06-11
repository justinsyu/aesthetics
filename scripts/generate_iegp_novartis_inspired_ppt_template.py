from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("outputs/iegp_powerpoint_template")
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "IEGP_Novartis_Inspired_Clinical_Briefing_Template.pptx"

PRS = Presentation()
PRS.slide_width = Inches(13.333)
PRS.slide_height = Inches(7.5)
W, H = PRS.slide_width, PRS.slide_height

COLORS = {
    "white": RGBColor(255, 255, 255),
    "navy": RGBColor(0, 37, 96),
    "blue": RGBColor(0, 93, 171),
    "sky": RGBColor(95, 166, 215),
    "orange": RGBColor(239, 101, 35),
    "gray": RGBColor(112, 119, 128),
    "mid": RGBColor(186, 195, 203),
    "panel": RGBColor(243, 246, 249),
    "line": RGBColor(218, 224, 230),
    "dark": RGBColor(32, 38, 46),
    "green": RGBColor(88, 156, 96),
}


def c(name):
    return COLORS[name]


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def line(shape, color, width=0.7):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def bg():
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    fill(rect, c("white"))
    rect.line.fill.background()
    return slide


def text(slide, x, y, w, h, value, size=12, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.text = value
    if align:
        p.alignment = align
    for r in p.runs:
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or c("dark")
    return box


def body(slide, x, y, w, h, lines, size=10, color=None, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, line_value in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_value
        p.space_after = Pt(4)
        if bullet:
            p._p.get_or_add_pPr().set("marL", "171450")
            p._p.get_or_add_pPr().set("indent", "-114300")
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(size)
            r.font.color.rgb = color or c("dark")
    return box


def rect(slide, x, y, w, h, color, border=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, c(color))
    if border:
        line(shp, c(border), 0.7)
    else:
        shp.line.fill.background()
    return shp


def nav(slide, active="Overview"):
    rect(slide, 0, 0, 0.82, 7.5, "panel")
    sections = ["Overview", "Context", "Needs", "Gaps", "Roadmap", "Control"]
    for i, section in enumerate(sections):
        y = 0.65 + i * 0.72
        if section == active:
            rect(slide, 0.0, y - 0.05, 0.08, 0.38, "orange")
            color = c("blue")
            bold = True
        else:
            color = c("gray")
            bold = False
        text(slide, 0.15, y, 0.55, 0.13, section, size=5.6, bold=bold, color=color)
    text(slide, 0.2, 7.08, 0.35, 0.12, f"{len(PRS.slides):02d}", size=5.6, color=c("gray"), align=PP_ALIGN.CENTER)


def footer(slide, active):
    text(slide, 1.0, 7.08, 4.8, 0.16, f"IEGP | {active} | [Product / indication]", size=6.6, color=c("gray"))
    text(slide, 11.65, 7.08, 0.85, 0.16, f"{len(PRS.slides):02d}", size=6.6, color=c("gray"), align=PP_ALIGN.RIGHT)


def title(slide, headline, active, sub=None):
    nav(slide, active)
    text(slide, 1.02, 0.34, 10.6, 0.48, headline, size=22, bold=True, color=c("navy"))
    if sub:
        text(slide, 1.04, 0.92, 9.9, 0.22, sub, size=9.2, color=c("gray"))
    rect(slide, 1.02, 1.26, 11.5, 0.012, "line")
    footer(slide, active)


def organic_mask(slide, x, y, w, h, color="panel"):
    # Editable rectangles create a soft Novartis-like image field without copying brand assets.
    rect(slide, x, y, w, h, color)
    rect(slide, x + 0.18, y + 0.4, w * 0.42, h * 0.72, "white")
    rect(slide, x + w * 0.55, y + h * 0.55, w * 0.34, h * 0.22, "orange")


def panel(slide, x, y, w, h, header=None, accent="blue"):
    rect(slide, x, y, w, h, "panel", "line")
    rect(slide, x, y, 0.07, h, accent)
    if header:
        text(slide, x + 0.22, y + 0.14, w - 0.4, 0.18, header, size=8.4, bold=True, color=c("navy"))


def callout(slide, x, y, w, h, header, lines, accent="orange"):
    panel(slide, x, y, w, h, header, accent)
    if h < 0.85:
        text(slide, x + 0.22, y + 0.32, w - 0.42, 0.16, "   ".join(lines), size=8.6, color=c("dark"))
    else:
        body(slide, x + 0.22, y + 0.5, w - 0.42, h - 0.62, lines, size=8.6, color=c("dark"))


def table(slide, x, y, w, h, rows, cols, headers, widths=None, header_color="navy", font_size=6.7):
    shp = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shp.table
    if widths:
        for i, width in enumerate(widths):
            tbl.columns[i].width = Inches(width)
    for r in range(rows):
        for col in range(cols):
            cell = tbl.cell(r, col)
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = c(header_color if r == 0 else ("white" if r % 2 else "panel"))
            cell.text = headers[col] if r == 0 else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0
                    run.font.color.rgb = c("white") if r == 0 else c("dark")
    return shp


def cover():
    slide = bg()
    organic_mask(slide, 6.7, 0.0, 6.63, 7.5, "panel")
    rect(slide, 6.95, 0.85, 5.5, 4.15, "blue")
    # simple building/lab-like abstract image field
    for i in range(8):
        rect(slide, 7.25 + i * 0.56, 1.2, 0.28, 3.4, "sky")
        rect(slide, 7.39 + i * 0.56, 1.2, 0.08, 3.4, "white")
    rect(slide, 7.1, 5.22, 3.0, 1.05, "orange")
    text(slide, 0.9, 0.78, 3.2, 0.15, "INTEGRATED EVIDENCE GENERATION PLAN", size=6.7, bold=True, color=c("blue"))
    text(slide, 0.88, 1.45, 5.25, 0.92, "Clinical evidence briefing template", size=32, bold=True, color=c("navy"))
    text(slide, 0.92, 2.72, 4.85, 0.42, "A Novartis-inspired, data-first PowerPoint system for evidence strategy, gap prioritization, and roadmap governance.", size=12, color=c("gray"))
    rect(slide, 0.92, 4.2, 3.5, 0.08, "orange")
    text(slide, 0.92, 5.35, 5.2, 0.22, "[Product] | [Indication] | [Region]", size=13, bold=True, color=c("navy"))
    text(slide, 0.92, 5.72, 4.2, 0.16, "Owner: [team]  |  Version: [v0.1]  |  Date: [date]", size=7.5, color=c("gray"))


def tokens():
    slide = bg()
    title(slide, "Brand tokens can be swapped without changing the layout system", "Overview", "Default colors are Novartis-inspired, not copied brand assets.")
    items = [("Primary Brand", "navy"), ("Secondary Brand", "blue"), ("Signal / highlight", "orange"), ("Comparator gray", "gray"), ("Panel gray", "panel"), ("Risk / gap", "orange")]
    for i, (label, col) in enumerate(items):
        x = 1.02 + (i % 3) * 3.65
        y = 2.0 + (i // 3) * 1.35
        rect(slide, x, y, 0.78, 0.78, col)
        text(slide, x + 0.95, y + 0.12, 1.8, 0.16, label, size=8.6, bold=True, color=c("navy"))
        text(slide, x + 0.95, y + 0.42, 1.85, 0.14, f"RGB {c(col)[0]}, {c(col)[1]}, {c(col)[2]}", size=6.4, color=c("gray"))
    callout(slide, 1.02, 5.3, 10.95, 0.78, "Adaptation rule", ["Keep white space, deep primary color, selective highlight color, and one dominant evidence object per slide."], "orange")


def section(name, active, note, color="blue"):
    slide = bg()
    nav(slide, active)
    organic_mask(slide, 7.0, 0, 6.33, 7.5, "panel")
    rect(slide, 7.35, 1.1, 4.5, 4.3, color)
    for i in range(6):
        rect(slide, 7.75 + i * 0.56, 1.45, 0.2, 3.5, "white")
        rect(slide, 7.95 + i * 0.56, 1.45, 0.1, 3.5, "sky")
    text(slide, 1.02, 1.35, 5.35, 0.75, name, size=34, bold=True, color=c("navy"))
    text(slide, 1.05, 2.45, 4.9, 0.35, note, size=13, color=c("gray"))
    rect(slide, 1.05, 3.45, 3.4, 0.08, "orange")
    footer(slide, active)


def executive():
    slide = bg()
    title(slide, "Evidence priorities focus decision makers on the highest-value gaps", "Overview", "Use this slide as the senior-team briefing view.")
    for i, (h, col) in enumerate([("Priority 1", "blue"), ("Priority 2", "navy"), ("Priority 3", "orange")]):
        x = 1.02 + i * 3.72
        panel(slide, x, 1.8, 3.25, 3.7, h, col)
        text(slide, x + 0.25, 2.35, 2.55, 0.22, "[Strategic evidence priority]", size=13, bold=True, color=c("navy"))
        body(slide, x + 0.25, 2.9, 2.6, 1.65, ["Gap: [top gap]", "Action: [study / analysis]", "Decision: [reg/access/adoption]", "Timing: [window]"], size=8.8)
        rect(slide, x + 0.25, 4.95, 1.8, 0.08, "orange")
    callout(slide, 1.02, 5.92, 10.9, 0.6, "Decision ask", ["[Approve roadmap] | [Fund priority activities] | [Confirm owners and refresh cadence]"], "orange")


def decisions():
    slide = bg()
    title(slide, "Decisions needed are tied to evidence readiness and timing", "Overview")
    shp = table(slide, 1.02, 1.85, 10.95, 4.8, 7, 5, ["Decision", "Recommendation", "Evidence basis", "Owner", "Deadline"], [2.4, 2.0, 3.8, 1.35, 1.4])
    rows = [["Approve priority gaps", "[approve/revise]", "[risk and value rationale]", "[owner]", "[date]"], ["Fund evidence activity", "[approve/defer]", "[decision value]", "[owner]", "[date]"], ["Select study design", "[option A/B]", "[credibility and feasibility]", "[owner]", "[date]"], ["Resolve local adaptation", "[global/local split]", "[reuse vs local gap]", "[owner]", "[date]"], ["Confirm governance", "[RACI/cadence]", "[execution accountability]", "[owner]", "[date]"], ["Other", "[recommendation]", "[basis]", "[owner]", "[date]"]]
    for r, row in enumerate(rows, 1):
        for col, value in enumerate(row):
            shp.table.cell(r, col).text = value


def asset():
    slide = bg()
    title(slide, "Asset snapshot frames evidence needs before study planning begins", "Context")
    items = [("Asset / mechanism", "[name, modality, MoA]"), ("Population", "[indication, line, biomarker]"), ("Lifecycle stage", "[pre-launch / launch / post-launch]"), ("Geography", "[global and priority markets]"), ("TPP / TVP", "[clinical and value targets]"), ("Key milestones", "[readouts, submissions, HTA, launch]")]
    for i, (h, b) in enumerate(items):
        x = 1.02 + (i % 3) * 3.62
        y = 1.85 + (i // 3) * 1.65
        panel(slide, x, y, 3.08, 1.08, h, ["blue", "navy", "orange"][i % 3])
        text(slide, x + 0.24, y + 0.54, 2.45, 0.2, b, size=9.5, color=c("gray"))


def value_story():
    slide = bg()
    title(slide, "Value messages trace to evidence requirements, support, gaps, and tactics", "Context")
    labels = ["Value driver", "Evidence required", "Current support", "Gap", "Tactic"]
    for i, label in enumerate(labels):
        x = 1.12 + i * 2.16
        rect(slide, x, 2.08, 1.55, 0.35, ["blue", "navy", "gray", "orange", "blue"][i])
        text(slide, x, 2.65, 1.45, 0.18, label, size=9.0, bold=True, color=c("navy"))
        panel(slide, x, 3.12, 1.55, 1.2, None, ["blue", "navy", "gray", "orange", "blue"][i])
        text(slide, x + 0.16, 3.48, 1.1, 0.22, "[content]", size=8.4, color=c("gray"))
        if i < 4:
            rect(slide, x + 1.62, 3.715, 0.42, 0.012, "line")
    callout(slide, 1.02, 5.6, 10.9, 0.65, "Narrative test", ["Can every value message be supported by source-backed evidence or a funded activity?"], "orange")


def stakeholder():
    slide = bg()
    title(slide, "Stakeholder needs define the evidence standard for each decision", "Needs")
    positions = [("Regulators", 1.02, 1.85, "blue"), ("HTA / payers", 4.65, 1.85, "navy"), ("HCPs / KOLs", 8.28, 1.85, "sky"), ("Patients", 1.02, 4.08, "orange"), ("Policy / systems", 4.65, 4.08, "blue"), ("Internal teams", 8.28, 4.08, "navy")]
    for name, x, y, col in positions:
        panel(slide, x, y, 3.05, 1.55, name, col)
        body(slide, x + 0.24, y + 0.56, 2.45, 0.62, ["Decision: [ ]", "Evidence standard: [ ]", "Timing: [ ]"], size=8.0)


def inventory():
    slide = bg()
    title(slide, "Current evidence inventory separates usable support from unresolved limits", "Needs")
    shp = table(slide, 0.95, 1.85, 11.15, 4.9, 9, 6, ["Source", "Status", "Population", "Endpoint / output", "Stakeholder", "Limits"], [1.9, 1.15, 1.9, 2.7, 1.55, 1.95])
    rows = ["Clinical trials", "RWE", "HEOR model", "SLR / ITC", "COA / PRO", "Epidemiology", "Publications", "Other"]
    for r, label in enumerate(rows, 1):
        vals = [label, "[status]", "[population]", "[output]", "[stakeholder]", "[quality / usability limits]"]
        for col, value in enumerate(vals):
            shp.table.cell(r, col).text = value


def gap_matrix():
    slide = bg()
    title(slide, "Evidence gaps are prioritized by decision consequence and actionability", "Gaps")
    shp = table(slide, 0.95, 1.85, 9.6, 4.9, 8, 5, ["Gap", "Decision", "Current evidence", "Priority", "Action"], [2.1, 2.1, 2.45, 0.9, 2.05], header_color="navy")
    for r in range(1, 8):
        vals = [f"G{r}: [gap]", "[decision]", "[evidence]", "[P1]", "[tactic]"]
        for col, value in enumerate(vals):
            shp.table.cell(r, col).text = value
    callout(slide, 10.85, 1.85, 1.35, 4.9, "Decision consequence", ["[label risk]", "[HTA risk]", "[slow uptake]", "[unfunded lifecycle]"], "orange")


def prioritization():
    slide = bg()
    title(slide, "Priority scoring converts evidence gaps into funded roadmap choices", "Gaps")
    metrics = [("5", "Strategic impact", "blue"), ("4", "Urgency", "navy"), ("3", "Feasibility", "orange"), ("2", "Reuse potential", "gray")]
    for i, (num, label, col) in enumerate(metrics):
        x = 1.08 + i * 2.7
        text(slide, x, 2.0, 1.3, 0.45, num, size=27, bold=True, color=c(col), align=PP_ALIGN.CENTER)
        rect(slide, x + 0.1, 2.67, 1.05, 0.06, col)
        text(slide, x - 0.25, 2.9, 1.8, 0.2, label, size=8, color=c("gray"), align=PP_ALIGN.CENTER)
    shp = table(slide, 1.02, 4.1, 10.95, 2.25, 5, 5, ["Gap", "Impact", "Urgency", "Feasibility", "Decision"], [2.8, 1.5, 1.5, 1.75, 3.4], header_color="navy")
    for r in range(1, 5):
        vals = [f"G{r}: [gap]", "[1-5]", "[1-5]", "[1-5]", "[fund / revise / hold]"]
        for col, value in enumerate(vals):
            shp.table.cell(r, col).text = value


def roadmap():
    slide = bg()
    title(slide, "Integrated roadmap aligns workstreams to evidence readouts and decision gates", "Roadmap")
    quarters = ["Y1 Q1", "Q2", "Q3", "Q4", "Y2 Q1", "Q2", "Q3", "Q4"]
    for i, q in enumerate(quarters):
        text(slide, 3.38 + i * 0.94, 1.85, 0.7, 0.14, q, size=6.2, bold=True, color=c("gray"), align=PP_ALIGN.CENTER)
    lanes = [("Clinical / regulatory", "blue"), ("RWE", "navy"), ("HEOR / access", "sky"), ("Patient / COA", "orange"), ("Publications", "blue")]
    for i, (lane, col) in enumerate(lanes):
        y = 2.42 + i * 0.72
        text(slide, 1.02, y + 0.07, 1.95, 0.15, lane, size=7.2, bold=True, color=c("navy"))
        rect(slide, 3.1, y + 0.17, 8.8, 0.012, "line")
        rect(slide, 3.35 + (i % 3) * 0.5, y, 1.55 + (i % 2) * 0.5, 0.3, col)
        text(slide, 3.48 + (i % 3) * 0.5, y + 0.08, 1.05, 0.08, "[activity]", size=5.5, bold=True, color=c("white"), align=PP_ALIGN.CENTER)
        rect(slide, 5.78 + i * 0.7, y - 0.03, 0.15, 0.36, "orange")
    callout(slide, 1.02, 6.22, 10.9, 0.5, "Decision gates", ["[protocol] | [data cut] | [abstract] | [regulatory] | [HTA] | [launch]"], "orange")


def activity():
    slide = bg()
    title(slide, "Evidence activity tracker functions as the working source of truth", "Roadmap")
    shp = table(slide, 0.82, 1.85, 11.55, 4.9, 9, 8, ["ID", "Activity", "Gap", "Method / data", "Owner", "Timing", "Status", "Priority"], [0.5, 1.8, 0.7, 2.45, 1.05, 1.15, 1.6, 0.9], font_size=6.0)
    for r in range(1, 9):
        vals = [f"A{r}", "[study / analysis]", "G[ ]", "[design / source]", "[owner]", "[start-readout]", "[idea/planned/active]", "[P1]"]
        for col, value in enumerate(vals):
            shp.table.cell(r, col).text = value


def study():
    slide = bg()
    title(slide, "Study concept one-pager keeps design, decision use, and pull-through visible", "Roadmap")
    items = [("Study / activity", "[name and ID]", "blue"), ("Evidence need", "[gap, ORQ, stakeholder]", "navy"), ("Objective", "[objective and hypothesis]", "sky"), ("Design", "[population, comparator, endpoint, data]", "blue"), ("Feasibility", "[cost, complexity, dependency]", "orange"), ("Pull-through", "[publication, dossier, claims]", "navy")]
    for i, (h, b, col) in enumerate(items):
        x = 1.02 + (i % 3) * 3.62
        y = 1.85 + (i // 3) * 1.65
        panel(slide, x, y, 3.08, 1.08, h, col)
        text(slide, x + 0.24, y + 0.54, 2.45, 0.2, b, size=9.0, color=c("gray"))
    callout(slide, 1.02, 6.0, 10.9, 0.5, "Decision", ["[approve / revise / defer]    Next step: [action]    Date: [ ]"], "orange")


def workstreams():
    slide = bg()
    title(slide, "Specialized workstreams remain distinct but connected to the integrated plan", "Roadmap")
    work = [("HEOR / access", ["SLR / ITC", "economic model", "budget impact", "payer research"], "navy"), ("RWE / data", ["retrospective database", "prospective study", "registry", "external control"], "blue"), ("Patient-centered", ["COA / PRO", "HRQoL", "preference", "caregiver burden"], "orange"), ("Medical communications", ["abstract", "manuscript", "field medical", "claims matrix"], "sky")]
    for i, (h, lines, col) in enumerate(work):
        x = 1.02 + (i % 2) * 5.38
        y = 1.95 + (i // 2) * 2.0
        panel(slide, x, y, 4.85, 1.45, h, col)
        body(slide, x + 0.24, y + 0.52, 4.0, 0.55, lines, size=8.5)


def governance():
    slide = bg()
    title(slide, "Governance cadence keeps the plan current after each evidence or landscape event", "Control")
    shp = table(slide, 1.02, 1.85, 5.95, 4.65, 7, 4, ["Forum", "Purpose", "Cadence", "Rights"], [1.35, 2.25, 1.1, 1.25], font_size=6.5)
    rows = [["Core team", "[maintain plan]", "[monthly]", "[recommend]"], ["Evidence council", "[prioritize/fund]", "[quarterly]", "[approve]"], ["Global-local", "[adapt needs]", "[quarterly]", "[recommend]"], ["Leadership", "[trade-offs]", "[semiannual]", "[approve]"], ["Publications", "[pull-through]", "[monthly]", "[recommend]"], ["Other", "[purpose]", "[cadence]", "[rights]"]]
    for r, row in enumerate(rows, 1):
        for col, value in enumerate(row):
            shp.table.cell(r, col).text = value
    callout(slide, 7.42, 1.85, 4.6, 4.65, "Update triggers", ["Clinical readout", "Regulatory feedback", "HTA / payer change", "Competitor approval or data", "Guideline update", "Budget or feasibility change"], "orange")


def kpis():
    slide = bg()
    title(slide, "Impact dashboard measures plan execution and evidence pull-through", "Control")
    metrics = [("72%", "Gap closure", "blue"), ("18/24", "Milestones on track", "navy"), ("6", "Markets reusing evidence", "orange"), ("12", "Pull-through outputs", "sky")]
    for i, (num, label, col) in enumerate(metrics):
        x = 1.05 + i * 2.72
        text(slide, x, 2.0, 1.45, 0.42, num, size=24, bold=True, color=c(col), align=PP_ALIGN.CENTER)
        rect(slide, x + 0.18, 2.62, 1.05, 0.06, col)
        text(slide, x - 0.15, 2.85, 1.75, 0.2, label, size=7.5, color=c("gray"), align=PP_ALIGN.CENTER)
    shp = table(slide, 1.02, 4.05, 10.95, 2.25, 5, 5, ["KPI", "Definition", "Baseline", "Target", "Owner"], [1.9, 4.1, 1.45, 1.45, 2.05], font_size=6.9)
    for r in range(1, 5):
        vals = ["[KPI]", "[definition]", "[baseline]", "[target]", "[owner]"]
        for col, value in enumerate(vals):
            shp.table.cell(r, col).text = value


def risk():
    slide = bg()
    title(slide, "Risks and dependencies are managed as evidence-design choices", "Control")
    shp = table(slide, 0.95, 1.85, 11.15, 4.9, 8, 5, ["Risk / dependency", "Activity", "Impact", "Mitigation", "Decision needed"], [2.55, 1.75, 1.15, 3.45, 2.25], header_color="navy")
    for r in range(1, 8):
        vals = ["[risk]", "[activity]", "[H/M/L]", "[mitigation or alternate evidence path]", "[decision / owner]"]
        for col, value in enumerate(vals):
            shp.table.cell(r, col).text = value


def appendix():
    slide = bg()
    title(slide, "Appendix worksheets preserve traceability without crowding briefing slides", "Control")
    items = [("Regional adaptation", "local needs, reusable global evidence, local gaps, owner", "blue"), ("Evidence source log", "source, date, owner, confidence, limitation, claim link", "navy"), ("Methodology", "inputs, workshops, scoring rubric, assumptions, exclusions", "orange"), ("Glossary", "IEGP/IEP, ORQ, PICOT, RWE, HEOR, HTA, COA/PRO, IIR", "sky")]
    for i, (h, b, col) in enumerate(items):
        x = 1.02 + (i % 2) * 5.38
        y = 1.95 + (i // 2) * 1.85
        panel(slide, x, y, 4.85, 1.25, h, col)
        text(slide, x + 0.24, y + 0.58, 4.1, 0.24, b, size=8.7, color=c("gray"))


slides = [
    cover,
    tokens,
    lambda: section("Overview", "Overview", "Summarize the evidence priorities and decisions needed now.", "blue"),
    executive,
    decisions,
    lambda: section("Asset context", "Context", "Frame the product, value story, and stakeholder requirements.", "navy"),
    asset,
    value_story,
    lambda: section("Evidence needs", "Needs", "Map stakeholder evidence standards and the current evidence base.", "blue"),
    stakeholder,
    inventory,
    lambda: section("Gap analysis", "Gaps", "Convert evidence needs into prioritized, actionable gaps.", "navy"),
    gap_matrix,
    prioritization,
    lambda: section("Roadmap", "Roadmap", "Translate priority gaps into activities, readouts, and pull-through.", "blue"),
    roadmap,
    activity,
    study,
    workstreams,
    lambda: section("Control", "Control", "Operate the IEGP as a living governance and impact system.", "navy"),
    governance,
    kpis,
    risk,
    appendix,
]

for make in slides:
    make()

PRS.save(OUT_FILE)
print(OUT_FILE.resolve())
