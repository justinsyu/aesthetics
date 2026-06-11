from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("outputs/iegp_powerpoint_template")
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "IEGP_Novartis_Adapted_Slide_System_Template.pptx"

PRS = Presentation()
PRS.slide_width = Inches(13.333)
PRS.slide_height = Inches(7.5)
W, H = PRS.slide_width, PRS.slide_height


COLORS = {
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
    "ink": RGBColor(24, 24, 24),
    "gray": RGBColor(112, 112, 112),
    "mid": RGBColor(169, 169, 169),
    "line": RGBColor(222, 222, 222),
    "panel": RGBColor(244, 244, 244),
    "nav": RGBColor(250, 250, 250),
    "blue": RGBColor(0, 101, 164),
    "deep_blue": RGBColor(0, 64, 103),
    "sky": RGBColor(82, 146, 223),
    "orange": RGBColor(255, 82, 18),
    "amber": RGBColor(242, 153, 74),
    "peach": RGBColor(248, 204, 187),
    "lavender": RGBColor(219, 219, 234),
    "dark_pill": RGBColor(59, 59, 59),
}


SECTIONS = [
    "Content",
    "Briefing",
    "Context",
    "Evidence needs",
    "Gap analysis",
    "Roadmap",
    "Governance",
    "Appendix",
]


def c(name):
    return COLORS[name]


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def line(shape, color, width=0.6):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def blank():
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    fill(bg, c("white"))
    bg.line.fill.background()
    return slide


def txt(slide, x, y, w, h, value, size=12, bold=False, color=None, align=None, font="Arial"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.01)
    tf.margin_right = Inches(0.01)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.text = value
    if align:
        p.alignment = align
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or c("ink")
    return box


def body(slide, x, y, w, h, lines, size=9.0, color=None, bullet=False, font="Arial"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, value in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = value
        p.space_after = Pt(4)
        if bullet:
            p._p.get_or_add_pPr().set("marL", "171450")
            p._p.get_or_add_pPr().set("indent", "-114300")
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color or c("ink")
    return box


def rect(slide, x, y, w, h, color, outline=None, width=0.6):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, c(color))
    if outline:
        line(shp, c(outline), width)
    else:
        shp.line.fill.background()
    return shp


def pill(slide, x, y, w, h, label, color="dark_pill", font_size=7.5):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, c(color))
    shp.line.fill.background()
    try:
        shp.adjustments[0] = 0.18
    except Exception:
        pass
    txt(slide, x + 0.06, y + 0.065, w - 0.12, h - 0.08, label.upper(), size=font_size, color=c("white"), align=PP_ALIGN.CENTER)
    return shp


def soft_gradient(slide, x, y, w, h, steps=32):
    for i in range(steps):
        t = i / max(steps - 1, 1)
        r = int(232 + (255 - 232) * t)
        g = int(239 + (196 - 239) * t)
        b = int(248 + (176 - 248) * t)
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + w * i / steps), Inches(y), Inches(w / steps + 0.01), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(r, g, b)
        shp.line.fill.background()


def left_nav(slide, active):
    rect(slide, 0, 0, 1.72, 7.5, "nav", "line", 0.25)
    circ1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.28), Inches(0.27), Inches(0.21), Inches(0.21))
    circ2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.55), Inches(0.27), Inches(0.21), Inches(0.21))
    for circ, label in [(circ1, "<"), (circ2, ">")]:
        circ.fill.background()
        line(circ, c("black"), 0.8)
        txt(slide, circ.left / 914400 + 0.045, circ.top / 914400 + 0.025, 0.12, 0.07, label, size=5.5, align=PP_ALIGN.CENTER)

    txt(slide, 0.28, 0.82, 0.8, 0.18, "Content", size=9, bold=True)
    body(slide, 0.28, 1.06, 1.0, 0.28, ["Click below to", "navigate through", "the document"], size=5.7, color=c("ink"))

    y0 = 1.62
    item_h = 0.49
    for i, label in enumerate(SECTIONS[1:]):
        y = y0 + i * item_h
        if label == active:
            rect(slide, 0, y, 0.12, item_h, "orange")
            soft_gradient(slide, 0.12, y, 0.24, item_h, 18)
            rect(slide, 0.36, y, 1.36, item_h, "panel")
            txt(slide, 0.28, y + 0.17, 1.18, 0.1, label, size=5.7, bold=True)
        else:
            rect(slide, 0, y, 1.72, item_h, "white", "line", 0.2)
            txt(slide, 0.28, y + 0.17, 1.18, 0.1, label, size=5.5, color=c("ink"))


def footer(slide):
    txt(slide, 10.2, 7.12, 1.8, 0.12, "IEGP template | [date]", size=5.5, color=c("gray"), align=PP_ALIGN.RIGHT)
    txt(slide, 12.18, 7.12, 0.25, 0.12, f"{len(PRS.slides)}", size=6.0, bold=True, color=c("ink"), align=PP_ALIGN.RIGHT)


def page(slide_title, section, section_pill=None, subtitle=None):
    slide = blank()
    left_nav(slide, section)
    if section_pill:
        pill(slide, 11.76, 0.35, 1.28, 0.25, section_pill)
    txt(slide, 2.20, 0.82, 9.8, 0.55, slide_title, size=24, bold=True)
    if subtitle:
        txt(slide, 2.20, 1.42, 9.0, 0.18, subtitle, size=8.5, color=c("gray"))
    rect(slide, 2.20, 1.72, 10.8, 0.01, "line")
    footer(slide)
    return slide


def panel(slide, x, y, w, h, header=None, accent="orange", fill_name="panel"):
    rect(slide, x, y, w, h, fill_name, "line", 0.4)
    rect(slide, x, y, 0.09, h, accent)
    if header:
        txt(slide, x + 0.25, y + 0.16, w - 0.4, 0.16, header, size=8.6, bold=True)


def table(slide, x, y, w, h, rows, cols, headers, widths=None, header_fill="deep_blue", font_size=6.2):
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    if widths:
        for i, width in enumerate(widths):
            tbl.columns[i].width = Inches(width)
    for r in range(rows):
        for col in range(cols):
            cell = tbl.cell(r, col)
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = c(header_fill if r == 0 else "white")
            cell.text = headers[col] if r == 0 else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0
                    run.font.color.rgb = c("white") if r == 0 else c("ink")
    return shape


def image_window(slide, x, y, w, h, label="[Insert relevant product, patient, lab, or market image]"):
    soft_gradient(slide, x, y, w, h, 28)
    rect(slide, x, y, w, h, "white", "line", 0.4).fill.transparency = 100000
    rect(slide, x + w - 0.15, y + h - 0.48, 0.15, 0.48, "orange")
    txt(slide, x + 0.22, y + h - 0.42, w - 0.55, 0.14, label, size=6.5, color=c("gray"))


def metric(slide, x, y, number, label, color="orange"):
    txt(slide, x, y, 1.5, 0.42, number, size=26, bold=True, color=c(color))
    txt(slide, x, y + 0.52, 2.1, 0.14, label, size=7.3, bold=True)
    rect(slide, x, y + 0.77, 1.0, 0.035, color)


def cover():
    slide = blank()
    left_nav(slide, "Briefing")
    txt(slide, 2.20, 0.78, 3.2, 0.25, "Investor-style evidence planning template", size=10, bold=True)
    soft_gradient(slide, 2.20, 2.45, 4.72, 3.85, 42)
    image_window(slide, 6.92, 2.45, 5.75, 3.85, "[Insert brand-neutral hero image]")
    rect(slide, 6.78, 1.78, 0.16, 1.15, "orange")
    txt(slide, 2.20, 3.42, 4.55, 1.2, "Integrated Evidence Generation Plan", size=36, bold=True)
    txt(slide, 2.24, 5.55, 3.7, 0.18, "[Product] | [Indication] | [Region]", size=10.5, bold=True)
    txt(slide, 2.24, 5.92, 3.7, 0.14, "Owner: [team]  |  Version: [v0.1]  |  Date: [date]", size=6.4, color=c("gray"))
    txt(slide, 2.22, 7.03, 2.0, 0.18, "[Company wordmark]", size=16, bold=True)


def tokens():
    slide = page("Brand adaptation controls", "Briefing", "SETUP", "Swap these tokens before building a company-specific deck.")
    labels = [
        ("Primary blue", "blue"),
        ("Deep evidence blue", "deep_blue"),
        ("Signal orange", "orange"),
        ("Support sky", "sky"),
        ("Soft panel", "panel"),
        ("Neutral ink", "ink"),
    ]
    for i, (label, col) in enumerate(labels):
        x = 2.25 + (i % 3) * 3.35
        y = 2.2 + (i // 3) * 1.4
        rect(slide, x, y, 0.72, 0.72, col)
        txt(slide, x + 0.9, y + 0.15, 1.6, 0.14, label, size=8.5, bold=True)
        rgb = c(col)
        txt(slide, x + 0.9, y + 0.43, 1.7, 0.12, f"RGB {rgb[0]}, {rgb[1]}, {rgb[2]}", size=6.2, color=c("gray"))
    panel(slide, 2.25, 5.25, 9.8, 0.92, "Adaptation rule", "orange", "white")
    txt(slide, 2.55, 5.58, 9.0, 0.18, "Keep the left navigation, light evidence panels, direct labels, and gradient edge accents. Replace colors, imagery, and wordmark only.", size=10.2, bold=True)


def section_divider(name, section, note):
    slide = blank()
    left_nav(slide, section)
    image_window(slide, 2.08, 0.55, 10.9, 3.0, "[Optional section image]")
    rect(slide, 8.98, 3.08, 0.18, 0.92, "orange")
    txt(slide, 2.18, 4.02, 7.6, 0.62, name, size=34, bold=True)
    txt(slide, 2.20, 4.85, 7.2, 0.28, note, size=11.5, color=c("ink"))
    footer(slide)


def content_map():
    slide = page("Content", "Briefing", None, "Click-style navigation adapted into an editable planning overview.")
    for i, sec in enumerate(SECTIONS[1:]):
        y = 2.04 + i * 0.56
        rect(slide, 2.26, y, 9.3, 0.46, "white", "line", 0.35)
        rect(slide, 2.26, y, 0.08, 0.46, "orange" if i in [0, 4] else "blue")
        txt(slide, 2.52, y + 0.15, 2.2, 0.12, sec, size=7.8, bold=True)
        txt(slide, 5.15, y + 0.15, 5.9, 0.12, "[section purpose / key pages]", size=7.2, color=c("gray"))


def executive():
    slide = page("Executive evidence priorities", "Briefing", "BRIEFING", "Summarize the evidence decisions leadership needs to make now.")
    for i, (num, label, color) in enumerate([("01", "Evidence priority", "orange"), ("02", "Decision impact", "blue"), ("03", "Timing risk", "orange")]):
        x = 2.25 + i * 3.3
        panel(slide, x, 2.18, 2.72, 2.85, label, color, "panel")
        txt(slide, x + 0.24, 2.72, 0.7, 0.35, num, size=23, bold=True, color=c(color))
        body(slide, x + 0.24, 3.3, 2.1, 1.1, ["Gap: [top gap]", "Action: [study / analysis]", "Decision use: [launch/access]"], size=8.0)
    panel(slide, 2.25, 5.65, 9.95, 0.72, "Decision ask", "orange", "white")
    txt(slide, 2.55, 5.92, 9.1, 0.16, "[Approve roadmap]  |  [Fund priority activities]  |  [Confirm owners and refresh cadence]", size=10.4, bold=True)


def decisions():
    slide = page("Decisions needed", "Briefing", "BRIEFING", "Move from evidence planning to accountable choices.")
    shp = table(slide, 2.20, 2.0, 10.55, 4.65, 7, 5, ["Decision", "Recommendation", "Rationale", "Owner", "Deadline"], [2.1, 1.85, 3.55, 1.35, 1.7], "deep_blue", 6.7)
    rows = [
        ["Approve priority gaps", "[approve / revise]", "[value and risk basis]", "[owner]", "[date]"],
        ["Fund activity", "[approve / defer]", "[decision value]", "[owner]", "[date]"],
        ["Select study design", "[option A/B]", "[credibility and feasibility]", "[owner]", "[date]"],
        ["Resolve local needs", "[global/local split]", "[reuse vs local gap]", "[owner]", "[date]"],
        ["Confirm governance", "[RACI/cadence]", "[accountability]", "[owner]", "[date]"],
        ["Other", "[recommendation]", "[basis]", "[owner]", "[date]"],
    ]
    for r, row in enumerate(rows, 1):
        for col, val in enumerate(row):
            shp.table.cell(r, col).text = val


def asset_snapshot():
    slide = page("Asset snapshot", "Context", "CONTEXT", "Frame the asset before defining stakeholder evidence needs.")
    image_window(slide, 8.25, 2.05, 3.85, 3.65, "[Insert product or disease image]")
    fields = [
        ("Asset / mechanism", "[name, modality, MoA]"),
        ("Population", "[indication, line, biomarker]"),
        ("Lifecycle stage", "[pre-launch / launch / post-launch]"),
        ("Priority geographies", "[global, US, EU5, JP, local markets]"),
        ("TPP / TVP", "[clinical and value targets]"),
    ]
    for i, (h, b) in enumerate(fields):
        y = 2.02 + i * 0.68
        panel(slide, 2.25, y, 5.45, 0.48, h, "orange" if i == 0 else "blue", "white")
        txt(slide, 4.75, y + 0.16, 2.65, 0.1, b, size=7.2)


def value_story():
    slide = page("Value narrative evidence story", "Context", "CONTEXT", "Link each value message to evidence required, support, gaps, and tactics.")
    labels = ["Value driver", "Evidence required", "Current support", "Gap", "Planned tactic"]
    colors = ["blue", "deep_blue", "sky", "orange", "blue"]
    for i, label in enumerate(labels):
        x = 2.18 + i * 2.03
        rect(slide, x, 2.25, 1.58, 0.34, colors[i])
        txt(slide, x + 0.05, 2.36, 1.48, 0.08, label, size=5.9, bold=True, color=c("white"), align=PP_ALIGN.CENTER)
        rect(slide, x, 2.88, 1.58, 1.35, "panel", "line", 0.35)
        txt(slide, x + 0.16, 3.38, 1.24, 0.16, "[content]", size=8.2, color=c("gray"), align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.62), Inches(3.55), Inches(x + 1.95), Inches(3.55))
            line(conn, c("mid"), 1.0)
    panel(slide, 2.20, 5.25, 10.0, 0.7, "Narrative test", "orange", "white")
    txt(slide, 2.50, 5.52, 9.0, 0.14, "Can every value message be supported by source-backed evidence or a funded evidence activity?", size=10.2, bold=True)


def stakeholder_needs():
    slide = page("Stakeholder evidence needs", "Evidence needs", "EVIDENCE NEEDS", "Map decision makers, questions, evidence standards, and timing.")
    stakeholders = [("Regulators", "blue"), ("HTA / payers", "deep_blue"), ("HCPs / KOLs", "sky"), ("Patients", "orange"), ("Policy / systems", "blue"), ("Internal teams", "deep_blue")]
    for i, (name, color) in enumerate(stakeholders):
        x = 2.20 + (i % 3) * 3.38
        y = 2.05 + (i // 3) * 1.85
        panel(slide, x, y, 2.78, 1.25, name, color, "white")
        body(slide, x + 0.24, y + 0.50, 2.18, 0.46, ["Question: [ ]", "Standard: [ ]", "Timing: [ ]"], size=7.2)


def patient_journey():
    slide = page("Patient journey evidence map", "Evidence needs", "EVIDENCE NEEDS", "Identify where evidence changes diagnosis, access, treatment, adherence, or outcomes.")
    stages = ["Pre-diagnosis", "Diagnosis", "Treatment choice", "Access", "Follow-up", "Outcomes"]
    for i, stage in enumerate(stages):
        x = 2.25 + i * 1.62
        rect(slide, x, 2.30, 1.25, 0.42, "blue" if i % 2 == 0 else "deep_blue")
        txt(slide, x + 0.05, 2.45, 1.15, 0.08, stage, size=5.5, bold=True, color=c("white"), align=PP_ALIGN.CENTER)
        panel(slide, x, 3.0, 1.25, 1.45, None, "orange", "panel")
        txt(slide, x + 0.12, 3.38, 1.0, 0.26, "[evidence need]", size=7.2, color=c("gray"), align=PP_ALIGN.CENTER)
    panel(slide, 2.25, 5.25, 9.75, 0.65, "Patient-centered evidence focus: PRO, COA, HRQoL, preferences, burden, adherence, caregiver evidence", "orange", "white")


def inventory():
    slide = page("Current evidence inventory", "Evidence needs", "EVIDENCE NEEDS", "Catalogue completed, ongoing, and planned evidence before defining new activities.")
    shp = table(slide, 2.10, 2.0, 10.85, 4.9, 9, 6, ["Source", "Status", "Population", "Endpoint / output", "Stakeholder", "Limits"], [1.55, 1.05, 1.75, 2.45, 1.45, 2.6], "deep_blue", 5.9)
    rows = ["Clinical trials", "RWE", "HEOR model", "SLR / ITC", "COA / PRO", "Epidemiology", "Publications", "Other"]
    for r, label in enumerate(rows, 1):
        vals = [label, "[status]", "[population]", "[output]", "[stakeholder]", "[quality / usability limits]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def gap_matrix():
    slide = page("Evidence gap matrix", "Gap analysis", "GAP ANALYSIS", "Prioritize by stakeholder decision, current coverage, consequence, and action.")
    shp = table(slide, 2.10, 2.0, 9.1, 4.85, 8, 5, ["Gap", "Decision", "Current evidence", "Priority", "Action"], [1.85, 2.0, 2.4, 0.9, 1.95], "deep_blue", 5.9)
    for r in range(1, 8):
        vals = [f"G{r}: [gap]", "[decision]", "[evidence]", "[P1/P2/P3]", "[tactic]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val
    panel(slide, 11.45, 2.0, 1.25, 4.85, "Consequence", "orange", "panel")
    body(slide, 11.66, 2.62, 0.82, 2.8, ["[label delay]", "[HTA risk]", "[slow uptake]", "[unfunded lifecycle]"], size=6.3)


def prioritization():
    slide = page("Gap prioritization scorecard", "Gap analysis", "GAP ANALYSIS", "Use a transparent rubric before converting gaps into funded activities.")
    for i, (num, label, color) in enumerate([("5", "Strategic impact", "orange"), ("4", "Urgency", "blue"), ("3", "Feasibility", "deep_blue"), ("2", "Reuse potential", "orange")]):
        metric(slide, 2.25 + i * 2.35, 2.08, num, label, color)
    shp = table(slide, 2.20, 4.10, 10.35, 2.3, 5, 5, ["Gap", "Impact", "Urgency", "Feasibility", "Decision"], [2.75, 1.5, 1.5, 1.7, 2.9], "deep_blue", 6.5)
    for r in range(1, 5):
        vals = [f"G{r}: [gap]", "[1-5]", "[1-5]", "[1-5]", "[fund / revise / hold]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def orqs():
    slide = page("Open research questions", "Gap analysis", "GAP ANALYSIS", "Translate priority gaps into PICOT-ready research questions.")
    shp = table(slide, 2.10, 2.0, 10.85, 4.9, 8, 7, ["ORQ", "Population", "Intervention", "Comparator", "Outcome", "Timing", "Decision use"], [0.7, 1.6, 1.65, 1.6, 1.75, 1.15, 2.4], "deep_blue", 5.5)
    for r in range(1, 8):
        vals = [f"Q{r}", "[population]", "[intervention]", "[comparator]", "[outcome]", "[timing]", "[decision]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def roadmap():
    slide = page("Integrated evidence roadmap", "Roadmap", "ROADMAP", "Function-level swimlanes with readouts, decisions, and pull-through milestones.")
    quarters = ["Y1 Q1", "Q2", "Q3", "Q4", "Y2 Q1", "Q2", "Q3", "Q4"]
    for i, q in enumerate(quarters):
        txt(slide, 4.04 + i * 0.92, 2.0, 0.65, 0.12, q, size=5.8, bold=True, color=c("gray"), align=PP_ALIGN.CENTER)
    lanes = [("Clinical / regulatory", "blue"), ("RWE", "deep_blue"), ("HEOR / access", "sky"), ("Patient / COA", "orange"), ("Publications", "blue")]
    for i, (lane, color) in enumerate(lanes):
        y = 2.45 + i * 0.72
        txt(slide, 2.20, y + 0.08, 1.6, 0.10, lane, size=6.6, bold=True)
        rect(slide, 3.95, y + 0.18, 7.55, 0.012, "line")
        rect(slide, 4.10 + (i % 3) * 0.55, y, 1.35 + (i % 2) * 0.52, 0.32, color)
        txt(slide, 4.22 + (i % 3) * 0.55, y + 0.10, 1.05, 0.07, "[activity]", size=5.1, bold=True, color=c("white"), align=PP_ALIGN.CENTER)
        rect(slide, 6.30 + i * 0.68, y - 0.04, 0.15, 0.40, "orange")
    panel(slide, 2.20, 6.35, 9.9, 0.43, "Decision gates: [protocol] | [data cut] | [abstract] | [regulatory] | [HTA] | [launch]", "orange", "white")


def tracker():
    slide = page("Evidence activity tracker", "Roadmap", "ROADMAP", "A single working table for studies, analyses, syntheses, and dissemination activities.")
    shp = table(slide, 2.05, 2.0, 10.95, 4.85, 9, 8, ["ID", "Activity", "Gap", "Method / data", "Owner", "Timing", "Status", "Priority"], [0.45, 1.7, 0.65, 2.35, 0.95, 1.15, 1.55, 0.85], "deep_blue", 5.3)
    for r in range(1, 9):
        vals = [f"A{r}", "[study / analysis]", "G[ ]", "[design / source]", "[owner]", "[start-readout]", "[idea/planned/active]", "[P1]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def study_concept():
    slide = page("Study concept one-pager", "Roadmap", "ROADMAP", "Duplicate this page for each priority evidence-generation activity.")
    items = [
        ("Study / activity", "[name and ID]", "blue"),
        ("Gap and decision use", "[gap, ORQ, stakeholder]", "orange"),
        ("Objective", "[objective and hypothesis]", "deep_blue"),
        ("Design", "[population, comparator, endpoint, data]", "blue"),
        ("Feasibility", "[cost, complexity, dependency]", "orange"),
        ("Pull-through", "[publication, dossier, claims]", "deep_blue"),
    ]
    for i, (h, b, color) in enumerate(items):
        x = 2.20 + (i % 3) * 3.38
        y = 2.05 + (i // 3) * 1.58
        panel(slide, x, y, 2.78, 1.05, h, color, "white")
        txt(slide, x + 0.24, y + 0.55, 2.25, 0.18, b, size=7.4, color=c("gray"))
    panel(slide, 2.20, 5.78, 9.95, 0.52, "Decision: [approve / revise / defer]    Next step: [action]    Date: [ ]", "orange", "white")


def workstreams():
    slide = page("Specialized evidence workstreams", "Roadmap", "ROADMAP", "Make HEOR, RWE, patient-centered, regulatory, and communications evidence explicit.")
    work = [
        ("HEOR / access", ["SLR / ITC", "economic model", "budget impact", "payer research"], "deep_blue"),
        ("RWE / data", ["retrospective database", "prospective study", "registry", "external control"], "blue"),
        ("Patient-centered", ["COA / PRO", "HRQoL", "preference", "caregiver burden"], "orange"),
        ("Publications / comms", ["abstract", "manuscript", "field medical", "claims matrix"], "sky"),
    ]
    for i, (h, lines, color) in enumerate(work):
        x = 2.20 + (i % 2) * 5.0
        y = 2.05 + (i // 2) * 1.95
        panel(slide, x, y, 4.35, 1.38, h, color, "white")
        body(slide, x + 0.24, y + 0.52, 3.65, 0.62, lines, size=7.8)


def governance():
    slide = page("Governance and refresh cadence", "Governance", "GOVERNANCE", "Operate the IEGP as a living plan, not a one-time deliverable.")
    shp = table(slide, 2.20, 2.0, 5.5, 4.55, 7, 4, ["Forum", "Purpose", "Cadence", "Rights"], [1.2, 2.1, 1.0, 1.2], "deep_blue", 5.9)
    rows = [["Core team", "[maintain plan]", "[monthly]", "[recommend]"], ["Evidence council", "[prioritize/fund]", "[quarterly]", "[approve]"], ["Global-local", "[adapt needs]", "[quarterly]", "[recommend]"], ["Leadership", "[trade-offs]", "[semiannual]", "[approve]"], ["Publications", "[pull-through]", "[monthly]", "[recommend]"], ["Other", "[purpose]", "[cadence]", "[rights]"]]
    for r, row in enumerate(rows, 1):
        for col, val in enumerate(row):
            shp.table.cell(r, col).text = val
    panel(slide, 8.15, 2.0, 3.75, 4.55, "Update triggers", "orange", "panel")
    body(slide, 8.45, 2.62, 3.0, 2.8, ["Clinical readout", "Regulatory feedback", "HTA / payer change", "Competitor approval or data", "Guideline update", "Budget or feasibility change"], size=8.8)


def kpis():
    slide = page("Evidence impact dashboard", "Governance", "GOVERNANCE", "Track execution and evidence impact with business-readable measures.")
    metrics = [("72%", "Gap closure", "orange"), ("18/24", "On-track activities", "blue"), ("6", "Markets reused", "deep_blue"), ("12", "Pull-through outputs", "orange")]
    for i, (num, label, color) in enumerate(metrics):
        metric(slide, 2.25 + i * 2.38, 2.08, num, label, color)
    shp = table(slide, 2.20, 4.10, 10.35, 2.3, 5, 5, ["KPI", "Definition", "Baseline", "Target", "Owner"], [1.65, 4.0, 1.35, 1.35, 2.0], "deep_blue", 6.3)
    for r in range(1, 5):
        vals = ["[KPI]", "[definition]", "[baseline]", "[target]", "[owner]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def risk_register():
    slide = page("Risk and dependency register", "Governance", "GOVERNANCE", "Document constraints that could reduce evidence usefulness or delay decisions.")
    shp = table(slide, 2.10, 2.0, 10.85, 4.8, 8, 5, ["Risk / dependency", "Activity", "Impact", "Mitigation", "Decision needed"], [2.35, 1.55, 1.05, 3.25, 2.65], "deep_blue", 5.9)
    for r in range(1, 8):
        vals = ["[risk]", "[activity]", "[H/M/L]", "[mitigation or alternate path]", "[decision / owner]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def appendix():
    slide = page("Appendix worksheets", "Appendix", "APPENDIX", "Keep source provenance, regional adaptation, methodology, and final review audit-ready.")
    items = [
        ("Regional adaptation", "local needs, reusable global evidence, local gaps, owner", "blue"),
        ("Evidence source log", "source, date, owner, confidence, limitation, claim link", "deep_blue"),
        ("Methodology", "inputs, workshops, scoring rubric, assumptions, exclusions", "orange"),
        ("Glossary", "IEGP/IEP, ORQ, PICOT, RWE, HEOR, HTA, COA/PRO, IIR", "sky"),
    ]
    for i, (h, desc, color) in enumerate(items):
        x = 2.20 + (i % 2) * 5.0
        y = 2.05 + (i // 2) * 1.72
        panel(slide, x, y, 4.35, 1.1, h, color, "white")
        txt(slide, x + 0.24, y + 0.58, 3.6, 0.2, desc, size=7.4, color=c("gray"))


SLIDES = [
    cover,
    tokens,
    content_map,
    lambda: section_divider("Briefing", "Briefing", "Set evidence priorities and decisions needed now."),
    executive,
    decisions,
    lambda: section_divider("Context", "Context", "Frame the asset, value story, and stakeholder landscape."),
    asset_snapshot,
    value_story,
    lambda: section_divider("Evidence needs", "Evidence needs", "Map decision needs before building new tactics."),
    stakeholder_needs,
    patient_journey,
    inventory,
    lambda: section_divider("Gap analysis", "Gap analysis", "Convert needs into prioritized and researchable gaps."),
    gap_matrix,
    prioritization,
    orqs,
    lambda: section_divider("Roadmap", "Roadmap", "Translate priority gaps into activities, timing, and owners."),
    roadmap,
    tracker,
    study_concept,
    workstreams,
    lambda: section_divider("Governance", "Governance", "Run the plan as a living evidence system."),
    governance,
    kpis,
    risk_register,
    appendix,
]


for make in SLIDES:
    make()

PRS.save(OUT_FILE)
print(OUT_FILE.resolve())
