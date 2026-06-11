from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("outputs/iegp_powerpoint_template")
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "IEGP_Integrated_Evidence_Generation_Plan_Template.pptx"

PRS = Presentation()
PRS.slide_width = Inches(13.333)
PRS.slide_height = Inches(7.5)

W = PRS.slide_width
H = PRS.slide_height

COLORS = {
    "ink": RGBColor(35, 38, 43),
    "muted": RGBColor(99, 105, 112),
    "line": RGBColor(213, 217, 222),
    "soft": RGBColor(247, 248, 246),
    "panel": RGBColor(255, 255, 255),
    "navy": RGBColor(31, 57, 82),
    "teal": RGBColor(27, 126, 124),
    "lime": RGBColor(190, 214, 69),
    "amber": RGBColor(231, 170, 74),
    "rose": RGBColor(190, 83, 93),
    "blue": RGBColor(87, 126, 171),
}


def rgb(name):
    return COLORS[name]


def blank_slide():
    return PRS.slides.add_slide(PRS.slide_layouts[6])


def set_fill(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def set_line(shape, color, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, x, y, w, h, text="", font_size=16, color=None, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or rgb("ink")
    return box


def add_body(slide, x, y, w, h, lines, font_size=13, color=None, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        if bullet:
            p._p.get_or_add_pPr().set("marL", "171450")
            p._p.get_or_add_pPr().set("indent", "-114300")
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color or rgb("ink")
    return box


def add_title(slide, title, subtitle=None, section=None):
    add_textbox(slide, 0.55, 0.32, 10.7, 0.45, title, font_size=24, bold=True, color=rgb("ink"))
    if subtitle:
        add_textbox(slide, 0.57, 0.82, 10.9, 0.33, subtitle, font_size=10.5, color=rgb("muted"))
    if section:
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.25), Inches(0.35), Inches(1.35), Inches(0.3))
        set_fill(pill, rgb("soft"))
        set_line(pill, rgb("line"), 0.7)
        add_textbox(slide, 11.32, 0.405, 1.2, 0.16, section.upper(), font_size=6.8, bold=True, color=rgb("navy"), align=PP_ALIGN.CENTER)
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.18), Inches(12.2), Inches(0.012))
    set_fill(rule, rgb("line"))
    rule.line.fill.background()


def add_footer(slide, page_label):
    add_textbox(slide, 0.56, 7.08, 5.2, 0.18, "Integrated Evidence Generation Plan template", font_size=7.5, color=rgb("muted"))
    add_textbox(slide, 11.7, 7.08, 1.0, 0.18, f"{len(PRS.slides):02d}", font_size=7.5, color=rgb("muted"), align=PP_ALIGN.RIGHT)


def add_panel(slide, x, y, w, h, title=None, fill="panel"):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(rect, rgb(fill))
    set_line(rect, rgb("line"), 0.8)
    if title:
        add_textbox(slide, x + 0.18, y + 0.14, w - 0.35, 0.22, title, font_size=9.5, bold=True, color=rgb("navy"))
    return rect


def add_table(slide, x, y, w, h, rows, cols, headers=None, widths=None, font_size=8.2):
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            set_fill(cell, rgb("panel") if r else rgb("navy"))
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Aptos"
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else rgb("ink")
            run.font.bold = r == 0
            cell.text = headers[c] if (r == 0 and headers) else ""
    return table_shape


def add_placeholder(slide, x, y, w, h, title, body=None, accent="teal"):
    add_panel(slide, x, y, w, h)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    set_fill(bar, rgb(accent))
    bar.line.fill.background()
    add_textbox(slide, x + 0.2, y + 0.14, w - 0.35, 0.25, title, font_size=10, bold=True, color=rgb("navy"))
    if body:
        add_body(slide, x + 0.18, y + 0.48, w - 0.32, h - 0.62, body, font_size=8.5, color=rgb("muted"), bullet=False)


def add_section_divider(title, subtitle, page):
    slide = blank_slide()
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, rgb("navy"))
    bg.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.02), Inches(0.12), Inches(5.25))
    set_fill(accent, rgb("lime"))
    accent.line.fill.background()
    add_textbox(slide, 0.92, 1.22, 9.8, 0.55, title, font_size=31, bold=True, color=RGBColor(255, 255, 255))
    add_body(slide, 0.95, 2.03, 8.9, 1.1, [subtitle], font_size=16, color=RGBColor(232, 237, 239), bullet=False)
    add_textbox(slide, 10.95, 6.85, 1.55, 0.22, page, font_size=8, color=RGBColor(232, 237, 239), align=PP_ALIGN.RIGHT)
    return slide


def style_chart_axis_label(slide, x, y, text, rotate=False):
    box = add_textbox(slide, x, y, 1.4 if not rotate else 0.3, 0.2 if not rotate else 1.2, text, font_size=7.5, bold=True, color=rgb("muted"), align=PP_ALIGN.CENTER)
    if rotate:
        box.rotation = 270
    return box


def cover():
    slide = blank_slide()
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, RGBColor(250, 250, 247))
    bg.line.fill.background()
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(5.6), W, Inches(1.9))
    set_fill(band, rgb("navy"))
    band.line.fill.background()
    add_textbox(slide, 0.75, 0.7, 2.4, 0.26, "TEMPLATE", font_size=8.5, bold=True, color=rgb("teal"))
    add_textbox(slide, 0.72, 1.15, 10.5, 1.35, "Integrated Evidence Generation Plan", font_size=39, bold=True, color=rgb("ink"))
    add_textbox(slide, 0.78, 2.52, 8.8, 0.42, "Reusable PowerPoint starting point for asset, indication, and lifecycle evidence planning", font_size=16, color=rgb("muted"))
    for i, label in enumerate(["Strategy", "Evidence needs", "Gap prioritization", "Roadmap", "Governance"]):
        add_placeholder(slide, 0.78 + i * 2.42, 3.45, 2.08, 1.15, label, ["Replace with product-specific content"], accent=["teal", "blue", "amber", "lime", "rose"][i])
    add_textbox(slide, 0.78, 6.12, 7.7, 0.34, "[Product / Asset] | [Indication] | [Region] | [Plan date]", font_size=18, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(slide, 0.8, 6.63, 8.0, 0.2, "Prepared by: [Team / function]  |  Version: [v0.1]", font_size=9.5, color=RGBColor(232, 237, 239))


def how_to_use():
    slide = blank_slide()
    add_title(slide, "How To Use This Template", "Keep the template editable; replace bracketed placeholders with product-specific evidence, owners, and decisions.", "setup")
    add_placeholder(slide, 0.7, 1.55, 3.8, 2.1, "1. Start With Strategy", ["Anchor evidence planning in asset strategy, TPP/TVP, lifecycle stage, geography, and decision windows.", "Define what decisions this IEGP must support."], "teal")
    add_placeholder(slide, 4.8, 1.55, 3.8, 2.1, "2. Convert Needs Into Gaps", ["Catalogue available evidence.", "Map stakeholder evidence needs.", "Prioritize gaps by strategic value and feasibility."], "amber")
    add_placeholder(slide, 8.9, 1.55, 3.35, 2.1, "3. Build The Roadmap", ["Translate gaps into study concepts and activities.", "Assign owners, timing, budget class, and pull-through plan."], "blue")
    add_panel(slide, 0.7, 4.1, 11.55, 1.95, "Template conventions")
    add_body(slide, 0.95, 4.55, 10.9, 1.1, [
        "Replace [bracketed text] with product-specific content.",
        "Keep source-backed facts distinct from assumptions and internal strategy decisions.",
        "Use the evidence activity tracker as the single source of truth for study concepts, status, and owners.",
        "Refresh after major regulatory, competitor, clinical, access, guideline, or data-readout events."
    ], font_size=11)
    add_footer(slide, "02")


def evidence_model():
    slide = blank_slide()
    add_title(slide, "IEGP Operating Model", "Public guidance converges on a living, cross-functional plan that turns evidence needs into prioritized activities.", "setup")
    steps = [
        ("Strategic context", "TPP/TVP, value narrative, lifecycle stage, decision gates"),
        ("Evidence needed", "Stakeholder questions by regulator, HTA/payer, HCP, patient, policy, internal team"),
        ("Evidence available", "Trial data, RWE, HEOR, SLRs, COA/PROs, publications, competitor evidence"),
        ("Gaps and priorities", "Risk, value, feasibility, timing, geography, funding and ownership"),
        ("Evidence roadmap", "Studies, analyses, data sources, publications, owners, milestones, KPIs")
    ]
    x0 = 0.65
    for i, (name, desc) in enumerate(steps):
        x = x0 + i * 2.5
        add_placeholder(slide, x, 1.85, 2.1, 2.15, name, [desc], accent=["teal", "blue", "amber", "rose", "lime"][i])
        if i < 4:
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x + 2.12), Inches(2.62), Inches(0.34), Inches(0.42))
            set_fill(arr, rgb("line"))
            arr.line.fill.background()
    add_panel(slide, 0.85, 4.65, 11.6, 1.1, "Design principle")
    add_textbox(slide, 1.12, 5.08, 10.9, 0.35, "Evidence to generate = evidence needed for strategic decisions minus evidence already available, adjusted for future landscape change and feasibility.", font_size=15, bold=True, color=rgb("navy"), align=PP_ALIGN.CENTER)
    add_footer(slide, "03")


def executive_summary():
    slide = blank_slide()
    add_title(slide, "Executive Summary", "Summarize the evidence strategy, key gaps, proposed investments, and decisions needed.", "summary")
    add_placeholder(slide, 0.72, 1.45, 3.75, 2.0, "Strategic evidence priorities", [
        "1. [Priority evidence objective]",
        "2. [Priority evidence objective]",
        "3. [Priority evidence objective]"
    ], "teal")
    add_placeholder(slide, 4.8, 1.45, 3.75, 2.0, "Top evidence gaps", [
        "G1 [Gap and decision risk]",
        "G2 [Gap and decision risk]",
        "G3 [Gap and decision risk]"
    ], "amber")
    add_placeholder(slide, 8.88, 1.45, 3.35, 2.0, "Major decisions needed", [
        "[Approve roadmap]",
        "[Fund priority studies]",
        "[Resolve owner / timing]"
    ], "rose")
    add_table(slide, 0.72, 4.05, 11.5, 1.9, 4, 5, headers=["Proposed Activity", "Gap Addressed", "Decision Supported", "Timing", "Investment Class"], widths=[3.1, 2.1, 3.0, 1.5, 1.8], font_size=7.8)
    table = slide.shapes[-1].table
    for r in range(1, 4):
        vals = ["[Study / analysis / synthesis]", "G[ ]", "[Regulatory / access / adoption / lifecycle]", "[Quarter / year]", "[$ / $$ / $$$]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def key_decisions():
    slide = blank_slide()
    add_title(slide, "Key Decisions Needed", "Use this slide for governance meetings and funding reviews.", "summary")
    headers = ["Decision", "Recommendation", "Rationale", "Owner", "Deadline"]
    add_table(slide, 0.7, 1.38, 11.9, 5.2, 7, 5, headers=headers, widths=[2.8, 2.4, 3.65, 1.45, 1.6], font_size=7.8)
    table = slide.shapes[-1].table
    rows = [
        ["[Approve priority gap list]", "[Approve / revise]", "[Why this matters for asset strategy]", "[Function]", "[Date]"],
        ["[Fund evidence activity A]", "[Approve / defer]", "[Decision value and risk of inaction]", "[Function]", "[Date]"],
        ["[Select study design option]", "[Option 1 / 2]", "[Feasibility, credibility, timing]", "[Function]", "[Date]"],
        ["[Resolve regional adaptation]", "[Global / local split]", "[Reusable evidence vs local need]", "[Function]", "[Date]"],
        ["[Confirm owner and cadence]", "[RACI approval]", "[Execution accountability]", "[Function]", "[Date]"],
        ["[Other decision]", "[Recommendation]", "[Rationale]", "[Owner]", "[Date]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def asset_snapshot():
    slide = blank_slide()
    add_title(slide, "Asset Strategy Snapshot", "Summarize the product context before defining evidence needs.", "context")
    fields = [
        ("Asset / mechanism", "[Name, mechanism, modality]"),
        ("Indication / population", "[Target population, line of therapy, biomarker status]"),
        ("Lifecycle stage", "[Pre-Ph2 / Ph3 / launch / post-launch / mature]"),
        ("Target geographies", "[Global, US, EU5, JP, priority local markets]"),
        ("Target product profile", "[Clinical, safety, convenience, durability, differentiation]"),
        ("Target value proposition", "[Expected payer/HCP/patient value drivers]")
    ]
    for i, (label, text) in enumerate(fields):
        x = 0.75 + (i % 2) * 6.0
        y = 1.45 + (i // 2) * 1.35
        add_placeholder(slide, x, y, 5.45, 1.0, label, [text], accent="teal" if i % 2 == 0 else "blue")
    add_panel(slide, 0.75, 5.75, 11.45, 0.82, "Key decision windows")
    add_textbox(slide, 1.0, 6.08, 10.9, 0.22, "[Regulatory interaction]  |  [Phase 3 design lock]  |  [HTA submission]  |  [Launch]  |  [Guideline update]  |  [Lifecycle expansion]", font_size=11.5, color=rgb("ink"), align=PP_ALIGN.CENTER)
    add_footer(slide, "04")


def strategic_imperatives():
    slide = blank_slide()
    add_title(slide, "Strategic Imperatives And Value Narrative", "Translate business, clinical, and access strategy into evidence-generation imperatives.", "context")
    add_table(slide, 0.7, 1.45, 11.9, 4.15, 6, 5, headers=["Imperative", "Decision It Supports", "Primary Stakeholders", "Evidence Needed", "Current Confidence"], widths=[2.2, 2.5, 2.2, 3.0, 2.0], font_size=8.0)
    prompts = [
        ["[e.g., demonstrate differentiated clinical benefit]", "[Label / guideline / uptake decision]", "[Regulators, HCPs]", "[Endpoint, comparator, subgroup evidence]", "[High / Med / Low]"],
        ["[e.g., prove economic value]", "[HTA / formulary / contracting]", "[Payers, HTA]", "[Cost, resource use, budget impact]", "[ ]"],
        ["[e.g., support patient-centered value]", "[Adoption / adherence / persistence]", "[Patients, HCPs]", "[PRO, COA, burden, QoL data]", "[ ]"],
        ["[e.g., de-risk local launch]", "[Market access sequencing]", "[Regional/local teams]", "[Local comparator, epidemiology, pathways]", "[ ]"],
        ["[Add imperative]", "[Decision]", "[Stakeholder]", "[Evidence]", "[ ]"],
    ]
    table = slide.shapes[-1].table
    for r, row in enumerate(prompts, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "05")


def value_narrative():
    slide = blank_slide()
    add_title(slide, "Value Narrative And Evidence Story", "Connect the asset value proposition to the claims and evidence required to support it.", "context")
    add_placeholder(slide, 0.7, 1.4, 3.6, 1.55, "Draft value proposition", [
        "[Concise product value narrative]",
        "[Differentiation vs standard of care and competitors]"
    ], "teal")
    add_placeholder(slide, 4.55, 1.4, 3.6, 1.55, "Core claims / messages", [
        "[Clinical claim]",
        "[Economic/access claim]",
        "[Patient-centered claim]"
    ], "blue")
    add_placeholder(slide, 8.4, 1.4, 3.85, 1.55, "Evidence standard", [
        "[Evidence required for label, payer, medical, publication, or internal use]"
    ], "amber")
    headers = ["Value Driver", "Evidence Needed", "Current Support", "Gap", "Planned Tactic"]
    add_table(slide, 0.7, 3.4, 11.55, 2.55, 5, 5, headers=headers, widths=[2.3, 3.0, 2.25, 1.85, 2.15], font_size=7.7)
    table = slide.shapes[-1].table
    rows = [
        ["Clinical differentiation", "[Endpoint / comparator / subgroup]", "[Evidence]", "[Gap]", "[Tactic]"],
        ["Economic value", "[Cost, utilization, model inputs]", "[Evidence]", "[Gap]", "[Tactic]"],
        ["Patient value", "[PRO/COA, burden, preference]", "[Evidence]", "[Gap]", "[Tactic]"],
        ["Implementation value", "[Pathway, adoption, service delivery]", "[Evidence]", "[Gap]", "[Tactic]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def milestones():
    slide = blank_slide()
    add_title(slide, "Development And Commercial Milestones", "Anchor evidence activities to readouts, submissions, launch windows, and external decision points.", "context")
    headers = ["Milestone", "Date / Window", "Decision Impact", "Evidence Required Before Milestone", "Owner"]
    add_table(slide, 0.7, 1.35, 11.9, 4.95, 8, 5, headers=headers, widths=[2.4, 1.65, 2.85, 3.5, 1.5], font_size=7.5)
    table = slide.shapes[-1].table
    rows = [
        ["[Phase 2/3 readout]", "[Date]", "[Design, claims, value story]", "[Evidence package]", "[Owner]"],
        ["[Regulatory interaction]", "[Date]", "[Label / endpoint / safety]", "[Evidence package]", "[Owner]"],
        ["[HTA / payer submission]", "[Date]", "[Access / reimbursement]", "[Economic and comparative evidence]", "[Owner]"],
        ["[Launch]", "[Date]", "[Adoption readiness]", "[Publication / medical / access evidence]", "[Owner]"],
        ["[Guideline opportunity]", "[Date]", "[Scientific credibility]", "[Published evidence]", "[Owner]"],
        ["[Lifecycle expansion]", "[Date]", "[New indication / population]", "[Evidence package]", "[Owner]"],
        ["[Other]", "[Date]", "[Impact]", "[Evidence]", "[Owner]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def stakeholder_needs():
    slide = blank_slide()
    add_title(slide, "Stakeholder Evidence Needs Matrix", "Capture the evidence questions that must be answered for each external and internal decision maker.", "needs")
    headers = ["Stakeholder", "Decision / Question", "Evidence Standard", "Timing", "Risk If Unanswered"]
    add_table(slide, 0.65, 1.35, 12.0, 5.3, 8, 5, headers=headers, widths=[1.7, 3.25, 2.85, 1.5, 2.7], font_size=7.6)
    rows = [
        ["Regulators", "[Benefit-risk, endpoint, safety, subgroup, external control]", "[RCT / supportive RWE / validated endpoint]", "[Milestone]", "[Label delay / restricted claim]"],
        ["HTA / payers", "[Comparative value, budget impact, durability]", "[Comparator relevance, model inputs, local data]", "[Milestone]", "[Access delay / unfavorable price]"],
        ["HCPs / KOLs", "[Place in therapy, patient selection, sequencing]", "[Head-to-head, real-world use, guidelines]", "[Milestone]", "[Slow adoption]"],
        ["Patients / caregivers", "[QoL, burden, preference, adherence]", "[COA/PRO, qualitative, preference studies]", "[Milestone]", "[Low perceived value]"],
        ["Policy / systems", "[Capacity, pathway, health equity, implementation]", "[Service delivery, epidemiology, outcomes]", "[Milestone]", "[Implementation barrier]"],
        ["Internal teams", "[Claims, launch readiness, lifecycle decisions]", "[Approved evidence package and claims matrix]", "[Milestone]", "[Fragmented activity]"],
        ["[Other]", "[Question]", "[Standard]", "[Timing]", "[Risk]"],
    ]
    table = slide.shapes[-1].table
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "06")


def patient_journey():
    slide = blank_slide()
    add_title(slide, "Patient Journey Evidence Map", "Identify evidence needs across diagnosis, treatment, access, persistence, and outcomes.", "needs")
    stages = [
        ("Pre-diagnosis", "Awareness, symptoms, unmet need"),
        ("Diagnosis", "Testing, referral, severity, biomarker"),
        ("Treatment choice", "Eligibility, sequencing, shared decision"),
        ("Access", "Coverage, affordability, site of care"),
        ("Follow-up", "Adherence, persistence, monitoring"),
        ("Outcomes", "QoL, durability, resource use")
    ]
    for i, (stage, body) in enumerate(stages):
        add_placeholder(slide, 0.65 + i * 2.05, 1.55, 1.75, 1.5, stage, [body, "[Evidence need]"], ["teal", "blue", "amber", "rose", "lime", "teal"][i])
        if i < 5:
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(2.28 + i * 2.05), Inches(2.05), Inches(0.3), Inches(0.34))
            set_fill(arr, rgb("line"))
            arr.line.fill.background()
    add_table(slide, 0.7, 4.0, 11.8, 1.9, 4, 4, headers=["Journey Barrier", "Evidence Question", "Potential Method", "Stakeholder"], widths=[3.0, 4.0, 2.7, 2.1], font_size=7.8)
    table = slide.shapes[-1].table
    for r in range(1, 4):
        vals = ["[Barrier]", "[Evidence question]", "[COA / preference / RWE / qualitative]", "[Patient / HCP / payer]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def geography_needs():
    slide = blank_slide()
    add_title(slide, "Geography And Market Evidence Needs", "Capture market-specific evidence requirements early enough to influence global studies or local tactics.", "needs")
    headers = ["Market / Region", "Decision Body", "Comparator / Standard", "Local Evidence Need", "Timing", "Owner"]
    add_table(slide, 0.62, 1.35, 12.1, 5.4, 8, 6, headers=headers, widths=[1.55, 1.95, 2.25, 3.0, 1.35, 2.0], font_size=7.2)
    table = slide.shapes[-1].table
    rows = [
        ["Global", "[Core evidence council]", "[Global SOC]", "[Core evidence requirement]", "[Date]", "[Owner]"],
        ["US", "[Payers / guidelines]", "[US SOC]", "[Coverage / pathway / outcomes]", "[Date]", "[Owner]"],
        ["EU5", "[HTA bodies]", "[Local comparator]", "[Indirect comparison, model, local data]", "[Date]", "[Owner]"],
        ["Japan", "[PMDA / payer]", "[Local SOC]", "[Local clinical/access evidence]", "[Date]", "[Owner]"],
        ["China", "[NMPA / access]", "[Local SOC]", "[Local epidemiology / registry / access]", "[Date]", "[Owner]"],
        ["[Market]", "[Body]", "[Comparator]", "[Need]", "[Date]", "[Owner]"],
        ["[Market]", "[Body]", "[Comparator]", "[Need]", "[Date]", "[Owner]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def inventory():
    slide = blank_slide()
    add_title(slide, "Current Evidence Inventory", "Catalogue available and planned evidence before defining the incremental IEGP roadmap.", "needs")
    headers = ["Evidence Source", "Status", "Key Findings / Asset Relevance", "Stakeholder Link", "Quality / Limits"]
    add_table(slide, 0.65, 1.35, 12.0, 5.45, 9, 5, headers=headers, widths=[2.2, 1.25, 3.5, 2.2, 2.85], font_size=7.4)
    rows = [
        ["Clinical trials", "[Available / ongoing]", "[Endpoints, comparator, population, subgroup]", "[Reg / HCP / payer]", "[Bias, maturity, generalizability]"],
        ["Real-world evidence", "[Available / planned]", "[Utilization, outcomes, safety, pathway]", "[Payer / HCP]", "[Data source, missingness, confounding]"],
        ["HEOR models", "[Available / planned]", "[CE model, BIA, burden, resource use]", "[HTA / payer]", "[Inputs, assumptions, transferability]"],
        ["SLR / indirect comparison", "[Available / planned]", "[Comparator landscape and evidence base]", "[HTA / internal]", "[Network, heterogeneity, recency]"],
        ["COA / PRO / preference", "[Available / planned]", "[Patient burden, QoL, preference]", "[Patient / HCP / payer]", "[Instrument validity, sample]"],
        ["Epidemiology / market landscape", "[Available / planned]", "[Population size, pathway, competitor shifts]", "[Global / local]", "[Date, geography, data gaps]"],
        ["Publications / congresses", "[Available / planned]", "[Evidence dissemination and claim support]", "[HCP / internal]", "[Status, timing, embargo]"],
        ["[Other]", "[Status]", "[Finding]", "[Stakeholder]", "[Limit]"],
    ]
    table = slide.shapes[-1].table
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "07")


def evidence_quality():
    slide = blank_slide()
    add_title(slide, "Evidence Quality And Usability", "Assess whether existing evidence is decision-ready, transferable, and claim-supporting.", "needs")
    headers = ["Evidence Source", "Strength", "Publication / Disclosure Status", "Jurisdiction Relevance", "Claim Support", "Limitations"]
    add_table(slide, 0.6, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[2.05, 1.2, 2.3, 2.0, 2.0, 2.55], font_size=7.0)
    table = slide.shapes[-1].table
    for r in range(1, 8):
        vals = ["[Source / study]", "[High/Med/Low]", "[Published / planned / internal]", "[Global / local / limited]", "[Claim/message]", "[Bias, comparator, sample, endpoint]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def need_have_generate():
    slide = blank_slide()
    add_title(slide, "Evidence Need vs Evidence Have", "Use this logic to isolate what must be generated, updated, or synthesized.", "prioritize")
    boxes = [
        ("Evidence needed", ["Stakeholder decisions", "TPP/TVP and value narrative", "Future landscape expectations"], "teal"),
        ("Evidence available", ["Completed / ongoing studies", "Published evidence", "Existing RWE, HEOR, COA, SLR"], "blue"),
        ("Landscape evolution", ["Competitors", "Regulatory / HTA expectations", "Guidelines and SoC change"], "amber"),
        ("Evidence to generate", ["New studies", "Analyses / synthesis", "Publication / pull-through"], "lime"),
    ]
    for i, (title, bullets, accent) in enumerate(boxes):
        add_placeholder(slide, 0.8 + i * 3.0, 1.75, 2.55, 2.15, title, bullets, accent)
        if i < 3:
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(3.25 + i * 3.0), Inches(2.45), Inches(0.32), Inches(0.42))
            set_fill(arr, rgb("line"))
            arr.line.fill.background()
    add_panel(slide, 0.85, 4.65, 11.45, 1.1, "Working formula")
    add_textbox(slide, 1.1, 5.05, 11.0, 0.3, "Prioritized evidence to generate = unmet stakeholder needs + future landscape changes - usable available evidence", font_size=16, bold=True, color=rgb("navy"), align=PP_ALIGN.CENTER)
    add_footer(slide, "00")


def future_landscape():
    slide = blank_slide()
    add_title(slide, "Future Landscape Assumptions", "IEGPs should anticipate the evidence environment at decision time, not only today.", "needs")
    columns = [("Clinical standard", "New comparators, sequencing, endpoints, guidelines"), ("Access environment", "HTA rules, payer evidence expectations, affordability pressure"), ("Competitive events", "Readouts, approvals, label expansions, biosimilars/generics"), ("Data and policy", "RWE availability, registries, privacy, health-system priorities")]
    for i, (title, text) in enumerate(columns):
        add_placeholder(slide, 0.7 + i * 3.0, 1.55, 2.65, 2.25, title, [text, "[Add product-specific assumptions]"], accent=["teal", "amber", "rose", "blue"][i])
    add_table(slide, 0.7, 4.35, 11.85, 1.6, 4, 4, headers=["Assumption", "Evidence Implication", "Trigger To Revisit", "Owner"], widths=[3.2, 4.0, 2.4, 2.0], font_size=8.0)
    table = slide.shapes[-1].table
    for r in range(1, 4):
        for c, text in enumerate(["[Assumption]", "[How this changes evidence needs]", "[Event / date]", "[Function]"]):
            table.cell(r, c).text = text
    add_footer(slide, "08")


def gap_matrix():
    slide = blank_slide()
    add_title(slide, "Evidence Gap Matrix", "Prioritize gaps by stakeholder importance, timing, feasibility, and strategic risk.", "prioritize")
    headers = ["Gap ID", "Evidence Gap", "Stakeholder / Decision", "Current Evidence", "Priority", "Recommended Action"]
    add_table(slide, 0.6, 1.32, 12.1, 5.55, 9, 6, headers=headers, widths=[0.75, 3.0, 2.4, 2.35, 1.1, 2.5], font_size=7.3)
    rows = [
        ["G1", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Study / analysis / synthesis]"],
        ["G2", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Action]"],
        ["G3", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Action]"],
        ["G4", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Action]"],
        ["G5", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Action]"],
        ["G6", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Action]"],
        ["G7", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Action]"],
        ["G8", "[Gap statement]", "[Stakeholder and decision]", "[What exists today]", "[H/M/L]", "[Action]"],
    ]
    table = slide.shapes[-1].table
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "09")


def prioritization():
    slide = blank_slide()
    add_title(slide, "Gap Prioritization Heatmap", "Use a consistent scoring rubric before converting gaps into funded evidence activities.", "prioritize")
    add_panel(slide, 0.72, 1.45, 5.75, 4.9, "Impact x feasibility map")
    colors = [[rgb("amber"), rgb("lime"), rgb("lime")], [rgb("amber"), rgb("amber"), rgb("lime")], [rgb("line"), rgb("amber"), rgb("rose")]]
    labels = [["P2", "P1", "P1"], ["P3", "P2", "P1"], ["Hold", "P3", "P2"]]
    for r in range(3):
        for c in range(3):
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.18 + c * 1.45), Inches(2.1 + r * 0.92), Inches(1.35), Inches(0.82))
            set_fill(cell, colors[2-r][c])
            set_line(cell, RGBColor(255, 255, 255), 1.5)
            add_textbox(slide, 1.47 + c * 1.45, 2.36 + r * 0.92, 0.7, 0.18, labels[2-r][c], font_size=12, bold=True, color=rgb("navy"), align=PP_ALIGN.CENTER)
    style_chart_axis_label(slide, 2.05, 5.05, "Feasibility / ability to generate")
    style_chart_axis_label(slide, 0.88, 2.7, "Strategic impact", rotate=True)
    add_body(slide, 4.95, 2.0, 1.2, 2.9, ["High", "Medium", "Low"], font_size=8.5, bullet=False, color=rgb("muted"))
    add_panel(slide, 6.85, 1.45, 5.65, 4.9, "Scoring criteria")
    add_table(slide, 7.08, 1.9, 5.15, 3.75, 6, 3, headers=["Criterion", "Score", "Notes"], widths=[2.1, 1.0, 2.05], font_size=7.2)
    table = slide.shapes[-1].table
    criteria = [
        ["Strategic importance", "1-5", "[Decision value / risk]"],
        ["Stakeholder urgency", "1-5", "[Timing and dependency]"],
        ["Evidence feasibility", "1-5", "[Data, design, cost]"],
        ["Differentiation potential", "1-5", "[Claim / value impact]"],
        ["Global-local leverage", "1-5", "[Reusable across markets]"],
    ]
    for r, row in enumerate(criteria, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "10")


def priority_gap_deep_dive():
    slide = blank_slide()
    add_title(slide, "Priority Gap Deep Dive", "Duplicate this worksheet for each top-priority evidence gap.", "prioritize")
    add_placeholder(slide, 0.7, 1.35, 3.65, 1.05, "Gap statement", ["[Concise gap and why existing evidence is insufficient]"], "rose")
    add_placeholder(slide, 4.6, 1.35, 3.65, 1.05, "Affected decisions", ["[Regulatory, HTA, HCP, patient, internal, lifecycle]"], "amber")
    add_placeholder(slide, 8.5, 1.35, 3.65, 1.05, "Risk of inaction", ["[Decision consequence, timing, market impact]"], "blue")
    add_placeholder(slide, 0.7, 2.75, 5.55, 1.45, "Current evidence and limitations", ["[Summarize available evidence]", "[Quality, transferability, comparator, endpoint, timing limits]"], "teal")
    add_placeholder(slide, 6.6, 2.75, 5.55, 1.45, "Evidence options", ["[Option A: study / analysis]", "[Option B: synthesis / RWE / model]", "[Option C: defer / monitor]"], "lime")
    add_table(slide, 0.7, 4.65, 11.45, 1.35, 4, 5, headers=["Option", "Credibility", "Timing", "Cost / Complexity", "Recommendation"], widths=[2.25, 2.0, 1.6, 2.3, 3.3], font_size=7.5)
    table = slide.shapes[-1].table
    for r in range(1, 4):
        vals = [f"Option {r}", "[High/Med/Low]", "[Date]", "[$ / $$ / $$$]", "[Select / hold / reject]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def research_questions():
    slide = blank_slide()
    add_title(slide, "Open Research Questions", "Convert priority gaps into answerable questions before selecting study designs.", "prioritize")
    headers = ["ORQ ID", "Gap Link", "Research Question", "PICOT / Design Notes", "Decision Use"]
    add_table(slide, 0.65, 1.32, 12.0, 5.5, 8, 5, headers=headers, widths=[0.85, 1.1, 4.1, 3.6, 2.35], font_size=7.4)
    rows = [
        ["ORQ1", "G[ ]", "[What evidence question must be answered?]", "P: [ ] I: [ ] C: [ ] O: [ ] T: [ ]", "[Decision]"],
        ["ORQ2", "G[ ]", "[Question]", "Population, endpoint, comparator, follow-up, data source", "[Decision]"],
        ["ORQ3", "G[ ]", "[Question]", "[Design considerations]", "[Decision]"],
        ["ORQ4", "G[ ]", "[Question]", "[Design considerations]", "[Decision]"],
        ["ORQ5", "G[ ]", "[Question]", "[Design considerations]", "[Decision]"],
        ["ORQ6", "G[ ]", "[Question]", "[Design considerations]", "[Decision]"],
        ["ORQ7", "G[ ]", "[Question]", "[Design considerations]", "[Decision]"],
    ]
    table = slide.shapes[-1].table
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "11")


def evidence_options():
    slide = blank_slide()
    add_title(slide, "Evidence Generation Options", "Compare possible evidence tactics before selecting the roadmap.", "roadmap")
    headers = ["Evidence Option", "Best Use", "Strengths", "Limitations", "Typical Owner"]
    add_table(slide, 0.65, 1.32, 12.0, 5.45, 8, 5, headers=headers, widths=[2.0, 2.6, 2.75, 3.15, 1.5], font_size=7.2)
    table = slide.shapes[-1].table
    rows = [
        ["RCT / Phase IV", "[Causal efficacy/safety, label/lifecycle]", "[High internal validity]", "[Cost, time, feasibility]", "[Clinical]"],
        ["RWE retrospective", "[Utilization, outcomes, safety, burden]", "[Speed, external validity]", "[Confounding, data gaps]", "[RWE]"],
        ["Registry / prospective", "[Natural history, long-term outcomes]", "[Rich longitudinal data]", "[Operational burden]", "[RWE/Medical]"],
        ["HEOR model / BIA", "[HTA, payer, budget decisions]", "[Decision-ready economics]", "[Assumption-sensitive]", "[HEOR]"],
        ["SLR / meta-analysis", "[Evidence landscape, comparator context]", "[Efficient synthesis]", "[Dependent on source data]", "[HEOR/Medical]"],
        ["Post-hoc / secondary analysis", "[Subgroups, endpoints, exploratory claims]", "[Uses existing data]", "[Multiplicity, credibility]", "[Clinical/Stats]"],
        ["IIR / external research", "[Independent science, unmet questions]", "[External credibility]", "[Control, timing]", "[Medical]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def roadmap():
    slide = blank_slide()
    add_title(slide, "Evidence Generation Roadmap", "Map prioritized activities across lifecycle milestones and decision windows.", "roadmap")
    months = ["Q1", "Q2", "Q3", "Q4", "Q5", "Q6", "Q7", "Q8"]
    for i, q in enumerate(months):
        add_textbox(slide, 3.05 + i * 1.08, 1.35, 0.8, 0.2, q, font_size=8.5, bold=True, color=rgb("navy"), align=PP_ALIGN.CENTER)
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.42 + i * 1.08), Inches(1.65), Inches(0.012), Inches(4.65))
        set_fill(line, rgb("line"))
        line.line.fill.background()
    lanes = [("Clinical / regulatory", "teal"), ("RWE / HEOR", "blue"), ("COA / patient", "amber"), ("Access / HTA", "rose"), ("Publications / pull-through", "lime")]
    for r, (lane, accent) in enumerate(lanes):
        y = 1.78 + r * 0.9
        add_textbox(slide, 0.72, y + 0.22, 2.1, 0.22, lane, font_size=8.8, bold=True, color=rgb("ink"))
        lane_rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.9), Inches(y + 0.72), Inches(8.85), Inches(0.012))
        set_fill(lane_rule, rgb("line"))
        lane_rule.line.fill.background()
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.25 + (r % 3) * 0.7), Inches(y + 0.12), Inches(1.95 + (r % 2) * 0.65), Inches(0.34))
        set_fill(bar, rgb(accent))
        set_line(bar, rgb(accent), 0.5)
        add_textbox(slide, 3.36 + (r % 3) * 0.7, y + 0.19, 1.7 + (r % 2) * 0.6, 0.13, "[Activity]", font_size=6.9, bold=True, color=rgb("navy"), align=PP_ALIGN.CENTER)
    add_panel(slide, 0.72, 6.38, 11.7, 0.46, "Milestones: [Protocol finalization] | [Data cut] | [Abstract submission] | [Regulatory interaction] | [HTA submission] | [Launch]")
    add_footer(slide, "12")


def activity_tracker():
    slide = blank_slide()
    add_title(slide, "Evidence Activity Tracker", "Use this as the working table for each study, analysis, synthesis, or dissemination activity.", "roadmap")
    headers = ["ID", "Activity", "Evidence Need", "Method / Data Source", "Owner", "Timing", "Status", "Priority"]
    add_table(slide, 0.55, 1.32, 12.25, 5.55, 9, 8, headers=headers, widths=[0.55, 2.05, 2.15, 2.45, 1.15, 1.25, 1.25, 1.0], font_size=6.6)
    table = slide.shapes[-1].table
    for r in range(1, 9):
        vals = [f"A{r}", "[Study / analysis]", "G[ ] / ORQ[ ]", "[Design, dataset, synthesis]", "[Function]", "[Start-readout]", "[Idea / planned / active / complete]", "[P1/P2/P3]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "13")


def claims_matrix():
    slide = blank_slide()
    add_title(slide, "Claims Matrix", "Connect planned evidence to value messages and permitted use cases.", "roadmap")
    headers = ["Claim / Message", "Stakeholder", "Required Evidence", "Current Support", "Gap", "Planned Tactic"]
    add_table(slide, 0.6, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[2.4, 1.55, 2.8, 2.05, 1.65, 1.65], font_size=7.0)
    table = slide.shapes[-1].table
    rows = [
        ["[Clinical benefit claim]", "[HCP / regulator]", "[Endpoint, comparator, population]", "[Study/source]", "[Gap]", "[Activity]"],
        ["[Safety/tolerability claim]", "[HCP / regulator]", "[Safety dataset, long-term evidence]", "[Study/source]", "[Gap]", "[Activity]"],
        ["[Economic value claim]", "[Payer / HTA]", "[Resource use, model, BIA]", "[Study/source]", "[Gap]", "[Activity]"],
        ["[Patient value claim]", "[Patient / HCP]", "[PRO/COA/preference/burden]", "[Study/source]", "[Gap]", "[Activity]"],
        ["[Implementation claim]", "[System / policy]", "[Pathway, capacity, outcomes]", "[Study/source]", "[Gap]", "[Activity]"],
        ["[Other]", "[Stakeholder]", "[Evidence]", "[Support]", "[Gap]", "[Tactic]"],
        ["[Other]", "[Stakeholder]", "[Evidence]", "[Support]", "[Gap]", "[Tactic]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def study_concept():
    slide = blank_slide()
    add_title(slide, "Study Concept One-Pager", "Duplicate this slide for each priority evidence-generation activity.", "roadmap")
    add_placeholder(slide, 0.65, 1.35, 3.7, 1.0, "Study / activity name", ["[Name and ID]"], "teal")
    add_placeholder(slide, 4.58, 1.35, 3.7, 1.0, "Evidence gap and decision use", ["[Gap ID, ORQ, stakeholder decision]"], "amber")
    add_placeholder(slide, 8.5, 1.35, 3.7, 1.0, "Priority and owner", ["[P1/P2/P3, accountable function]"], "blue")
    add_placeholder(slide, 0.65, 2.6, 5.6, 1.35, "Objective and hypothesis", ["[Primary objective]", "[Hypothesis / expected value]"], "teal")
    add_placeholder(slide, 6.55, 2.6, 5.65, 1.35, "Design summary", ["[Design, population, comparator, endpoint, duration, data source]"], "blue")
    add_placeholder(slide, 0.65, 4.22, 3.7, 1.3, "Feasibility", ["[Data access, operational complexity, cost class, dependencies]"], "amber")
    add_placeholder(slide, 4.58, 4.22, 3.7, 1.3, "Milestones", ["[Protocol, start, interim, readout, publication]"], "rose")
    add_placeholder(slide, 8.5, 4.22, 3.7, 1.3, "Pull-through", ["[Claims matrix, publication, field medical, HTA dossier, guideline]"], "lime")
    add_panel(slide, 0.65, 5.85, 11.55, 0.55, "Decision: [Approve / revise / defer]  |  Required next step: [Action]  |  Date: [ ]")
    add_footer(slide, "14")


def heor_access_plan():
    slide = blank_slide()
    add_title(slide, "HEOR And Market Access Evidence Plan", "Define economic, comparative, payer, and HTA evidence activities.", "roadmap")
    headers = ["Activity", "Purpose", "Inputs Needed", "Market / Stakeholder", "Timing", "Owner"]
    add_table(slide, 0.62, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[2.1, 2.65, 2.6, 2.15, 1.25, 1.35], font_size=7.0)
    table = slide.shapes[-1].table
    rows = [
        ["SLR / evidence synthesis", "[Comparator landscape, model inputs]", "[Search strategy, criteria]", "[HTA / payer]", "[Date]", "[HEOR]"],
        ["Indirect treatment comparison", "[Relative effectiveness]", "[Trial network, endpoints]", "[HTA]", "[Date]", "[HEOR]"],
        ["Cost-effectiveness model", "[Value demonstration]", "[Clinical, utility, cost inputs]", "[HTA / payer]", "[Date]", "[HEOR]"],
        ["Budget impact model", "[Affordability planning]", "[Population, uptake, cost]", "[Payer]", "[Date]", "[HEOR]"],
        ["Payer research / advisory", "[Evidence expectation validation]", "[Profiles, stimuli]", "[Payers]", "[Date]", "[Access]"],
        ["Dossier / value story", "[Submission / negotiation support]", "[Claims, sources, model outputs]", "[HTA / payer]", "[Date]", "[Access]"],
        ["[Other]", "[Purpose]", "[Inputs]", "[Stakeholder]", "[Date]", "[Owner]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def rwe_strategy():
    slide = blank_slide()
    add_title(slide, "RWE And Data Strategy", "Document data sources, feasibility, endpoints, bias risks, and intended use.", "roadmap")
    headers = ["Use Case", "Data Source", "Population / Sample", "Endpoint Availability", "Bias / Limits", "Decision Use"]
    add_table(slide, 0.62, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[2.0, 2.2, 2.1, 2.2, 2.0, 1.6], font_size=7.0)
    table = slide.shapes[-1].table
    rows = [
        ["Treatment patterns", "[Claims / EHR / registry]", "[Population]", "[Available / proxy]", "[Selection, coding]", "[Access / HCP]"],
        ["Comparative effectiveness", "[Dataset]", "[Population]", "[Endpoint]", "[Confounding]", "[HTA / payer]"],
        ["Safety monitoring", "[Dataset]", "[Population]", "[Safety outcome]", "[Under-capture]", "[Reg / medical]"],
        ["External control", "[Dataset / registry]", "[Population]", "[Endpoint]", "[Comparability]", "[Reg / clinical]"],
        ["Burden / resource use", "[Dataset]", "[Population]", "[Cost/utilization]", "[Generalizability]", "[HEOR]"],
        ["Long-term outcomes", "[Registry / follow-up]", "[Population]", "[Durability/QoL]", "[Attrition]", "[Lifecycle]"],
        ["[Other]", "[Source]", "[Population]", "[Endpoint]", "[Limit]", "[Use]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def patient_centered_evidence():
    slide = blank_slide()
    add_title(slide, "Patient-Centered Evidence Plan", "Define COA, PRO, HRQoL, preference, burden, adherence, and caregiver evidence.", "roadmap")
    add_placeholder(slide, 0.7, 1.42, 3.75, 1.65, "Patient-relevant questions", [
        "[Symptoms, functioning, burden, preferences]",
        "[Treatment trade-offs and adherence drivers]"
    ], "teal")
    add_placeholder(slide, 4.75, 1.42, 3.75, 1.65, "Measures / methods", [
        "[Validated COA/PRO, utilities, interviews]",
        "[Preference or qualitative studies]"
    ], "blue")
    add_placeholder(slide, 8.8, 1.42, 3.35, 1.65, "Decision use", [
        "[Label, HTA, medical education, access, lifecycle]"
    ], "amber")
    headers = ["Evidence Need", "Method / Instrument", "Population", "Timing", "Owner"]
    add_table(slide, 0.7, 3.55, 11.45, 2.2, 5, 5, headers=headers, widths=[2.9, 3.25, 2.15, 1.4, 1.75], font_size=7.6)
    table = slide.shapes[-1].table
    for r in range(1, 5):
        vals = ["[Need]", "[COA/PRO/preference/qualitative]", "[Population]", "[Date]", "[Owner]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def regulatory_considerations():
    slide = blank_slide()
    add_title(slide, "Regulatory Evidence Considerations", "Capture regulatory questions, commitments, and evidence acceptability considerations.", "roadmap")
    headers = ["Regulatory Question", "Evidence Needed", "Current Support", "Planned Activity", "Milestone / Interaction", "Owner"]
    add_table(slide, 0.62, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[2.25, 2.45, 2.0, 2.15, 2.0, 1.25], font_size=7.0)
    table = slide.shapes[-1].table
    for r in range(1, 8):
        vals = ["[Question / commitment]", "[Endpoint, safety, subgroup, RWE, COA]", "[Evidence]", "[Activity]", "[Milestone]", "[Owner]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def med_comms():
    slide = blank_slide()
    add_title(slide, "Medical And Scientific Communications Pull-Through", "Plan how generated evidence will be disclosed, published, and used in scientific exchange.", "roadmap")
    headers = ["Evidence Output", "Target Channel", "Audience", "Timing", "Claim / Message Link", "Owner"]
    add_table(slide, 0.62, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[2.25, 2.1, 1.7, 1.35, 3.2, 1.5], font_size=7.0)
    table = slide.shapes[-1].table
    rows = [
        ["[Study readout]", "[Congress / manuscript]", "[HCP / KOL]", "[Date]", "[Claim/value message]", "[Publications]"],
        ["[HEOR model]", "[HTA dossier / payer tool]", "[Payer / HTA]", "[Date]", "[Economic value]", "[Access]"],
        ["[RWE study]", "[Manuscript / field medical]", "[HCP / payer]", "[Date]", "[Real-world value]", "[Medical]"],
        ["[COA analysis]", "[Congress / manuscript]", "[HCP / patient]", "[Date]", "[Patient-centered value]", "[Medical]"],
        ["[Evidence synthesis]", "[Internal narrative / publication]", "[Internal / external]", "[Date]", "[Evidence story]", "[Medical]"],
        ["[Other]", "[Channel]", "[Audience]", "[Date]", "[Message]", "[Owner]"],
        ["[Other]", "[Channel]", "[Audience]", "[Date]", "[Message]", "[Owner]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def external_research_iir():
    slide = blank_slide()
    add_title(slide, "External Research And IIR Areas", "Define external research interests and link them to strategic evidence gaps.", "roadmap")
    headers = ["Area Of Interest", "Evidence Gap Link", "Strategic Rationale", "Eligible Methods", "Governance Notes", "Owner"]
    add_table(slide, 0.62, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[2.25, 1.55, 3.0, 2.1, 1.9, 1.3], font_size=7.0)
    table = slide.shapes[-1].table
    for r in range(1, 8):
        vals = ["[Area]", "G[ ]", "[Why external research could address gap]", "[Clinical / RWE / COA / HEOR]", "[Review, compliance, budget]", "[Medical]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def budget_resources():
    slide = blank_slide()
    add_title(slide, "Budget And Resource View", "Estimate resources by workstream and year to support prioritization and governance decisions.", "governance")
    headers = ["Workstream", "Year 1", "Year 2", "Year 3", "FTE / Vendor Need", "Notes"]
    add_table(slide, 0.65, 1.32, 12.0, 4.3, 7, 6, headers=headers, widths=[2.2, 1.35, 1.35, 1.35, 2.65, 3.1], font_size=7.5)
    table = slide.shapes[-1].table
    rows = [
        ["Clinical / regulatory", "[$]", "[$]", "[$]", "[Internal / vendor]", "[Notes]"],
        ["RWE", "[$]", "[$]", "[$]", "[Internal / vendor]", "[Notes]"],
        ["HEOR / access", "[$]", "[$]", "[$]", "[Internal / vendor]", "[Notes]"],
        ["COA / patient", "[$]", "[$]", "[$]", "[Internal / vendor]", "[Notes]"],
        ["Publications / pull-through", "[$]", "[$]", "[$]", "[Internal / vendor]", "[Notes]"],
        ["Contingency / TBD", "[$]", "[$]", "[$]", "[Need]", "[Notes]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_panel(slide, 0.65, 6.0, 12.0, 0.55, "Budget decision: [Approve all priority activities] | [Phase by decision window] | [Defer low-priority tactics]")
    add_footer(slide, "00")


def global_local():
    slide = blank_slide()
    add_title(slide, "Global-Regional Evidence Adaptation", "Separate globally reusable evidence from local evidence requirements and adaptations.", "roadmap")
    headers = ["Region / Market", "Local Decision Need", "Reusable Global Evidence", "Local Gap", "Local Tactic / Owner"]
    add_table(slide, 0.65, 1.32, 12.0, 5.45, 8, 5, headers=headers, widths=[1.65, 2.75, 2.75, 2.35, 2.5], font_size=7.4)
    rows = [
        ["Global", "[Core evidence need]", "[Clinical / RWE / HEOR / COA]", "[Gap]", "[Global owner]"],
        ["US", "[Coverage / guideline / pathway]", "[Evidence]", "[Gap]", "[US owner]"],
        ["EU5", "[HTA / comparator / local model]", "[Evidence]", "[Gap]", "[EU owner]"],
        ["Japan", "[Regulatory / reimbursement]", "[Evidence]", "[Gap]", "[JP owner]"],
        ["China", "[Regulatory / access / local data]", "[Evidence]", "[Gap]", "[CN owner]"],
        ["[Market]", "[Need]", "[Evidence]", "[Gap]", "[Owner]"],
        ["[Market]", "[Need]", "[Evidence]", "[Gap]", "[Owner]"],
    ]
    table = slide.shapes[-1].table
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "15")


def specialized_workstreams():
    slide = blank_slide()
    add_title(slide, "Specialized Evidence Workstreams", "Use these prompts to ensure RWE, HEOR, access, and patient-centered evidence are represented.", "roadmap")
    cards = [
        ("RWE", ["Retrospective database studies", "Prospective observational studies", "External controls / natural history", "Safety, utilization, persistence, outcomes"], "blue"),
        ("HEOR / access", ["SLR and indirect comparison", "Cost-effectiveness and budget impact", "Burden and resource use", "HTA and payer value evidence"], "amber"),
        ("COA / patient", ["PROs, HRQoL, utility measures", "Patient preference and burden", "Caregiver impact", "Digital endpoints and interviews"], "teal"),
        ("Evidence synthesis", ["Evidence catalog and source log", "Gap updates and literature surveillance", "Claims matrix and publication support", "Evidence quality notes"], "rose"),
    ]
    for i, (title, bullets, accent) in enumerate(cards):
        add_placeholder(slide, 0.75 + (i % 2) * 6.0, 1.55 + (i // 2) * 2.45, 5.45, 2.0, title, bullets, accent)
    add_footer(slide, "16")


def governance():
    slide = blank_slide()
    add_title(slide, "Governance, Cadence, And Decision Rights", "A living IEGP needs explicit ownership, review cadence, and update triggers.", "governance")
    add_table(slide, 0.7, 1.35, 6.0, 4.65, 7, 4, headers=["Forum", "Purpose", "Cadence", "Decision Rights"], widths=[1.3, 2.1, 1.1, 1.5], font_size=7.2)
    table = slide.shapes[-1].table
    rows = [
        ["Core team", "[Maintain plan, resolve dependencies]", "[Biweekly/monthly]", "[Recommend]"],
        ["Evidence council", "[Prioritize gaps, approve studies]", "[Quarterly]", "[Approve]"],
        ["Global-local review", "[Adapt plan, capture local needs]", "[Quarterly]", "[Recommend]"],
        ["Leadership review", "[Funding, trade-offs, risk]", "[Semiannual]", "[Approve funding]"],
        ["Publication / claims", "[Evidence pull-through]", "[Monthly]", "[Recommend]"],
        ["[Other]", "[Purpose]", "[Cadence]", "[Rights]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_panel(slide, 7.05, 1.35, 5.1, 4.65, "Update triggers")
    add_body(slide, 7.35, 1.83, 4.55, 3.65, [
        "Major clinical readout or protocol change",
        "Regulatory feedback or label change",
        "HTA / payer evidence expectation change",
        "Competitor approval, readout, guideline, or pricing event",
        "New RWE dataset, registry, or evidence synthesis",
        "Budget, timeline, or operational feasibility change",
        "Regional launch sequencing or access need change"
    ], font_size=10.2)
    add_footer(slide, "17")


def raci():
    slide = blank_slide()
    add_title(slide, "RACI And Cross-Functional Ownership", "Clarify who is accountable, consulted, and informed for major evidence decisions.", "governance")
    headers = ["Activity / Decision", "Medical", "Clinical", "HEOR/RWE", "Market Access", "Regulatory", "Commercial / Comms"]
    add_table(slide, 0.55, 1.32, 12.25, 5.45, 8, 7, headers=headers, widths=[2.5, 1.55, 1.55, 1.65, 1.75, 1.6, 1.65], font_size=6.8)
    table = slide.shapes[-1].table
    rows = [
        ["Strategic imperatives", "A/R", "C", "C", "C", "C", "C"],
        ["Evidence gap prioritization", "A/R", "C", "R", "R", "C", "C"],
        ["Study concept approval", "A", "R/C", "R/C", "C", "C", "I"],
        ["HEOR / access evidence", "C", "C", "R", "A/R", "C", "I"],
        ["Regulatory evidence", "C", "A/R", "C", "C", "R", "I"],
        ["Publication / pull-through", "A/R", "C", "C", "C", "C", "R/C"],
        ["Plan refresh", "A/R", "C", "C", "C", "C", "C"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def refresh_cadence():
    slide = blank_slide()
    add_title(slide, "Refresh Cadence And Version Control", "Set an update rhythm and document trigger-based changes to preserve the plan as a living artifact.", "governance")
    add_placeholder(slide, 0.7, 1.42, 3.7, 1.6, "Annual planning cycle", ["[Portfolio strategy input]", "[Gap review]", "[Budget and roadmap approval]"], "teal")
    add_placeholder(slide, 4.65, 1.42, 3.7, 1.6, "Quarterly check-ins", ["[Milestone status]", "[Risk updates]", "[New evidence / landscape events]"], "blue")
    add_placeholder(slide, 8.6, 1.42, 3.55, 1.6, "Trigger events", ["[Readout, regulatory feedback, competitor event, HTA change, guideline update]"], "amber")
    headers = ["Version", "Date", "Change Trigger", "Major Updates", "Approved By"]
    add_table(slide, 0.7, 3.55, 11.45, 2.25, 5, 5, headers=headers, widths=[1.2, 1.45, 2.85, 4.0, 1.95], font_size=7.5)
    table = slide.shapes[-1].table
    for r in range(1, 5):
        vals = [f"v0.{r}", "[Date]", "[Trigger]", "[Summary of changed gaps, activities, owners, timing]", "[Forum]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def kpis():
    slide = blank_slide()
    add_title(slide, "IEGP KPIs And Impact Dashboard", "Track both plan execution and evidence impact, not only study completion.", "governance")
    metrics = [
        ("Gap closure", "[# high-priority gaps closed / total]", "teal"),
        ("On-time delivery", "[% milestones on track]", "blue"),
        ("Evidence reuse", "[# markets using global evidence]", "lime"),
        ("Pull-through", "[# claims, dossiers, publications supported]", "amber"),
    ]
    for i, (label, value, accent) in enumerate(metrics):
        add_placeholder(slide, 0.72 + i * 3.0, 1.5, 2.65, 1.15, label, [value], accent)
    add_table(slide, 0.72, 3.05, 11.75, 3.05, 6, 5, headers=["KPI", "Definition", "Baseline", "Target", "Owner"], widths=[2.0, 4.4, 1.55, 1.55, 1.85], font_size=7.8)
    table = slide.shapes[-1].table
    rows = [
        ["Evidence-gap closure", "[High-priority gaps with approved evidence solution]", "[ ]", "[ ]", "[ ]"],
        ["Decision readiness", "[Decision windows with complete evidence package]", "[ ]", "[ ]", "[ ]"],
        ["Publication / dissemination", "[Completed outputs vs plan]", "[ ]", "[ ]", "[ ]"],
        ["Market access support", "[HTA/payer submissions supported by roadmap evidence]", "[ ]", "[ ]", "[ ]"],
        ["Plan refresh quality", "[Updates completed after trigger events]", "[ ]", "[ ]", "[ ]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "18")


def risks():
    slide = blank_slide()
    add_title(slide, "Risks, Dependencies, And Trade-Offs", "Document constraints that could prevent evidence generation or reduce decision usefulness.", "governance")
    headers = ["Risk / Dependency", "Evidence Activity Affected", "Impact", "Mitigation", "Decision Needed"]
    add_table(slide, 0.65, 1.32, 12.0, 5.4, 8, 5, headers=headers, widths=[2.65, 2.35, 1.6, 3.25, 2.15], font_size=7.5)
    table = slide.shapes[-1].table
    for r in range(1, 8):
        row = ["[Risk or dependency]", "[Activity ID]", "[H/M/L]", "[Mitigation or alternate evidence path]", "[Decision / owner]"]
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "19")


def regional_appendix():
    slide = blank_slide()
    add_title(slide, "Appendix: Regional Adaptation Worksheet", "Use one copy per region or priority affiliate to capture local needs and adaptation decisions.", "appendix")
    add_placeholder(slide, 0.72, 1.35, 3.6, 1.0, "Region / affiliate", ["[Region and lead]"], "teal")
    add_placeholder(slide, 4.6, 1.35, 3.6, 1.0, "Local decision windows", ["[HTA, payer, guideline, launch]"], "blue")
    add_placeholder(slide, 8.48, 1.35, 3.6, 1.0, "Local constraints", ["[Comparator, dataset, policy, budget]"], "amber")
    headers = ["Local Need", "Global Evidence Reuse", "Local Gap", "Adaptation / New Tactic", "Affiliate Owner"]
    add_table(slide, 0.72, 2.8, 11.35, 3.15, 6, 5, headers=headers, widths=[2.7, 2.7, 2.05, 2.45, 1.45], font_size=7.4)
    table = slide.shapes[-1].table
    for r in range(1, 6):
        vals = ["[Need]", "[Evidence / asset]", "[Gap]", "[Adapt / generate / monitor]", "[Owner]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def evidence_source_log():
    slide = blank_slide()
    add_title(slide, "Appendix: Evidence Source Log", "Track source provenance, dates, confidence, and limitations for source-backed evidence claims.", "appendix")
    headers = ["Source ID", "Source / Dataset", "Date", "Owner", "Use In IEGP", "Confidence / Limitation"]
    add_table(slide, 0.62, 1.32, 12.1, 5.45, 8, 6, headers=headers, widths=[1.1, 2.75, 1.25, 1.55, 2.8, 2.65], font_size=7.0)
    table = slide.shapes[-1].table
    for r in range(1, 8):
        vals = [f"S{r}", "[Study, publication, registry, model, interview, internal source]", "[Date]", "[Owner]", "[Slide / claim / gap link]", "[Confidence and limitation]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def methodology():
    slide = blank_slide()
    add_title(slide, "Appendix: IEGP Development Methodology", "Document how the plan was created so future refreshes remain auditable.", "appendix")
    add_placeholder(slide, 0.7, 1.4, 5.45, 1.65, "Inputs reviewed", [
        "[TPP/TVP, clinical plans, publications, landscape, access insights, source inventory]"
    ], "teal")
    add_placeholder(slide, 6.55, 1.4, 5.45, 1.65, "Stakeholder engagement", [
        "[Functions interviewed, workshops held, affiliate input, external advisors]"
    ], "blue")
    add_placeholder(slide, 0.7, 3.45, 5.45, 1.65, "Scoring method", [
        "[Criteria, weights, evidence standards, tie-break rules]"
    ], "amber")
    add_placeholder(slide, 6.55, 3.45, 5.45, 1.65, "Assumptions and limits", [
        "[Known evidence gaps, unvalidated assumptions, data limitations, scope exclusions]"
    ], "rose")
    add_footer(slide, "00")


def glossary():
    slide = blank_slide()
    add_title(slide, "Appendix: Glossary", "Common terms for integrated evidence generation planning.", "appendix")
    headers = ["Term", "Definition / Template Usage", "Notes"]
    add_table(slide, 0.7, 1.32, 11.85, 5.45, 10, 3, headers=headers, widths=[1.65, 7.3, 2.9], font_size=7.4)
    table = slide.shapes[-1].table
    rows = [
        ["IEGP / IEP", "Integrated evidence generation/evidence plan across lifecycle, stakeholders, and functions", "[Use organization terminology]"],
        ["ORQ", "Open research question converted from a priority evidence gap", "[Often PICOT-framed]"],
        ["PICOT", "Population, intervention, comparator, outcome, timing", "[For answerable questions]"],
        ["RWE", "Real-world evidence from data outside conventional randomized trials", "[Claims/EHR/registry/prospective]"],
        ["HEOR", "Health economics and outcomes research", "[SLR, models, BIA, burden]"],
        ["HTA", "Health technology assessment", "[Market access decision body]"],
        ["COA / PRO", "Clinical outcome assessment / patient-reported outcome", "[Patient-centered evidence]"],
        ["IIR / IIT", "Investigator-initiated research/trial", "[External research]"],
        ["Pull-through", "Use of generated evidence in publications, dossiers, claims, medical exchange, or decisions", "[Track as KPI]"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "00")


def source_basis():
    slide = blank_slide()
    add_title(slide, "Source-Informed Template Basis", "This template combines public IEGP/IEP methods, case studies, and evidence-planning guidance.", "appendix")
    add_table(slide, 0.65, 1.32, 12.0, 5.4, 9, 3, headers=["Source", "Template Contribution", "Source Type"], widths=[3.25, 6.65, 2.1], font_size=7.4)
    rows = [
        ["Lumanity", "Evidence needed vs available, future landscape, gap prioritization, global harmonization, KPIs", "Framework / PDF"],
        ["Cencora", "Market access, RWE, HEOR, SLR/evidence-synthesis components", "Service / article"],
        ["ZS", "IEP deliverable structure, 3-5 year roadmap, governance and execution discipline", "How-to article"],
        ["Veranex", "Claims matrix, scorecard, roadmap, budget, ROI dashboard concepts", "Framework"],
        ["DiMe", "Stage-based evidence plan resources and concrete DHT case examples", "Toolkit / case PDFs"],
        ["MAPS / Prescient", "Medical affairs standards, stakeholder relevance, adaptive evidence plans, refresh cadence", "Guidance / whitepaper"],
        ["McKinsey / Springer-Tufts", "Evidence catalog, prioritization economics, value-driver framing", "Thought leadership / academic"],
        ["Oxford PharmaGenesis", "PICOT refinement, data source selection, study mapping workflow", "Process model"],
    ]
    table = slide.shapes[-1].table
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    add_footer(slide, "20")


def appendix_tracker():
    slide = blank_slide()
    add_title(slide, "Appendix: Claims And Evidence Linkage", "Optional table for connecting evidence outputs to claims, messages, and local adaptations.", "appendix")
    headers = ["Claim / Message", "Evidence Source", "Evidence Strength", "Allowed Use", "Market / Audience", "Gap / Next Step"]
    add_table(slide, 0.6, 1.32, 12.1, 5.5, 8, 6, headers=headers, widths=[2.25, 2.3, 1.65, 2.0, 1.75, 2.15], font_size=7.1)
    table = slide.shapes[-1].table
    for r in range(1, 8):
        vals = ["[Claim]", "[Study / source]", "[Strong / moderate / weak]", "[Label / medical / payer / internal]", "[Audience]", "[Gap or action]"]
        for c, text in enumerate(vals):
            table.cell(r, c).text = text
    add_footer(slide, "21")


def final_checklist():
    slide = blank_slide()
    add_title(slide, "Final IEGP Review Checklist", "Use this slide before governance review or funding decisions.", "appendix")
    add_placeholder(slide, 0.75, 1.42, 5.5, 2.1, "Content completeness", [
        "Strategic imperatives tied to decisions",
        "Stakeholder evidence needs documented",
        "Current evidence inventory reviewed",
        "Evidence gaps prioritized using agreed criteria"
    ], "teal")
    add_placeholder(slide, 6.65, 1.42, 5.5, 2.1, "Execution readiness", [
        "Each activity has owner, timing, status, and decision use",
        "Budget class and dependencies are known",
        "Pull-through plan and claims linkage documented",
        "Governance cadence and update triggers agreed"
    ], "blue")
    add_placeholder(slide, 0.75, 4.0, 5.5, 1.6, "Evidence integrity", [
        "Source-backed facts are separated from assumptions",
        "Evidence limitations and quality notes retained"
    ], "amber")
    add_placeholder(slide, 6.65, 4.0, 5.5, 1.6, "Decision", [
        "[Approve roadmap]  [Revise]  [Defer]",
        "Next review date: [ ]"
    ], "lime")
    add_footer(slide, "22")


for fn in [
    cover,
    how_to_use,
    evidence_model,
    executive_summary,
    key_decisions,
    asset_snapshot,
    strategic_imperatives,
    value_narrative,
    milestones,
    stakeholder_needs,
    patient_journey,
    geography_needs,
    inventory,
    evidence_quality,
    need_have_generate,
    future_landscape,
    gap_matrix,
    prioritization,
    priority_gap_deep_dive,
    research_questions,
    evidence_options,
    roadmap,
    activity_tracker,
    claims_matrix,
    study_concept,
    heor_access_plan,
    rwe_strategy,
    patient_centered_evidence,
    regulatory_considerations,
    med_comms,
    external_research_iir,
    budget_resources,
    global_local,
    specialized_workstreams,
    governance,
    raci,
    refresh_cadence,
    kpis,
    risks,
    regional_appendix,
    evidence_source_log,
    methodology,
    source_basis,
    appendix_tracker,
    glossary,
    final_checklist,
]:
    fn()

PRS.save(OUT_FILE)
print(OUT_FILE.resolve())
