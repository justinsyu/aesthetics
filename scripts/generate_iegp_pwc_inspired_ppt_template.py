from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("outputs/iegp_powerpoint_template")
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "IEGP_PwC_Inspired_Editorial_Template.pptx"

PRS = Presentation()
PRS.slide_width = Inches(13.333)
PRS.slide_height = Inches(7.5)
W, H = PRS.slide_width, PRS.slide_height

COLORS = {
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
    "charcoal": RGBColor(45, 45, 45),
    "gray": RGBColor(125, 125, 125),
    "light": RGBColor(222, 222, 222),
    "wash": RGBColor(245, 245, 243),
    "yellow": RGBColor(255, 182, 0),
    "tangerine": RGBColor(235, 140, 0),
    "orange": RGBColor(208, 74, 2),
    "rose": RGBColor(219, 83, 106),
    "red": RGBColor(224, 48, 30),
    "maroon": RGBColor(96, 34, 43),
}


def c(name):
    return COLORS[name]


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def line(shape, color, width=0.7):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def slide_bg(color="white"):
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    fill(bg, c(color))
    bg.line.fill.background()
    return slide


def text(slide, x, y, w, h, value, size=12, font="Arial", bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.01)
    tf.margin_right = Inches(0.01)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.text = value
    if align:
        p.alignment = align
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or c("black")
    return box


def body(slide, x, y, w, h, lines, size=10, color=None, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, value in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = value
        p.space_after = Pt(4)
        if bullet:
            p._p.get_or_add_pPr().set("marL", "171450")
            p._p.get_or_add_pPr().set("indent", "-114300")
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(size)
            r.font.color.rgb = color or c("charcoal")
    return box


def block(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, c(color))
    shp.line.fill.background()
    return shp


def block_stack(slide, x=11.35, y=0.35, scale=1.0, variant=0):
    palettes = [
        ["yellow", "tangerine", "orange", "rose", "red"],
        ["orange", "tangerine", "yellow", "maroon", "rose"],
        ["rose", "red", "orange", "yellow", "tangerine"],
    ]
    pal = palettes[variant % len(palettes)]
    specs = [
        (0.00, 0.00, 0.44, 0.44, pal[0]),
        (0.36, 0.20, 0.50, 0.50, pal[1]),
        (0.08, 0.55, 0.62, 0.38, pal[2]),
        (0.76, 0.04, 0.36, 0.74, pal[3]),
        (0.62, 0.78, 0.54, 0.28, pal[4]),
    ]
    for dx, dy, w, h, col in specs:
        block(slide, x + dx * scale, y + dy * scale, w * scale, h * scale, col)


def footer(slide, section, product="[Product / indication]"):
    text(slide, 0.55, 7.08, 4.6, 0.18, f"IEGP | {section} | {product}", size=6.8, color=c("gray"))
    text(slide, 11.9, 7.08, 0.6, 0.18, f"{len(PRS.slides):02d}", size=6.8, color=c("gray"), align=PP_ALIGN.RIGHT)


def title(slide, title_text, section, subtitle=None, variant=0):
    text(slide, 0.58, 0.30, 4.8, 0.17, f"IEGP | {section}".upper(), size=6.8, font="Arial", bold=True, color=c("orange"))
    text(slide, 0.56, 0.65, 9.6, 0.62, title_text, size=26, font="Georgia", bold=False, color=c("black"))
    if subtitle:
        text(slide, 0.6, 1.24, 9.6, 0.32, subtitle, size=10.5, color=c("charcoal"))
    rule = block(slide, 0.58, 1.68, 11.85, 0.012, "light")
    block_stack(slide, 11.35, 0.28, 0.82, variant)
    footer(slide, section)


def panel(slide, x, y, w, h, header=None, color="orange", bg="wash"):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(rect, c(bg))
    line(rect, c("light"), 0.6)
    block(slide, x, y, 0.08, h, color)
    if header:
        text(slide, x + 0.2, y + 0.14, w - 0.35, 0.18, header, size=8.2, font="Arial", bold=True, color=c("black"))
    return rect


def table(slide, x, y, w, h, rows, cols, headers, widths=None, font_size=6.8, header_color="orange"):
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    if widths:
        for i, width in enumerate(widths):
            tbl.columns[i].width = Inches(width)
    for r in range(rows):
        for col in range(cols):
            cell = tbl.cell(r, col)
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = c(header_color if r == 0 else "white")
            cell.text = headers[col] if r == 0 else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0
                    run.font.color.rgb = c("white") if r == 0 else c("charcoal")
    return shape


def big_number(slide, x, y, num, label, color="orange"):
    text(slide, x, y, 1.55, 0.45, num, size=28, font="Georgia", bold=False, color=c(color))
    block(slide, x, y + 0.62, 1.25, 0.045, color)
    text(slide, x, y + 0.78, 2.2, 0.32, label, size=7.8, color=c("charcoal"))


def cover():
    slide = slide_bg("black")
    # warm dot field approximated by editable squares
    for i in range(20):
        for j in range(9):
            if (i + j) % 3 != 0:
                continue
            x = 7.0 + i * 0.24
            y = 0.7 + j * 0.32 + i * 0.05
            if x > 12.5 or y > 6.8:
                continue
            col = ["yellow", "tangerine", "orange", "rose", "red"][(i + j) % 5]
            block(slide, x, y, 0.12, 0.12, col)
    block(slide, 0.7, 0.72, 1.2, 0.35, "orange")
    text(slide, 0.82, 0.82, 0.95, 0.09, "TEMPLATE", size=6.5, bold=True, color=c("white"))
    text(slide, 0.68, 1.55, 6.3, 1.08, "Integrated Evidence Generation Plan", size=32, font="Georgia", color=c("white"))
    text(slide, 0.72, 3.02, 5.3, 0.42, "Editorial consulting-style PowerPoint starting point for evidence strategy, gap prioritization, and lifecycle execution.", size=12, color=RGBColor(225, 225, 225))
    block_stack(slide, 0.72, 5.4, 1.0, 1)
    text(slide, 2.15, 5.62, 4.8, 0.22, "[Product] | [Indication] | [Region]", size=12, bold=True, color=c("white"))
    text(slide, 2.15, 6.0, 4.0, 0.16, "Owner: [team]  |  Version: [v0.1]  |  Date: [date]", size=7.2, color=RGBColor(210, 210, 210))


def tokens():
    slide = slide_bg()
    title(slide, "Brand adaptation tokens", "Setup", "Replace only these accent tokens to adapt the deck to a company palette.", 1)
    rows = [("Primary accent", "orange"), ("Secondary accent", "tangerine"), ("Warm highlight", "yellow"), ("Risk / gap", "red"), ("Neutral dark", "charcoal"), ("Neutral light", "light")]
    for i, (label, col) in enumerate(rows):
        x = 0.72 + (i % 3) * 3.8
        y = 2.2 + (i // 3) * 1.38
        block(slide, x, y, 0.78, 0.78, col)
        text(slide, x + 0.95, y + 0.13, 1.7, 0.18, label, size=8.5, bold=True)
        text(slide, x + 0.95, y + 0.45, 1.8, 0.15, f"RGB {c(col)[0]}, {c(col)[1]}, {c(col)[2]}", size=6.7, color=c("gray"))
    panel(slide, 0.72, 5.3, 11.4, 0.8, "Rule", "orange", "white")
    text(slide, 1.0, 5.62, 10.7, 0.2, "Use accent color for hierarchy and meaning, not decoration. Keep backgrounds mostly white and preserve table readability.", size=11.3, bold=True)


def section_divider(name, section, note, variant=0):
    slide = slide_bg("white")
    block_stack(slide, 0.72, 0.72, 1.25, variant)
    text(slide, 0.72, 2.0, 8.4, 0.75, name, size=42, font="Georgia")
    text(slide, 0.78, 3.05, 6.9, 0.38, note, size=13, color=c("charcoal"))
    block(slide, 0.78, 4.35, 5.9, 0.055, "orange")
    text(slide, 0.78, 6.92, 2.0, 0.15, section.upper(), size=7.2, bold=True, color=c("gray"))
    footer(slide, section)


def executive():
    slide = slide_bg()
    title(slide, "Executive summary", "Briefing", "Three evidence priorities, their consequences, and the actions needed now.", 0)
    for i, (h, col) in enumerate([("Priority 1", "orange"), ("Priority 2", "tangerine"), ("Priority 3", "red")]):
        x = 0.72 + i * 4.0
        block(slide, x, 2.08, 3.25, 0.2, col)
        text(slide, x, 2.45, 2.2, 0.22, h, size=15, font="Georgia")
        body(slide, x, 2.95, 3.1, 1.95, [
            "Gap: [top evidence gap]",
            "Action: [study / analysis / synthesis]",
            "Decision impact: [reg/access/adoption]",
            "Timing: [decision window]"
        ], size=9.2)
    panel(slide, 0.72, 5.55, 11.35, 0.82, "Decision ask", "orange", "white")
    text(slide, 1.0, 5.87, 10.5, 0.2, "[Approve roadmap]  |  [Fund priority evidence activities]  |  [Confirm owners and refresh cadence]", size=11.3, bold=True)


def decisions():
    slide = slide_bg()
    title(slide, "Decisions needed", "Briefing", "A governance-ready view of choices, recommendations, and timing.", 2)
    shp = table(slide, 0.62, 2.02, 12.05, 4.65, 7, 5, ["Decision", "Recommendation", "Rationale", "Owner", "Deadline"], [2.55, 2.0, 4.2, 1.55, 1.75], 7.2, "maroon")
    rows = [
        ["Approve priority gaps", "[approve / revise]", "[risk and value rationale]", "[owner]", "[date]"],
        ["Fund evidence activity", "[approve / defer]", "[decision value]", "[owner]", "[date]"],
        ["Select study design", "[option A/B]", "[credibility and feasibility]", "[owner]", "[date]"],
        ["Resolve local adaptation", "[global/local split]", "[reuse vs local gap]", "[owner]", "[date]"],
        ["Confirm governance", "[RACI/cadence]", "[execution accountability]", "[owner]", "[date]"],
        ["Other", "[recommendation]", "[basis]", "[owner]", "[date]"],
    ]
    for r, row in enumerate(rows, 1):
        for col, val in enumerate(row):
            shp.table.cell(r, col).text = val


def asset():
    slide = slide_bg()
    title(slide, "Asset snapshot", "Context", "Frame the asset before defining stakeholder evidence needs.", 1)
    fields = [
        ("Asset / mechanism", "[name, modality, MoA]"),
        ("Population", "[indication, line, biomarker]"),
        ("Lifecycle stage", "[pre-launch / launch / post-launch]"),
        ("Geography", "[global, US, EU5, JP, local markets]"),
        ("TPP / TVP", "[clinical and value targets]"),
        ("Key milestones", "[readouts, submissions, HTA, launch]")
    ]
    for i, (h, b) in enumerate(fields):
        x = 0.72 + (i % 3) * 3.92
        y = 2.05 + (i // 3) * 1.85
        panel(slide, x, y, 3.35, 1.18, h, ["orange", "tangerine", "red"][i % 3], "white")
        text(slide, x + 0.22, y + 0.58, 2.75, 0.26, b, size=10.5, color=c("charcoal"))


def value_story():
    slide = slide_bg()
    title(slide, "Value narrative evidence story", "Context", "Trace each value message to evidence required, current support, gaps, and tactics.", 0)
    labels = ["Value driver", "Evidence required", "Current support", "Gap", "Planned tactic"]
    colors = ["orange", "tangerine", "yellow", "red", "maroon"]
    for i, label in enumerate(labels):
        x = 0.85 + i * 2.35
        block(slide, x, 2.25, 1.55, 0.32, colors[i])
        text(slide, x, 2.75, 1.55, 0.25, label, size=10.5, font="Georgia")
        panel(slide, x, 3.35, 1.55, 1.25, None, colors[i], "white")
        text(slide, x + 0.16, 3.72, 1.15, 0.34, "[content]", size=9, color=c("charcoal"))
        if i < len(labels) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.65), Inches(3.95), Inches(x + 2.2), Inches(3.95))
            line(conn, c("light"), 1.1)
    text(slide, 0.85, 5.55, 10.8, 0.24, "Narrative test: Can every value message be supported by source-backed evidence or a funded activity?", size=13.5, font="Georgia")


def stakeholder():
    slide = slide_bg()
    title(slide, "Stakeholder evidence needs", "Needs", "Compact map of decision makers, their questions, and evidence standards.", 2)
    positions = [
        ("Regulators", 0.78, 2.05, "orange"),
        ("HTA / payers", 4.35, 2.05, "tangerine"),
        ("HCPs / KOLs", 7.92, 2.05, "yellow"),
        ("Patients", 0.78, 4.35, "red"),
        ("Policy / systems", 4.35, 4.35, "maroon"),
        ("Internal teams", 7.92, 4.35, "orange"),
    ]
    for name, x, y, col in positions:
        block(slide, x, y, 2.85, 0.28, col)
        text(slide, x, y + 0.52, 2.5, 0.2, name, size=13, font="Georgia")
        body(slide, x, y + 0.95, 2.75, 0.85, ["Decision: [ ]", "Evidence standard: [ ]", "Timing: [ ]"], size=8.3)


def inventory():
    slide = slide_bg()
    title(slide, "Current evidence inventory", "Needs", "Catalogue available and planned evidence before defining new evidence to generate.", 0)
    shp = table(slide, 0.55, 2.02, 12.25, 4.85, 9, 6, ["Source", "Status", "Population", "Endpoint / output", "Stakeholder", "Limits"], [2.0, 1.2, 2.0, 2.75, 1.65, 2.65], 6.6, "orange")
    rows = ["Clinical trials", "RWE", "HEOR model", "SLR / ITC", "COA / PRO", "Epidemiology", "Publications", "Other"]
    for r, label in enumerate(rows, 1):
        vals = [label, "[status]", "[population]", "[output]", "[stakeholder]", "[quality / usability limits]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def gap_matrix():
    slide = slide_bg()
    title(slide, "Evidence gap matrix", "Gaps", "Prioritize by stakeholder decision, current coverage, consequence, and action.", 2)
    shp = table(slide, 0.55, 2.02, 10.3, 4.9, 8, 5, ["Gap", "Decision", "Current evidence", "Priority", "Action"], [2.25, 2.25, 2.65, 1.0, 2.15], 6.8, "red")
    for r in range(1, 8):
        vals = [f"G{r}: [gap]", "[decision]", "[evidence]", "[P1/P2/P3]", "[tactic]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val
    panel(slide, 11.12, 2.02, 1.35, 4.9, "Consequence", "red", "wash")
    body(slide, 11.32, 2.65, 0.95, 3.1, ["[label delay]", "[HTA risk]", "[slow uptake]", "[unfunded lifecycle]"], size=7.4)


def prioritization():
    slide = slide_bg()
    title(slide, "Gap prioritization scorecard", "Gaps", "Use a transparent rubric before converting gaps into funded activities.", 1)
    for i, (num, label, col) in enumerate([("5", "strategic impact", "orange"), ("4", "urgency", "tangerine"), ("3", "feasibility", "yellow"), ("2", "reuse", "red")]):
        big_number(slide, 0.82 + i * 2.8, 2.05, num, label, col)
    shp = table(slide, 0.78, 4.0, 11.55, 2.45, 5, 5, ["Gap", "Impact", "Urgency", "Feasibility", "Decision"], [3.0, 1.7, 1.7, 1.9, 3.25], 7.2, "maroon")
    for r in range(1, 5):
        vals = [f"G{r}: [gap]", "[1-5]", "[1-5]", "[1-5]", "[fund / revise / hold]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def roadmap():
    slide = slide_bg()
    title(slide, "Integrated evidence roadmap", "Roadmap", "Function-level swimlanes with readouts, decisions, and pull-through milestones.", 0)
    quarters = ["Y1 Q1", "Q2", "Q3", "Q4", "Y2 Q1", "Q2", "Q3", "Q4"]
    for i, q in enumerate(quarters):
        text(slide, 3.15 + i * 1.0, 2.0, 0.75, 0.15, q, size=6.5, bold=True, color=c("gray"), align=PP_ALIGN.CENTER)
    lanes = [("Clinical / regulatory", "orange"), ("RWE", "tangerine"), ("HEOR / access", "yellow"), ("Patient / COA", "red"), ("Publications", "maroon")]
    for i, (lane, col) in enumerate(lanes):
        y = 2.45 + i * 0.75
        text(slide, 0.76, y + 0.07, 2.0, 0.16, lane, size=7.2, bold=True)
        block(slide, 2.95, y + 0.17, 8.8, 0.015, "light")
        block(slide, 3.2 + (i % 3) * 0.55, y, 1.65 + (i % 2) * 0.55, 0.32, col)
        text(slide, 3.32 + (i % 3) * 0.55, y + 0.09, 1.2, 0.08, "[activity]", size=5.8, bold=True, color=c("black"), align=PP_ALIGN.CENTER)
        block(slide, 5.75 + i * 0.75, y - 0.04, 0.18, 0.4, "red")
    panel(slide, 0.76, 6.35, 11.65, 0.44, "Decision gates: [protocol] | [data cut] | [abstract] | [regulatory] | [HTA] | [launch]", "orange", "white")


def activity_tracker():
    slide = slide_bg()
    title(slide, "Evidence activity tracker", "Roadmap", "Single working table for studies, analyses, syntheses, and dissemination activities.", 2)
    shp = table(slide, 0.48, 2.02, 12.35, 4.85, 9, 8, ["ID", "Activity", "Gap", "Method / data", "Owner", "Timing", "Status", "Priority"], [0.55, 1.9, 0.8, 2.6, 1.1, 1.25, 1.7, 1.0], 6.0, "orange")
    for r in range(1, 9):
        vals = [f"A{r}", "[study / analysis]", "G[ ]", "[design / source]", "[owner]", "[start-readout]", "[idea/planned/active]", "[P1]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def study_concept():
    slide = slide_bg()
    title(slide, "Study concept one-pager", "Roadmap", "Duplicate this page for each priority evidence-generation activity.", 1)
    items = [
        ("Study / activity", "[name and ID]", "orange"),
        ("Gap and decision use", "[gap, ORQ, stakeholder]", "tangerine"),
        ("Objective", "[objective and hypothesis]", "yellow"),
        ("Design", "[population, comparator, endpoint, data]", "orange"),
        ("Feasibility", "[cost, complexity, dependency]", "red"),
        ("Pull-through", "[publication, dossier, claims]", "maroon"),
    ]
    for i, (h, b, col) in enumerate(items):
        x = 0.72 + (i % 3) * 3.92
        y = 2.05 + (i // 3) * 1.75
        panel(slide, x, y, 3.35, 1.08, h, col, "white")
        text(slide, x + 0.22, y + 0.56, 2.75, 0.22, b, size=9.3, color=c("charcoal"))
    panel(slide, 0.72, 6.0, 11.35, 0.46, "Decision: [approve / revise / defer]    Next step: [action]    Date: [ ]", "orange", "white")


def workstreams():
    slide = slide_bg()
    title(slide, "Specialized evidence workstreams", "Roadmap", "Ensure HEOR, RWE, patient-centered, regulatory, and communications evidence are explicit.", 0)
    work = [
        ("HEOR / access", ["SLR / ITC", "economic model", "budget impact", "payer research"], "tangerine"),
        ("RWE / data", ["retrospective database", "prospective study", "registry", "external control"], "orange"),
        ("Patient-centered", ["COA / PRO", "HRQoL", "preference", "caregiver burden"], "red"),
        ("Publications / comms", ["abstract", "manuscript", "field medical", "claims matrix"], "maroon"),
    ]
    for i, (h, lines, col) in enumerate(work):
        x = 0.78 + (i % 2) * 5.75
        y = 2.0 + (i // 2) * 2.05
        panel(slide, x, y, 5.15, 1.55, h, col, "white")
        body(slide, x + 0.24, y + 0.52, 4.35, 0.65, lines, size=8.8)


def governance():
    slide = slide_bg()
    title(slide, "Governance and refresh cadence", "Control", "Operate the IEGP as a living plan, not a one-time deliverable.", 2)
    shp = table(slide, 0.7, 2.02, 6.3, 4.55, 7, 4, ["Forum", "Purpose", "Cadence", "Rights"], [1.45, 2.45, 1.1, 1.3], 6.8, "maroon")
    rows = [["Core team", "[maintain plan]", "[monthly]", "[recommend]"], ["Evidence council", "[prioritize/fund]", "[quarterly]", "[approve]"], ["Global-local", "[adapt needs]", "[quarterly]", "[recommend]"], ["Leadership", "[trade-offs]", "[semiannual]", "[approve]"], ["Publications", "[pull-through]", "[monthly]", "[recommend]"], ["Other", "[purpose]", "[cadence]", "[rights]"]]
    for r, row in enumerate(rows, 1):
        for col, val in enumerate(row):
            shp.table.cell(r, col).text = val
    panel(slide, 7.42, 2.02, 4.85, 4.55, "Update triggers", "red", "wash")
    body(slide, 7.72, 2.6, 4.0, 2.8, ["Clinical readout", "Regulatory feedback", "HTA / payer change", "Competitor approval or data", "Guideline update", "Budget or feasibility change"], size=10.2)


def kpis():
    slide = slide_bg()
    title(slide, "Evidence impact dashboard", "Control", "Track plan execution and evidence impact with concise business-readable measures.", 1)
    metrics = [("Gap closure", "72%", "orange"), ("On track", "18/24", "tangerine"), ("Evidence reuse", "6 mkts", "yellow"), ("Pull-through", "12", "red")]
    for i, (label, num, col) in enumerate(metrics):
        big_number(slide, 0.8 + i * 2.8, 2.0, num, label, col)
    shp = table(slide, 0.72, 4.0, 11.5, 2.45, 5, 5, ["KPI", "Definition", "Baseline", "Target", "Owner"], [2.0, 4.3, 1.5, 1.5, 2.2], 7.2, "orange")
    for r in range(1, 5):
        vals = ["[KPI]", "[definition]", "[baseline]", "[target]", "[owner]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def risk():
    slide = slide_bg()
    title(slide, "Risk and dependency register", "Control", "Document constraints that could reduce evidence usefulness or delay decisions.", 2)
    shp = table(slide, 0.62, 2.02, 12.0, 4.8, 8, 5, ["Risk / dependency", "Activity", "Impact", "Mitigation", "Decision needed"], [2.75, 1.85, 1.25, 3.6, 2.55], 6.9, "red")
    for r in range(1, 8):
        vals = ["[risk]", "[activity]", "[H/M/L]", "[mitigation or alternate evidence path]", "[decision / owner]"]
        for col, val in enumerate(vals):
            shp.table.cell(r, col).text = val


def appendix():
    slide = slide_bg()
    title(slide, "Appendix worksheets", "Appendix", "Keep source provenance, regional adaptation, methodology, and final review audit-ready.", 0)
    for i, (h, desc, col) in enumerate([
        ("Regional adaptation", "local needs, reusable global evidence, local gaps, owner", "orange"),
        ("Evidence source log", "source, date, owner, confidence, limitation, claim link", "tangerine"),
        ("Methodology", "inputs, workshops, scoring rubric, assumptions, exclusions", "red"),
        ("Glossary", "IEGP/IEP, ORQ, PICOT, RWE, HEOR, HTA, COA/PRO, IIR", "maroon"),
    ]):
        x = 0.78 + (i % 2) * 5.7
        y = 2.05 + (i // 2) * 1.85
        panel(slide, x, y, 5.05, 1.2, h, col, "white")
        text(slide, x + 0.24, y + 0.58, 4.35, 0.24, desc, size=9.2, color=c("charcoal"))


slides = [
    cover,
    tokens,
    lambda: section_divider("Briefing", "Briefing", "Set the evidence priorities and the decisions needed now.", 0),
    executive,
    decisions,
    lambda: section_divider("Context", "Context", "Frame the asset, value story, and stakeholder landscape.", 1),
    asset,
    value_story,
    stakeholder,
    inventory,
    lambda: section_divider("Gap analysis", "Gaps", "Convert evidence needs into prioritized, actionable gaps.", 2),
    gap_matrix,
    prioritization,
    lambda: section_divider("Roadmap", "Roadmap", "Translate priority gaps into activities, owners, timing, and pull-through.", 0),
    roadmap,
    activity_tracker,
    study_concept,
    workstreams,
    lambda: section_divider("Control", "Control", "Run the IEGP as a living governance and impact system.", 1),
    governance,
    kpis,
    risk,
    appendix,
]

for make in slides:
    make()

PRS.save(OUT_FILE)
print(OUT_FILE.resolve())
