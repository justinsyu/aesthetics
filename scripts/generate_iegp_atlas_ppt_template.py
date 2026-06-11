from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("outputs/iegp_powerpoint_template")
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "IEGP_Evidence_Atlas_Design_Template.pptx"

PRS = Presentation()
PRS.slide_width = Inches(13.333)
PRS.slide_height = Inches(7.5)
W = PRS.slide_width
H = PRS.slide_height


COLORS = {
    "paper": RGBColor(248, 247, 242),
    "porcelain": RGBColor(253, 252, 248),
    "ink": RGBColor(31, 34, 38),
    "soft_ink": RGBColor(80, 86, 93),
    "muted": RGBColor(126, 132, 137),
    "line": RGBColor(214, 214, 205),
    "mist": RGBColor(238, 238, 230),
    "primary": RGBColor(20, 116, 123),
    "secondary": RGBColor(156, 128, 68),
    "signal": RGBColor(176, 199, 80),
    "risk": RGBColor(180, 70, 80),
    "night": RGBColor(25, 31, 39),
    "white": RGBColor(255, 255, 255),
}


def c(name):
    return COLORS[name]


def blank():
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = c("paper")
    bg.line.fill.background()
    return slide


def fill(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def line(shape, color, width=0.8):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def txt(slide, x, y, w, h, text, size=12, bold=False, color=None, align=None, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or c("ink")
    return box


def body(slide, x, y, w, h, lines, size=9.5, color=None, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(3)
        if bullet:
            p._p.get_or_add_pPr().set("marL", "171450")
            p._p.get_or_add_pPr().set("indent", "-114300")
        for r in p.runs:
            r.font.name = "Aptos"
            r.font.size = Pt(size)
            r.font.color.rgb = color or c("ink")
    return box


def rail(slide, section="A", page=None, dark=False):
    rail_color = c("paper") if dark else c("night")
    muted = RGBColor(214, 219, 216) if dark else c("muted")
    x = 0.28
    main = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(0.58), Inches(0.012), Inches(6.15))
    fill(main, muted)
    main.line.fill.background()
    for i in range(9):
        tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x - 0.055), Inches(0.75 + i * 0.7), Inches(0.11), Inches(0.012))
        fill(tick, muted)
        tick.line.fill.background()
    txt(slide, 0.12, 0.25, 0.35, 0.2, section, size=8, bold=True, color=rail_color, align=PP_ALIGN.CENTER, font="Consolas")
    txt(slide, 0.1, 6.92, 0.45, 0.18, page or f"{len(PRS.slides):02d}", size=7, color=muted, align=PP_ALIGN.CENTER, font="Consolas")


def title(slide, heading, kicker, section="A"):
    rail(slide, section)
    txt(slide, 0.74, 0.34, 1.5, 0.18, kicker.upper(), size=6.7, bold=True, color=c("primary"), font="Consolas")
    txt(slide, 0.74, 0.62, 10.8, 0.48, heading, size=22.5, bold=True, color=c("ink"))
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.74), Inches(1.17), Inches(11.9), Inches(0.01))
    fill(rule, c("line"))
    rule.line.fill.background()


def panel(slide, x, y, w, h, title_text=None, accent="primary", fill_name="porcelain"):
    r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(r, c(fill_name))
    line(r, c("line"), 0.65)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.055), Inches(h))
    fill(band, c(accent))
    band.line.fill.background()
    if title_text:
        txt(slide, x + 0.18, y + 0.12, w - 0.32, 0.22, title_text, size=8.7, bold=True, color=c("ink"), font="Consolas")
    return r


def pin(slide, x, y, label, color="primary"):
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.28), Inches(0.28))
    fill(circ, c(color))
    circ.line.fill.background()
    txt(slide, x + 0.035, y + 0.055, 0.21, 0.1, label, size=6.2, bold=True, color=c("white"), align=PP_ALIGN.CENTER, font="Consolas")
    return circ


def signal_band(slide, x, y, w, label, color="primary"):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.075))
    fill(bar, c(color))
    bar.line.fill.background()
    txt(slide, x, y + 0.12, w, 0.14, label, size=6.5, color=c("muted"), align=PP_ALIGN.CENTER, font="Consolas")


def atlas_table(slide, x, y, w, h, rows, cols, headers, widths=None, font_size=6.7):
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for r in range(rows):
        for col in range(cols):
            cell = table.cell(r, col)
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            cell.margin_top = Inches(0.025)
            cell.margin_bottom = Inches(0.025)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = c("night") if r == 0 else (c("porcelain") if r % 2 else c("paper"))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = c("white") if r == 0 else c("ink")
                    run.font.bold = r == 0
            cell.text = headers[col] if r == 0 else ""
    return shape


def contour(slide, x, y, w, h, dark=False):
    color = RGBColor(74, 84, 89) if dark else c("line")
    for i in range(7):
        ox = i * 0.11
        oy = i * 0.08
        free = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x + ox), Inches(y + oy), Inches(w - ox * 2), Inches(h - oy * 2))
        free.rotation = i * 13
        free.fill.background()
        line(free, color, 0.55)


def cover():
    slide = blank()
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.65), 0, Inches(5.68), H)
    fill(block, c("night"))
    block.line.fill.background()
    contour(slide, 8.15, 0.75, 4.2, 4.8, dark=True)
    rail(slide, "00")
    txt(slide, 0.78, 0.72, 1.4, 0.2, "EVIDENCE ATLAS", size=7, bold=True, color=c("primary"), font="Consolas")
    txt(slide, 0.78, 1.18, 6.55, 1.35, "Integrated Evidence Generation Plan", size=34, bold=True)
    txt(slide, 0.82, 2.72, 5.85, 0.44, "A reusable, brand-adaptable PowerPoint system for mapping evidence, decisions, gaps, and lifecycle execution.", size=13.5, color=c("soft_ink"))
    for i, text in enumerate(["STRATEGY", "NEEDS", "GAPS", "ROADMAP", "CONTROL"]):
        pin(slide, 0.85 + i * 1.26, 4.05, str(i + 1), ["primary", "secondary", "risk", "signal", "primary"][i])
        txt(slide, 1.18 + i * 1.26, 4.1, 0.85, 0.12, text, size=5.5, bold=True, color=c("muted"), font="Consolas")
    panel(slide, 0.82, 5.58, 5.85, 0.85, None, "primary", "porcelain")
    txt(slide, 1.05, 5.82, 5.1, 0.22, "[Product] | [Indication] | [Lifecycle stage] | [Region]", size=12.5, bold=True)
    txt(slide, 1.05, 6.18, 4.5, 0.16, "Owner: [team]  |  Version: [v0.1]  |  Date: [date]", size=7.5, color=c("muted"), font="Consolas")
    txt(slide, 8.18, 6.45, 3.6, 0.24, "Designed as editable PowerPoint objects", size=9.5, color=RGBColor(226, 229, 225), font="Consolas")


def design_tokens():
    slide = blank()
    title(slide, "Design Tokens: Swap These To Brand The Deck", "setup", "01")
    body(slide, 0.78, 1.48, 4.2, 0.55, ["Replace the two accent colors and one risk color with company palette equivalents. The neutral system should remain stable for readability."], size=10, color=c("soft_ink"))
    swatches = [("Primary accent", "primary"), ("Secondary accent", "secondary"), ("Signal / positive", "signal"), ("Risk / urgency", "risk"), ("Ink", "ink"), ("Paper", "paper")]
    for i, (label, col) in enumerate(swatches):
        x = 0.82 + (i % 3) * 2.25
        y = 2.35 + (i // 3) * 1.35
        sq = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.72), Inches(0.72))
        fill(sq, c(col))
        line(sq, c("line"), 0.55)
        txt(slide, x + 0.85, y + 0.12, 1.25, 0.18, label, size=7.8, bold=True, font="Consolas")
        txt(slide, x + 0.85, y + 0.42, 1.1, 0.16, f"RGB {c(col)[0]},{c(col)[1]},{c(col)[2]}", size=6.2, color=c("muted"), font="Consolas")
    panel(slide, 7.55, 1.45, 4.8, 4.3, "Visual language", "primary")
    body(slide, 7.88, 1.92, 4.1, 3.4, [
        "Coordinate rails: orient sections and slide position.",
        "Decision pins: identify decisions, gaps, and milestones.",
        "Signal bands: encode maturity, urgency, stakeholder relevance, or workstream.",
        "Evidence cards: compact fielded modules instead of generic decorative cards.",
        "Contour lines: reserved for section dividers and gap/topography slides."
    ], size=10, bullet=False, color=c("ink"))
    rail(slide, "01")


def section(name, kicker, num, note):
    slide = blank()
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    fill(bg, c("night"))
    bg.line.fill.background()
    rail(slide, num, dark=True)
    contour(slide, 7.5, 0.55, 4.8, 5.5, dark=True)
    txt(slide, 0.78, 1.0, 1.15, 0.25, f"{num}", size=13, bold=True, color=c("signal"), font="Consolas")
    txt(slide, 0.78, 1.52, 8.8, 0.72, name, size=31, bold=True, color=c("white"))
    txt(slide, 0.82, 2.55, 6.4, 0.38, kicker.upper(), size=7, bold=True, color=RGBColor(220, 224, 220), font="Consolas")
    body(slide, 0.82, 3.05, 6.8, 0.7, [note], size=14, color=RGBColor(232, 234, 230))


def executive_summary():
    slide = blank()
    title(slide, "Executive Evidence Briefing", "summary", "A1")
    for i, (label, accent) in enumerate([("Priority 1", "primary"), ("Priority 2", "secondary"), ("Priority 3", "risk")]):
        x = 0.78 + i * 4.05
        panel(slide, x, 1.55, 3.65, 4.55, label, accent)
        pin(slide, x + 0.22, 2.0, str(i + 1), accent)
        txt(slide, x + 0.62, 2.02, 2.5, 0.18, "[Strategic evidence priority]", size=10.5, bold=True)
        signal_band(slide, x + 0.3, 2.55, 2.85, "gap | action | decision | timing", accent)
        body(slide, x + 0.32, 2.95, 2.95, 2.45, [
            "Gap: [top evidence gap]",
            "Action: [study / analysis / synthesis]",
            "Decision: [regulatory/access/adoption]",
            "Timing: [readout and decision window]",
            "Risk if delayed: [impact]"
        ], size=8.7, color=c("ink"))


def key_decisions():
    slide = blank()
    title(slide, "Decision Board", "summary", "A2")
    table = atlas_table(slide, 0.72, 1.42, 11.65, 4.8, 7, 5, ["Decision", "Recommendation", "Evidence Basis", "Owner", "Deadline"], [2.55, 2.15, 4.1, 1.3, 1.55])
    rows = [
        ["[Approve priority gap list]", "[Approve/revise]", "[Rationale and supporting evidence]", "[Owner]", "[Date]"],
        ["[Fund priority activity]", "[Approve/defer]", "[Decision value and risk of inaction]", "[Owner]", "[Date]"],
        ["[Select design option]", "[Option A/B]", "[Credibility, timing, feasibility]", "[Owner]", "[Date]"],
        ["[Resolve local need]", "[Global/local split]", "[Evidence reuse vs local gap]", "[Owner]", "[Date]"],
        ["[Confirm governance]", "[RACI/cadence]", "[Execution accountability]", "[Owner]", "[Date]"],
        ["[Other]", "[Recommendation]", "[Basis]", "[Owner]", "[Date]"],
    ]
    for r, row in enumerate(rows, 1):
        for col, val in enumerate(row):
            table.table.cell(r, col).text = val


def asset_snapshot():
    slide = blank()
    title(slide, "Asset Coordinate Sheet", "context", "B1")
    fields = [
        ("Asset / MoA", "[name, mechanism, modality]", "primary"),
        ("Population", "[indication, line, biomarker]", "secondary"),
        ("Lifecycle", "[stage and decision horizon]", "signal"),
        ("Geography", "[global and priority markets]", "primary"),
        ("TPP / TVP", "[clinical and value targets]", "secondary"),
        ("Decision windows", "[regulatory, access, launch, guideline]", "risk"),
    ]
    for i, (h, b, accent) in enumerate(fields):
        panel(slide, 0.78 + (i % 3) * 4.0, 1.5 + (i // 3) * 2.15, 3.55, 1.55, h, accent)
        txt(slide, 1.02 + (i % 3) * 4.0, 2.05 + (i // 3) * 2.15, 2.8, 0.28, b, size=12, color=c("soft_ink"))


def value_story():
    slide = blank()
    title(slide, "Value Story Evidence Spine", "context", "B2")
    x0 = 1.05
    for i, (node, accent) in enumerate([("Value driver", "primary"), ("Evidence required", "secondary"), ("Current support", "signal"), ("Gap", "risk"), ("Tactic", "primary")]):
        x = x0 + i * 2.3
        pin(slide, x, 2.0, str(i + 1), accent)
        txt(slide, x - 0.1, 2.45, 1.25, 0.22, node, size=8, bold=True, font="Consolas")
        panel(slide, x - 0.35, 2.9, 1.75, 1.55, None, accent)
        body(slide, x - 0.18, 3.22, 1.35, 0.8, ["[replace with product-specific content]"], size=8.6, color=c("soft_ink"))
        if i < 4:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.28), Inches(2.14), Inches(x + 2.1), Inches(2.14))
            line(conn, c("line"), 1.1)
    panel(slide, 0.85, 5.3, 11.5, 0.7, "Narrative test", "primary")
    txt(slide, 1.1, 5.56, 10.9, 0.2, "Can each external-facing value message be traced to a source-backed evidence package or a funded evidence activity?", size=11.5, bold=True)


def milestone_map():
    slide = blank()
    title(slide, "Milestone Map", "context", "B3")
    y = 3.15
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(y), Inches(10.8), Inches(0.025))
    fill(base, c("line"))
    base.line.fill.background()
    events = [
        ("Readout", "primary", 1.15),
        ("Regulatory", "secondary", 2.7),
        ("HTA", "risk", 4.55),
        ("Launch", "signal", 6.35),
        ("Guideline", "primary", 8.1),
        ("Lifecycle", "secondary", 9.8),
    ]
    for idx, (label, accent, x) in enumerate(events, 1):
        pin(slide, x, y - 0.15, str(idx), accent)
        txt(slide, x - 0.25, y - 0.65, 0.9, 0.17, label, size=7.2, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
        panel(slide, x - 0.42, y + 0.45, 1.25, 1.0, None, accent)
        body(slide, x - 0.3, y + 0.67, 1.0, 0.55, ["[Evidence required before milestone]"], size=7.1, color=c("soft_ink"))
    body(slide, 0.95, 1.45, 10.5, 0.4, ["Replace milestone labels with product-specific clinical readouts, regulatory interactions, HTA submissions, launch windows, guideline opportunities, and lifecycle decisions."], size=10.5, color=c("soft_ink"))


def stakeholder_map():
    slide = blank()
    title(slide, "Stakeholder Evidence Map", "needs", "C1")
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.35), Inches(2.55), Inches(1.8), Inches(1.8))
    fill(center, c("night"))
    center.line.fill.background()
    txt(slide, 5.62, 3.14, 1.25, 0.17, "ASSET", size=9.5, bold=True, color=c("white"), align=PP_ALIGN.CENTER, font="Consolas")
    nodes = [
        ("Regulators", 1.0, 1.45, "primary"),
        ("HTA / payers", 4.0, 1.15, "secondary"),
        ("HCPs / KOLs", 8.05, 1.35, "signal"),
        ("Patients", 1.35, 4.75, "risk"),
        ("Policy / systems", 5.0, 5.25, "primary"),
        ("Internal teams", 8.75, 4.65, "secondary"),
    ]
    for i, (name, x, y, accent) in enumerate(nodes, 1):
        panel(slide, x, y, 2.2, 0.9, name, accent)
        body(slide, x + 0.18, y + 0.45, 1.75, 0.25, ["[Decision / evidence question]"], size=6.7, color=c("soft_ink"))
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.1), Inches(y + 0.45), Inches(6.25), Inches(3.45))
        line(conn, c("line"), 0.8)


def patient_journey():
    slide = blank()
    title(slide, "Patient Journey Evidence Route", "needs", "C2")
    stages = [("Pre-Dx", "primary"), ("Dx", "secondary"), ("Choice", "signal"), ("Access", "risk"), ("Follow-up", "primary"), ("Outcomes", "secondary")]
    for i, (stage, accent) in enumerate(stages):
        x = 0.9 + i * 1.9
        pin(slide, x, 1.85 + (i % 2) * 0.42, str(i + 1), accent)
        panel(slide, x - 0.25, 2.6 + (i % 2) * 0.42, 1.45, 1.35, stage, accent)
        body(slide, x - 0.08, 3.08 + (i % 2) * 0.42, 1.05, 0.45, ["[barrier]", "[evidence need]"], size=7.2, color=c("soft_ink"))
        if i < 5:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.28), Inches(2.0 + (i % 2) * 0.42), Inches(x + 1.9), Inches(2.0 + ((i + 1) % 2) * 0.42))
            line(conn, c("line"), 1.0)
    panel(slide, 0.85, 5.4, 11.4, 0.75, "Evidence route question", "primary")
    txt(slide, 1.1, 5.7, 10.8, 0.18, "Where does missing evidence create a diagnosis, access, adherence, outcomes, or patient-value barrier?", size=11.5, bold=True)


def evidence_inventory():
    slide = blank()
    title(slide, "Evidence Inventory Instrument", "needs", "C3")
    shape = atlas_table(slide, 0.7, 1.38, 11.8, 5.35, 9, 6, ["Source", "Status", "Population", "Endpoint / Output", "Stakeholder", "Usability"], [2.0, 1.15, 2.0, 2.6, 1.65, 2.4])
    rows = [
        ["Clinical trials", "[available]", "[population]", "[endpoint]", "[reg/HCP]", "[limits]"],
        ["RWE", "[planned]", "[population]", "[outcome]", "[payer/HCP]", "[limits]"],
        ["HEOR model", "[planned]", "[inputs]", "[CE/BIA]", "[HTA]", "[limits]"],
        ["SLR / ITC", "[available]", "[network]", "[comparison]", "[HTA/internal]", "[limits]"],
        ["COA / PRO", "[TBD]", "[population]", "[QoL/burden]", "[patient/HCP]", "[limits]"],
        ["Epidemiology", "[available]", "[market]", "[prevalence]", "[access]", "[limits]"],
        ["Publications", "[planned]", "[audience]", "[message]", "[HCP]", "[limits]"],
        ["Other", "[status]", "[scope]", "[output]", "[stakeholder]", "[limits]"],
    ]
    for r, row in enumerate(rows, 1):
        for col, val in enumerate(row):
            shape.table.cell(r, col).text = val


def need_have_generate():
    slide = blank()
    title(slide, "Need-Have-Generate Equation", "gaps", "D1")
    labels = [("NEED", "primary"), ("HAVE", "secondary"), ("LANDSCAPE", "risk"), ("GENERATE", "signal")]
    for i, (label, accent) in enumerate(labels):
        x = 0.95 + i * 3.0
        panel(slide, x, 2.0, 2.35, 2.2, label, accent)
        body(slide, x + 0.2, 2.55, 1.85, 0.95, ["[stakeholder requirement]", "[available evidence]", "[future change]", "[priority action]"][i:i+1], size=10.5, color=c("soft_ink"))
        if i < 3:
            txt(slide, x + 2.55, 2.85, 0.35, 0.3, "−" if i == 0 else "+", size=19, bold=True, color=c("muted"), align=PP_ALIGN.CENTER)
    txt(slide, 0.95, 5.25, 11.2, 0.35, "Prioritized evidence to generate = stakeholder need minus usable evidence, adjusted for future landscape change.", size=15, bold=True, color=c("night"), align=PP_ALIGN.CENTER)


def gap_matrix():
    slide = blank()
    title(slide, "Gap Topography Matrix", "gaps", "D2")
    contour(slide, 9.0, 1.45, 2.75, 2.75)
    shape = atlas_table(slide, 0.7, 1.35, 11.65, 5.15, 8, 6, ["Gap", "Decision", "Current Evidence", "Priority", "Risk", "Action"], [2.35, 2.25, 2.6, 1.05, 1.35, 2.05])
    for r in range(1, 8):
        vals = [f"G{r}: [gap]", "[stakeholder decision]", "[evidence]", "[P1/P2/P3]", "[impact]", "[tactic]"]
        for col, val in enumerate(vals):
            shape.table.cell(r, col).text = val


def prioritization_heatmap():
    slide = blank()
    title(slide, "Gap Prioritization Field", "gaps", "D3")
    panel(slide, 0.85, 1.45, 5.7, 4.9, "Impact x feasibility", "primary")
    grid_colors = [[c("mist"), c("secondary"), c("risk")], [c("secondary"), c("secondary"), c("signal")], [c("secondary"), c("signal"), c("signal")]]
    for r in range(3):
        for col in range(3):
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.38 + col * 1.35), Inches(2.1 + r * 0.9), Inches(1.25), Inches(0.8))
            fill(cell, grid_colors[2 - r][col])
            line(cell, c("paper"), 1.2)
            txt(slide, 1.74 + col * 1.35, 2.38 + r * 0.9, 0.45, 0.12, f"P{max(1, 4-r-col)}", size=8, bold=True, color=c("night"), align=PP_ALIGN.CENTER, font="Consolas")
    txt(slide, 2.2, 5.05, 1.8, 0.16, "FEASIBILITY", size=7, bold=True, color=c("muted"), font="Consolas", align=PP_ALIGN.CENTER)
    ylbl = txt(slide, 0.95, 2.7, 0.3, 1.05, "IMPACT", size=7, bold=True, color=c("muted"), font="Consolas", align=PP_ALIGN.CENTER)
    ylbl.rotation = 270
    shape = atlas_table(slide, 6.95, 1.45, 5.35, 4.85, 6, 3, ["Criterion", "Score", "Notes"], [2.2, 1.0, 2.15])
    for r, row in enumerate([["Impact", "1-5", "[decision value]"], ["Urgency", "1-5", "[timing]"], ["Feasibility", "1-5", "[data/design]"], ["Reuse", "1-5", "[global-local]"], ["Risk", "1-5", "[inaction]"]], 1):
        for col, val in enumerate(row):
            shape.table.cell(r, col).text = val


def orq_picot():
    slide = blank()
    title(slide, "Open Research Question Ledger", "gaps", "D4")
    shape = atlas_table(slide, 0.65, 1.35, 12.0, 5.4, 8, 5, ["ORQ", "Gap", "Research Question", "PICOT / Design", "Decision Use"], [0.75, 0.8, 4.1, 4.05, 2.3])
    for r in range(1, 8):
        vals = [f"ORQ{r}", "G[ ]", "[answerable evidence question]", "P:[ ] I:[ ] C:[ ] O:[ ] T:[ ]", "[decision]"]
        for col, val in enumerate(vals):
            shape.table.cell(r, col).text = val


def roadmap():
    slide = blank()
    title(slide, "Integrated Evidence Roadmap", "roadmap", "E1")
    quarters = ["Y1 Q1", "Q2", "Q3", "Q4", "Y2 Q1", "Q2", "Q3", "Q4"]
    for i, q in enumerate(quarters):
        txt(slide, 3.2 + i * 1.02, 1.32, 0.75, 0.14, q, size=6.6, bold=True, color=c("muted"), font="Consolas", align=PP_ALIGN.CENTER)
        v = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.55 + i * 1.02), Inches(1.6), Inches(0.01), Inches(4.65))
        fill(v, c("line"))
        v.line.fill.background()
    lanes = [("Clinical/reg", "primary"), ("RWE", "secondary"), ("HEOR/access", "signal"), ("Patient/COA", "risk"), ("Publications", "primary")]
    for i, (lane, accent) in enumerate(lanes):
        y = 1.78 + i * 0.88
        txt(slide, 0.78, y + 0.18, 1.8, 0.18, lane, size=7.5, bold=True, font="Consolas")
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.25 + (i % 3) * 0.55), Inches(y + 0.12), Inches(1.85 + (i % 2) * 0.62), Inches(0.28))
        fill(bar, c(accent))
        bar.line.fill.background()
        txt(slide, 3.38 + (i % 3) * 0.55, y + 0.18, 1.45, 0.09, "[activity]", size=5.7, bold=True, color=c("night"), align=PP_ALIGN.CENTER, font="Consolas")
        pin(slide, 5.45 + i * 0.92, y + 0.05, str(i + 1), accent)
    panel(slide, 0.78, 6.28, 11.5, 0.48, "Decision gates: [protocol] | [data cut] | [abstract] | [regulatory] | [HTA] | [launch]", "primary")


def activity_tracker():
    slide = blank()
    title(slide, "Evidence Activity Tracker", "roadmap", "E2")
    shape = atlas_table(slide, 0.55, 1.35, 12.25, 5.4, 9, 8, ["ID", "Activity", "Gap", "Method / Data", "Owner", "Timing", "Status", "Priority"], [0.55, 2.0, 0.9, 2.65, 1.15, 1.25, 1.65, 1.1], font_size=6.2)
    for r in range(1, 9):
        vals = [f"A{r}", "[study / analysis]", "G[ ]", "[design / source]", "[owner]", "[start-readout]", "[idea/planned/active]", "[P1]"]
        for col, val in enumerate(vals):
            shape.table.cell(r, col).text = val


def claims_matrix():
    slide = blank()
    title(slide, "Claims Matrix Route Map", "roadmap", "E3")
    shape = atlas_table(slide, 0.6, 1.35, 12.1, 5.35, 8, 6, ["Claim", "Stakeholder", "Required Evidence", "Current Support", "Gap", "Tactic"], [2.35, 1.55, 2.75, 2.0, 1.55, 1.9], font_size=6.5)
    rows = ["Clinical benefit", "Safety", "Economic value", "Patient value", "Implementation", "Other", "Other"]
    for r, label in enumerate(rows, 1):
        vals = [f"[{label} claim]", "[audience]", "[required standard]", "[support]", "[gap]", "[activity]"]
        for col, val in enumerate(vals):
            shape.table.cell(r, col).text = val


def study_concept():
    slide = blank()
    title(slide, "Study Concept Card", "roadmap", "E4")
    cards = [
        ("Study / activity", "[name and ID]", "primary"),
        ("Gap and decision use", "[gap, ORQ, stakeholder]", "secondary"),
        ("Objective", "[objective and hypothesis]", "signal"),
        ("Design", "[population, comparator, endpoint, data source]", "primary"),
        ("Feasibility", "[cost, complexity, dependencies]", "risk"),
        ("Pull-through", "[publication, dossier, claims, medical]", "secondary"),
    ]
    for i, (h, b, accent) in enumerate(cards):
        panel(slide, 0.78 + (i % 3) * 4.0, 1.45 + (i // 3) * 2.25, 3.55, 1.65, h, accent)
        body(slide, 1.0 + (i % 3) * 4.0, 2.0 + (i // 3) * 2.25, 2.9, 0.55, [b], size=10, color=c("soft_ink"))
    panel(slide, 0.78, 6.15, 11.55, 0.42, "Decision: [approve / revise / defer]    Next step: [action]    Date: [ ]", "signal")


def heor_rwe_patient():
    slide = blank()
    title(slide, "Specialized Evidence Workstreams", "roadmap", "E5")
    work = [
        ("HEOR / access", ["SLR / ITC", "economic model", "budget impact", "payer research"], "secondary"),
        ("RWE / data", ["retrospective database", "prospective study", "registry", "external control"], "primary"),
        ("Patient-centered", ["COA / PRO", "HRQoL", "preference", "caregiver burden"], "risk"),
        ("Evidence synthesis", ["catalog", "source log", "claims support", "publication pull-through"], "signal"),
    ]
    for i, (h, lines, accent) in enumerate(work):
        panel(slide, 0.82 + (i % 2) * 5.85, 1.55 + (i // 2) * 2.45, 5.25, 1.95, h, accent)
        body(slide, 1.1 + (i % 2) * 5.85, 2.0 + (i // 2) * 2.45, 4.4, 0.95, lines, size=9.3, color=c("soft_ink"), bullet=False)


def governance():
    slide = blank()
    title(slide, "Evidence Control Room", "governance", "F1")
    shape = atlas_table(slide, 0.72, 1.4, 6.0, 4.7, 7, 4, ["Forum", "Purpose", "Cadence", "Rights"], [1.35, 2.25, 1.05, 1.35], font_size=6.6)
    for r, row in enumerate([["Core team", "[maintain plan]", "[monthly]", "[recommend]"], ["Evidence council", "[prioritize/fund]", "[quarterly]", "[approve]"], ["Global-local", "[adapt needs]", "[quarterly]", "[recommend]"], ["Leadership", "[trade-offs]", "[semiannual]", "[approve]"], ["Publications", "[pull-through]", "[monthly]", "[recommend]"], ["Other", "[purpose]", "[cadence]", "[rights]"]], 1):
        for col, val in enumerate(row):
            shape.table.cell(r, col).text = val
    panel(slide, 7.05, 1.4, 5.2, 4.7, "Update triggers", "risk")
    body(slide, 7.35, 1.9, 4.35, 3.45, ["Clinical readout", "Regulatory feedback", "HTA / payer expectation change", "Competitor approval or data", "Guideline update", "Budget or feasibility change", "Regional access need"], size=10.2, color=c("ink"))


def kpi_dashboard():
    slide = blank()
    title(slide, "Evidence Impact Dashboard", "governance", "F2")
    metrics = [("Gap closure", "primary"), ("Milestones", "secondary"), ("Evidence reuse", "signal"), ("Pull-through", "risk")]
    for i, (m, accent) in enumerate(metrics):
        panel(slide, 0.78 + i * 3.0, 1.45, 2.55, 1.35, m, accent)
        txt(slide, 1.05 + i * 3.0, 2.03, 1.45, 0.28, "[metric]", size=18, bold=True, color=c("night"), align=PP_ALIGN.CENTER)
        signal_band(slide, 1.05 + i * 3.0, 2.48, 1.45, "[target]", accent)
    shape = atlas_table(slide, 0.78, 3.35, 11.55, 2.8, 6, 5, ["KPI", "Definition", "Baseline", "Target", "Owner"], [2.0, 4.4, 1.5, 1.5, 2.15])
    for r in range(1, 6):
        vals = ["[KPI]", "[definition]", "[baseline]", "[target]", "[owner]"]
        for col, val in enumerate(vals):
            shape.table.cell(r, col).text = val


def risks():
    slide = blank()
    title(slide, "Risk And Dependency Register", "governance", "F3")
    shape = atlas_table(slide, 0.62, 1.35, 12.1, 5.4, 8, 5, ["Risk / Dependency", "Activity", "Impact", "Mitigation", "Decision Needed"], [2.65, 2.0, 1.3, 3.7, 2.45])
    for r in range(1, 8):
        vals = ["[risk]", "[activity]", "[H/M/L]", "[mitigation or alternate evidence path]", "[decision / owner]"]
        for col, val in enumerate(vals):
            shape.table.cell(r, col).text = val


def source_basis():
    slide = blank()
    title(slide, "Source-Informed Template Basis", "appendix", "G1")
    shape = atlas_table(slide, 0.65, 1.35, 12.0, 5.4, 9, 3, ["Source", "Template Contribution", "Type"], [3.25, 6.65, 2.1], font_size=6.7)
    rows = [
        ["Lumanity", "Need-have-generate logic, future landscape, global harmonization, KPIs", "Framework/PDF"],
        ["Cencora", "HEOR, RWE, SLR/evidence-synthesis, market-access components", "Articles/service"],
        ["ZS", "IEP deliverable structure, roadmap, governance, execution tracking", "How-to"],
        ["Veranex", "Claims matrix, scorecard, roadmap, budget and ROI views", "Framework"],
        ["DiMe", "Stage resources, stakeholder maps, concrete digital-health case examples", "Toolkit"],
        ["MAPS/Prescient", "Medical affairs standards, adaptive plans, refresh cadence", "Guidance"],
        ["McKinsey/Springer", "Evidence catalog, prioritization, economics, value-driver framing", "Thought leadership"],
        ["Oxford PharmaGenesis", "PICOT, data-source selection, study mapping workflow", "Process model"],
    ]
    for r, row in enumerate(rows, 1):
        for col, val in enumerate(row):
            shape.table.cell(r, col).text = val


def appendix_logs():
    slide = blank()
    title(slide, "Appendix Worksheets", "appendix", "G2")
    panels = [
        ("Regional adaptation", "local needs, reusable global evidence, local gaps, owner", "primary"),
        ("Evidence source log", "source, date, owner, confidence, limitation, claim link", "secondary"),
        ("Methodology", "inputs, workshops, scoring rubric, assumptions, exclusions", "signal"),
        ("Glossary", "IEGP/IEP, ORQ, PICOT, RWE, HEOR, HTA, COA/PRO, IIR", "risk"),
    ]
    for i, (h, b, accent) in enumerate(panels):
        panel(slide, 0.85 + (i % 2) * 5.75, 1.55 + (i // 2) * 2.4, 5.15, 1.8, h, accent)
        body(slide, 1.1 + (i % 2) * 5.75, 2.08 + (i // 2) * 2.4, 4.45, 0.5, [b], size=10, color=c("soft_ink"))


def checklist():
    slide = blank()
    title(slide, "Final Review Checklist", "appendix", "G3")
    items = [
        ("Strategy linked to decisions", "primary"),
        ("Stakeholder needs documented", "secondary"),
        ("Evidence inventory quality checked", "signal"),
        ("Gaps prioritized with rubric", "risk"),
        ("Activities have owner/timing/status", "primary"),
        ("Governance cadence agreed", "secondary"),
        ("Source-backed facts separated from assumptions", "signal"),
        ("Next review date set", "risk"),
    ]
    for i, (item, accent) in enumerate(items):
        x = 0.9 + (i % 2) * 5.75
        y = 1.35 + (i // 2) * 1.25
        pin(slide, x, y + 0.05, str(i + 1), accent)
        panel(slide, x + 0.42, y, 4.85, 0.62, None, accent)
        txt(slide, x + 0.65, y + 0.22, 4.25, 0.12, item, size=8.8, bold=True)


slides = [
    cover,
    design_tokens,
    lambda: section("Strategy Coordinates", "context and decisions", "A", "Define the decisions, strategy, and value story before building evidence activities."),
    executive_summary,
    key_decisions,
    lambda: section("Asset And Value Context", "asset framing", "B", "Document the asset, value narrative, milestones, and decision windows that govern evidence needs."),
    asset_snapshot,
    value_story,
    milestone_map,
    lambda: section("Stakeholder Evidence Map", "needs and inventory", "C", "Map who needs evidence, what decisions they make, and what evidence exists today."),
    stakeholder_map,
    patient_journey,
    evidence_inventory,
    lambda: section("Gap Topography", "prioritization", "D", "Turn evidence needs into prioritized gaps and answerable research questions."),
    need_have_generate,
    gap_matrix,
    prioritization_heatmap,
    orq_picot,
    lambda: section("Roadmap Architecture", "activities and pull-through", "E", "Convert priority gaps into studies, syntheses, claims support, and lifecycle milestones."),
    roadmap,
    activity_tracker,
    claims_matrix,
    study_concept,
    heor_rwe_patient,
    lambda: section("Evidence Control Room", "governance and impact", "F", "Operate the plan as a living system with owners, decisions, KPIs, risks, and refresh triggers."),
    governance,
    kpi_dashboard,
    risks,
    lambda: section("Appendix System", "source basis and worksheets", "G", "Keep provenance, methodology, regional adaptation, and final review materials audit-ready."),
    source_basis,
    appendix_logs,
    checklist,
]

for make in slides:
    make()

PRS.save(OUT_FILE)
print(OUT_FILE.resolve())
